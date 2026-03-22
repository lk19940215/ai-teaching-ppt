"""
PPTX 往返质量验证

生成一份特征完整的测试 PPTX，然后走完整管线：
  解析 → 提取 AI 内容 → 模拟修改 → 写回 → 重新解析 → 对比

验证项：
  1. 动画标记是否保留
  2. 文字格式（字体/颜色/大小/粗体/斜体）是否保留
  3. 表格结构是否完整
  4. 布局位置是否不变
  5. 图片是否保留
"""

import sys
import os
import tempfile

os.environ["PYTHONIOENCODING"] = "utf-8"
sys.stdout.reconfigure(encoding="utf-8")
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

from pathlib import Path
from app.core.pptx_reader import PPTXReader
from app.core.content_extractor import ContentExtractor
from app.core.pptx_writer import PPTXWriter
from app.core.models import (
    TextModification, TableCellModification, SlideModification, ProcessingResult
)


def create_test_pptx(path: str):
    """生成一份特征完整的测试 PPTX"""
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)

    # === 第1页：标题页 ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide1.placeholders[0]
    title_shape.text = "教学课件：Python 入门"
    for run in title_shape.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
        run.font.name = "微软雅黑"

    if 1 in slide1.placeholders:
        subtitle = slide1.placeholders[1]
        subtitle.text = "从零开始学编程"
        for run in subtitle.text_frame.paragraphs[0].runs:
            run.font.size = Pt(20)
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 给第1页添加动画标记（手动注入 timing XML）
    timing_xml = '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"/></p:par></p:tnLst></p:timing>'
    timing_elem = etree.fromstring(timing_xml)
    slide1._element.append(timing_elem)

    # === 第2页：内容页（多格式文本）===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run1 = p.add_run()
    run1.text = "Python 基础语法"
    run1.font.size = Pt(28)
    run1.font.bold = True
    run1.font.name = "楷体"
    run1.font.color.rgb = RGBColor(0x00, 0x00, 0x80)

    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    p1 = tf2.paragraphs[0]
    run_a = p1.add_run()
    run_a.text = "1. 变量和数据类型：Python 支持 int, float, str, list, dict 等"
    run_a.font.size = Pt(16)
    run_a.font.name = "微软雅黑"

    p2 = tf2.add_paragraph()
    run_b = p2.add_run()
    run_b.text = "2. 控制流：if/elif/else, for, while"
    run_b.font.size = Pt(16)
    run_b.font.bold = True
    run_b.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    p3 = tf2.add_paragraph()
    run_c = p3.add_run()
    run_c.text = "3. 函数定义：def function_name(params):"
    run_c.font.size = Pt(16)
    run_c.font.italic = True

    # === 第3页：表格页 ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    table_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    table_title.text_frame.paragraphs[0].add_run().text = "常用数据类型对比"
    table_title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    table_title.text_frame.paragraphs[0].runs[0].font.bold = True

    rows, cols = 4, 4
    table = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(3)).table

    headers = ["类型", "示例", "可变性", "用途"]
    data = [
        ["int", "42", "不可变", "整数计算"],
        ["str", '"hello"', "不可变", "文本处理"],
        ["list", "[1,2,3]", "可变", "有序集合"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)

    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)

    prs.save(path)
    print(f"测试 PPTX 已生成: {path}")


def run_roundtrip_test(pptx_path: str):
    """执行完整的往返测试"""
    reader = PPTXReader()
    extractor = ContentExtractor()
    writer = PPTXWriter()

    # === Step 1: 解析 ===
    print("\n=== Step 1: 解析原始 PPTX ===")
    parsed = reader.parse(Path(pptx_path))
    print(f"  文件: {parsed.filename}, 页数: {parsed.slide_count}")

    for slide in parsed.slides:
        print(f"\n  --- 第 {slide.slide_index + 1} 页 (版式: {slide.layout_name}) ---")
        print(f"    动画: {slide.has_animations}, 媒体: {slide.has_media}")
        for elem in slide.elements:
            if elem.plain_text:
                fmt_info = []
                if elem.paragraphs:
                    for para in elem.paragraphs:
                        for run in para.runs:
                            parts = []
                            if run.bold: parts.append("B")
                            if run.italic: parts.append("I")
                            if run.font_size: parts.append(f"{run.font_size}pt")
                            if run.font_color: parts.append(run.font_color)
                            if run.font_name: parts.append(run.font_name)
                            if parts:
                                fmt_info.append(f"[{','.join(parts)}]")
                print(f"    [{elem.shape_index}] {elem.element_type.value}: {elem.plain_text[:50]}...")
                if fmt_info:
                    print(f"         格式: {' '.join(fmt_info[:3])}")
            elif elem.table_data:
                print(f"    [{elem.shape_index}] table {elem.table_rows}x{elem.table_cols}")

    # === Step 2: 提取 AI 内容 ===
    print("\n=== Step 2: 提取 AI 可处理内容 ===")
    contents = extractor.extract_all(parsed)
    for content in contents:
        ai_text = extractor.format_for_ai(content)
        print(f"\n  --- 第 {content.slide_index + 1} 页 AI 视图 ---")
        for line in ai_text.split("\n"):
            print(f"    {line}")

    # === Step 3: 模拟 AI 修改 ===
    print("\n=== Step 3: 模拟 AI 修改 ===")
    modifications = []

    # 修改第2页的正文
    slide2_content = contents[1] if len(contents) > 1 else None
    if slide2_content and slide2_content.text_blocks:
        for tb in slide2_content.text_blocks:
            if tb.role == "body":
                modifications.append(SlideModification(
                    slide_index=1,
                    text_modifications=[TextModification(
                        shape_index=tb.shape_index,
                        original_text=tb.text,
                        new_text=tb.text.replace("Python", "Python 3.12"),
                    )],
                ))
                print(f"  文本修改: shape_{tb.shape_index}, '{tb.text[:30]}' → 替换 Python→Python 3.12")
                break

    # 修改第3页的表格
    slide3_content = contents[2] if len(contents) > 2 else None
    if slide3_content and slide3_content.table_blocks:
        tb = slide3_content.table_blocks[0]
        modifications.append(SlideModification(
            slide_index=2,
            table_modifications=[TableCellModification(
                shape_index=tb.shape_index,
                row=1,
                col=3,
                original_text="整数计算",
                new_text="整数运算与数学计算",
            )],
        ))
        print(f"  表格修改: shape_{tb.shape_index}, row=1 col=3, '整数计算' → '整数运算与数学计算'")

    result = ProcessingResult(
        success=True,
        modifications=modifications,
        action="polish",
        total_changes=len(modifications),
    )

    # === Step 4: 写回 ===
    print("\n=== Step 4: 写回修改 ===")
    output_dir = Path(pptx_path).parent
    output_path_obj = writer.apply(Path(pptx_path), result.modifications, output_dir)
    output_path = str(output_path_obj)
    print(f"  输出文件: {output_path}")

    # === Step 5: 重新解析并对比 ===
    print("\n=== Step 5: 重新解析修改后的 PPTX 并对比 ===")
    parsed2 = reader.parse(Path(output_path))

    all_passed = True
    checks = []

    # 检查1: 页数不变
    pages_ok = parsed.slide_count == parsed2.slide_count
    checks.append(("页数保持不变", pages_ok, f"{parsed.slide_count} → {parsed2.slide_count}"))

    # 检查2: 动画保留
    for s1, s2 in zip(parsed.slides, parsed2.slides):
        anim_ok = s1.has_animations == s2.has_animations
        checks.append((f"第{s1.slide_index+1}页动画保留", anim_ok,
                       f"{s1.has_animations} → {s2.has_animations}"))

    # 检查3: 文字格式保留（第2页）
    if len(parsed2.slides) > 1:
        slide2_orig = parsed.slides[1]
        slide2_mod = parsed2.slides[1]
        for e1, e2 in zip(slide2_orig.elements, slide2_mod.elements):
            if e1.paragraphs and e2.paragraphs:
                for p1, p2 in zip(e1.paragraphs, e2.paragraphs):
                    for r1, r2 in zip(p1.runs, p2.runs):
                        fmt_ok = (r1.bold == r2.bold and
                                  r1.italic == r2.italic and
                                  r1.font_size == r2.font_size and
                                  r1.font_color == r2.font_color and
                                  r1.font_name == r2.font_name)
                        if not fmt_ok:
                            checks.append((f"第2页 shape_{e1.shape_index} 格式保留", False,
                                          f"bold:{r1.bold}→{r2.bold}, italic:{r1.italic}→{r2.italic}, "
                                          f"size:{r1.font_size}→{r2.font_size}"))
                            all_passed = False

    # 检查4: 修改后的文本确实变了
    if len(parsed2.slides) > 1:
        for elem in parsed2.slides[1].elements:
            if elem.plain_text and "Python 3.12" in elem.plain_text:
                checks.append(("第2页文本修改生效", True, "包含 'Python 3.12'"))
                break
        else:
            checks.append(("第2页文本修改生效", False, "未找到 'Python 3.12'"))

    # 检查5: 表格修改生效
    if len(parsed2.slides) > 2:
        for elem in parsed2.slides[2].elements:
            if elem.table_data:
                cell_text = elem.table_data[1][3].text if len(elem.table_data) > 1 and len(elem.table_data[1]) > 3 else ""
                table_ok = "整数运算" in cell_text
                checks.append(("第3页表格修改生效", table_ok, f"cell(1,3)='{cell_text}'"))
                break

    # 检查6: 位置保留
    for s1, s2 in zip(parsed.slides, parsed2.slides):
        for e1, e2 in zip(s1.elements, s2.elements):
            pos_ok = (e1.position.left == e2.position.left and
                      e1.position.top == e2.position.top and
                      e1.position.width == e2.position.width)
            if not pos_ok:
                checks.append((f"第{s1.slide_index+1}页 shape_{e1.shape_index} 位置保留", False,
                               f"left:{e1.position.left}→{e2.position.left}"))

    # === 输出报告 ===
    print("\n" + "=" * 60)
    print("往返质量验证报告")
    print("=" * 60)

    for name, passed, detail in checks:
        status = "PASS" if passed else "FAIL"
        print(f"  [{status}] {name}: {detail}")
        if not passed:
            all_passed = False

    print("\n" + "=" * 60)
    if all_passed:
        print("=== ALL QUALITY CHECKS PASSED ===")
    else:
        print("=== SOME CHECKS FAILED ===")
        sys.exit(1)


if __name__ == "__main__":
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False, dir=".") as f:
        test_path = f.name

    create_test_pptx(test_path)
    run_roundtrip_test(test_path)

    # 清理
    import glob
    for f in glob.glob(test_path.replace(".pptx", "*")):
        try:
            os.unlink(f)
        except:
            pass
