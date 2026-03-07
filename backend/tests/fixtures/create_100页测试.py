import sys
sys.path.insert(0, '/e/Code/ai-teaching-ppt/backend')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 创建 100 页 PPT
prs = Presentation()

for i in range(100):
    # 使用标题和内容布局
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title = slide.shapes.title
    title.text = f"第 {i+1} 页 - 测试幻灯片"
    
    # 设置内容
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"这是第 {i+1} 页的测试内容"
    
    # 添加一些项目符号
    for j in range(3):
        p = tf.add_paragraph()
        p.text = f"要点 {j+1} - 页面 {i+1}"
        p.level = 0

# 保存
output_path = '/e/Code/ai-teaching-ppt/backend/tests/fixtures/large_test.pptx'
prs.save(output_path)
print(f"已创建 100 页 PPT: {output_path}")

# 验证
prs2 = Presentation(output_path)
print(f"验证页数：{len(prs2.slides)}")
