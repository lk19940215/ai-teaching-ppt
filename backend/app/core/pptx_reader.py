"""
PPTX 解析器

将 PPTX 文件解析为 ParsedPresentation 中间模型。
保留 shape_index 用于后续写回定位。

技术文档: docs/technical-spec.md §5.1
"""

import base64
import logging
from io import BytesIO
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

from .models import (
    TextRun, Paragraph, Position, TableCell,
    ElementType, SlideElement, ParsedSlide, ParsedPresentation,
)

logger = logging.getLogger(__name__)

PLACEHOLDER_TEXTS = frozenset([
    "添加标题", "添加标题内容", "添加副标题",
    "单击此处添加标题", "单击此处添加副标题",
    "单击此处编辑标题", "单击此处编辑副标题",
    "请输入标题", "请输入副标题",
    "点击此处添加文本", "在此处输入文本", "双击此处添加文本",
    "click to add title", "click to add text", "click to add subtitle",
])


class PPTXReader:
    """PPTX 解析器：PPTX 文件 → ParsedPresentation"""

    def __init__(self, max_image_size: int = 512):
        self.max_image_size = max_image_size

    def parse(self, file_path: Path) -> ParsedPresentation:
        """解析 PPTX 文件为 ParsedPresentation"""
        logger.info(f"开始解析: {file_path.name}")
        prs = Presentation(str(file_path))

        slides = []
        for slide_idx, slide in enumerate(prs.slides):
            parsed_slide = self._parse_slide(slide, slide_idx)
            slides.append(parsed_slide)

        doc_title = self._extract_doc_title(slides)

        result = ParsedPresentation(
            filename=file_path.name,
            slide_count=len(slides),
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            slides=slides,
            title=doc_title,
        )
        logger.info(f"解析完成: {len(slides)} 页, 标题: {doc_title}")
        return result

    def _parse_slide(self, slide, slide_idx: int) -> ParsedSlide:
        """解析单页幻灯片"""
        elements = []
        for shape_idx, shape in enumerate(slide.shapes):
            element = self._parse_shape(shape, shape_idx)
            if element is not None:
                elements.append(element)

        layout_name = None
        try:
            if slide.slide_layout:
                layout_name = slide.slide_layout.name
        except Exception:
            pass

        has_notes = False
        notes_text = None
        try:
            if slide.has_notes_slide:
                raw = slide.notes_slide.notes_text_frame.text.strip()
                if raw:
                    has_notes = True
                    notes_text = raw
        except Exception:
            pass

        return ParsedSlide(
            slide_index=slide_idx,
            elements=elements,
            layout_name=layout_name,
            has_notes=has_notes,
            notes_text=notes_text,
            has_animations=self._check_animations(slide),
            has_media=self._check_media(slide),
        )

    def _parse_shape(self, shape, shape_idx: int) -> Optional[SlideElement]:
        """解析单个 Shape → SlideElement"""
        position = Position(
            left=shape.left or 0,
            top=shape.top or 0,
            width=shape.width or 0,
            height=shape.height or 0,
        )

        shape_type = shape.shape_type

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._parse_image(shape, shape_idx, position)

        if shape.has_table:
            return self._parse_table(shape, shape_idx, position)

        if shape.has_text_frame:
            return self._parse_text(shape, shape_idx, position)

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            group_text = self._extract_group_text(shape)
            return SlideElement(
                shape_index=shape_idx,
                element_type=ElementType.GROUP,
                position=position,
                name=shape.name,
                plain_text=group_text if group_text else None,
            )

        if shape_type in (MSO_SHAPE_TYPE.CHART,):
            return SlideElement(
                shape_index=shape_idx,
                element_type=ElementType.CHART,
                position=position,
                name=shape.name,
            )

        if shape_type in (MSO_SHAPE_TYPE.MEDIA,):
            return SlideElement(
                shape_index=shape_idx,
                element_type=ElementType.MEDIA,
                position=position,
                name=shape.name,
            )

        return SlideElement(
            shape_index=shape_idx,
            element_type=ElementType.SHAPE,
            position=position,
            name=shape.name,
        )

    def _parse_text(self, shape, shape_idx: int, position: Position) -> Optional[SlideElement]:
        """解析文本 Shape，保留每个 Run 的格式信息"""
        paragraphs = []
        all_text_parts = []

        for para in shape.text_frame.paragraphs:
            runs = []
            for run in para.runs:
                if not run.text:
                    continue
                runs.append(TextRun(
                    text=run.text,
                    bold=run.font.bold,
                    italic=run.font.italic,
                    underline=run.font.underline,
                    font_size=run.font.size.pt if run.font.size else None,
                    font_color=self._extract_color(run.font),
                    font_name=run.font.name,
                ))

            if not runs and para.text.strip():
                runs.append(TextRun(text=para.text.strip()))

            if runs:
                paragraphs.append(Paragraph(
                    runs=runs,
                    alignment=self._extract_alignment(para),
                    level=para.level or 0,
                    line_spacing=self._extract_line_spacing(para),
                ))
                all_text_parts.append("".join(r.text for r in runs))

        if not paragraphs:
            return None

        plain_text = "\n".join(all_text_parts)

        if self._is_placeholder_text(plain_text):
            return None

        is_title = self._is_title_shape(shape)
        is_ph = shape.is_placeholder
        ph_type = self._get_placeholder_type(shape)
        has_link = self._check_hyperlinks(shape)

        return SlideElement(
            shape_index=shape_idx,
            element_type=ElementType.TEXT_BOX,
            position=position,
            name=shape.name,
            paragraphs=paragraphs,
            plain_text=plain_text,
            is_title=is_title,
            is_placeholder=is_ph,
            placeholder_type=ph_type,
            has_hyperlink=has_link,
        )

    def _parse_table(self, shape, shape_idx: int, position: Position) -> Optional[SlideElement]:
        """解析表格"""
        try:
            table = shape.table
            table_data = []
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_paras = []
                    for para in cell.text_frame.paragraphs:
                        runs = []
                        for run in para.runs:
                            if run.text:
                                runs.append(TextRun(
                                    text=run.text,
                                    bold=run.font.bold,
                                    italic=run.font.italic,
                                    font_size=run.font.size.pt if run.font.size else None,
                                    font_color=self._extract_color(run.font),
                                    font_name=run.font.name,
                                ))
                        if runs:
                            cell_paras.append(Paragraph(runs=runs))

                    row_cells.append(TableCell(
                        text=cell.text.strip() if cell.text else "",
                        paragraphs=cell_paras,
                    ))
                table_data.append(row_cells)

            return SlideElement(
                shape_index=shape_idx,
                element_type=ElementType.TABLE,
                position=position,
                name=shape.name,
                table_data=table_data,
                table_rows=len(table.rows),
                table_cols=len(table.columns),
            )
        except Exception as e:
            logger.warning(f"表格解析失败 shape_idx={shape_idx}: {e}")
            return SlideElement(
                shape_index=shape_idx,
                element_type=ElementType.TABLE,
                position=position,
                name=shape.name,
            )

    def _parse_image(self, shape, shape_idx: int, position: Position) -> Optional[SlideElement]:
        """解析图片，提取缩略 base64"""
        image_b64 = None
        image_fmt = None
        try:
            image = shape.image
            image_bytes = image.blob
            image_fmt = image.ext
            image_b64 = self._compress_image(image_bytes)
        except Exception as e:
            logger.warning(f"图片解析失败 shape_idx={shape_idx}: {e}")

        return SlideElement(
            shape_index=shape_idx,
            element_type=ElementType.IMAGE,
            position=position,
            name=shape.name,
            image_base64=image_b64,
            image_format=image_fmt,
        )

    # ---- 辅助方法 ----

    def _extract_color(self, font) -> Optional[str]:
        """提取字体颜色"""
        try:
            if font.color is None:
                return None
            color_type = getattr(font.color, "type", None)
            if color_type is not None and str(color_type) == "THEME":
                theme_color = getattr(font.color, "theme_color", None)
                if theme_color is not None:
                    names = {
                        0: "theme_dark1", 1: "theme_light1",
                        2: "theme_dark2", 3: "theme_light2",
                        4: "theme_accent1", 5: "theme_accent2",
                        6: "theme_accent3", 7: "theme_accent4",
                        8: "theme_accent5", 9: "theme_accent6",
                    }
                    return names.get(theme_color, f"theme_{theme_color}")
            if hasattr(font.color, "rgb") and font.color.rgb:
                rgb = font.color.rgb
                return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        except Exception:
            pass
        return None

    def _extract_alignment(self, para) -> Optional[str]:
        """提取段落对齐方式"""
        try:
            from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
            mapping = {
                PP_PARAGRAPH_ALIGNMENT.LEFT: "left",
                PP_PARAGRAPH_ALIGNMENT.CENTER: "center",
                PP_PARAGRAPH_ALIGNMENT.RIGHT: "right",
                PP_PARAGRAPH_ALIGNMENT.JUSTIFY: "justify",
            }
            if para.alignment is not None:
                return mapping.get(para.alignment)
        except Exception:
            pass
        return None

    def _extract_line_spacing(self, para) -> Optional[float]:
        """提取行距"""
        try:
            pf = para.paragraph_format
            if pf and pf.line_spacing is not None:
                return float(pf.line_spacing)
        except Exception:
            pass
        return None

    def _is_placeholder_text(self, text: str) -> bool:
        """判断是否为模板占位符文本"""
        if not text:
            return True
        normalized = text.strip().lower().replace(" ", "").replace("\t", "")
        if len(normalized) > 30:
            return False
        return normalized in {p.lower().replace(" ", "") for p in PLACEHOLDER_TEXTS}

    def _is_title_shape(self, shape) -> bool:
        """判断 Shape 是否为标题（不含副标题和正文）"""
        if not shape.is_placeholder:
            return False
        try:
            ph_idx = shape.placeholder_format.idx
            # 0=title, 10=center title; 1=body, 12=subtitle 不算标题
            return ph_idx in (0, 10)
        except Exception:
            return False

    def _get_placeholder_type(self, shape) -> Optional[str]:
        """获取占位符类型名（基于 OOXML 标准 placeholder idx）"""
        if not shape.is_placeholder:
            return None
        try:
            ph_idx = shape.placeholder_format.idx
            types = {
                0: "title",
                1: "body",
                10: "title",
                11: "subtitle",
                12: "subtitle",
                13: "body",
            }
            return types.get(ph_idx, f"placeholder_{ph_idx}")
        except Exception:
            return None

    def _check_hyperlinks(self, shape) -> bool:
        """检查是否包含超链接"""
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            return True
        except Exception:
            pass
        return False

    def _check_animations(self, slide) -> bool:
        """检查幻灯片是否有动画"""
        try:
            anim_el = slide._element.find(qn("p:timing"))
            return anim_el is not None
        except Exception:
            return False

    def _check_media(self, slide) -> bool:
        """检查幻灯片是否有音视频"""
        try:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    return True
        except Exception:
            pass
        return False

    def _compress_image(self, image_bytes: bytes) -> Optional[str]:
        """压缩图片为缩略 base64"""
        try:
            from PIL import Image
            img = Image.open(BytesIO(image_bytes))
            img.thumbnail((self.max_image_size, self.max_image_size))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            buf = BytesIO()
            img.save(buf, format="JPEG", quality=75)
            return base64.b64encode(buf.getvalue()).decode("ascii")
        except Exception as e:
            logger.warning(f"图片压缩失败: {e}")
            return None

    def _extract_group_text(self, shape) -> Optional[str]:
        """递归提取 Group 内所有文本（只读，不用于写回）"""
        texts = []
        try:
            for child_shape in shape.shapes:
                if child_shape.has_text_frame:
                    text = child_shape.text_frame.text.strip()
                    if text and not self._is_placeholder_text(text):
                        texts.append(text)
                elif child_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    sub_text = self._extract_group_text(child_shape)
                    if sub_text:
                        texts.append(sub_text)
        except Exception as e:
            logger.debug(f"Group 文本提取失败: {e}")
        return " | ".join(texts) if texts else None

    def _extract_doc_title(self, slides: list[ParsedSlide]) -> Optional[str]:
        """从第一页提取文档标题"""
        if not slides:
            return None
        for elem in slides[0].elements:
            if elem.is_title and elem.plain_text:
                return elem.plain_text.strip()
        for elem in slides[0].elements:
            if elem.element_type == ElementType.TEXT_BOX and elem.plain_text:
                return elem.plain_text.strip()[:50]
        return None
