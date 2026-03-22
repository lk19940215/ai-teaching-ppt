"""
样式应用器 — 将 AI 返回的 StyleHints 应用到 PPTX Shape

支持的样式属性（python-pptx 原生 API）：
  - bold / italic / underline
  - font_size_pt (磅值)
  - font_color (#RRGGBB)
  - font_name
  - alignment (left / center / right / justify)

设计为可插拔模块：pptx_writer.py 按需调用，
不影响核心文本替换逻辑。
"""

import logging
from typing import Optional

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from .models import StyleHints

logger = logging.getLogger(__name__)

_ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


class StyleApplicator:
    """将 StyleHints 应用到 python-pptx Shape 上"""

    def apply(self, shape, hints: StyleHints) -> int:
        """
        将样式提示应用到 Shape 的所有 Run。

        Returns:
            实际应用的属性数量
        """
        if not hasattr(shape, "text_frame"):
            logger.debug(f"Shape 无 text_frame，跳过样式应用")
            return 0

        applied = 0
        tf = shape.text_frame

        if hints.alignment:
            applied += self._apply_alignment(tf, hints.alignment)

        for para in tf.paragraphs:
            for run in para.runs:
                applied += self._apply_run_style(run, hints)

        if applied > 0:
            logger.debug(f"样式应用完成: {applied} 项属性")
        return applied

    def _apply_alignment(self, text_frame, alignment: str) -> int:
        """应用段落对齐"""
        pp_align = _ALIGNMENT_MAP.get(alignment)
        if pp_align is None:
            logger.warning(f"不支持的对齐方式: {alignment}")
            return 0

        for para in text_frame.paragraphs:
            para.alignment = pp_align
        return 1

    def _apply_run_style(self, run, hints: StyleHints) -> int:
        """应用 Run 级别的字体样式"""
        applied = 0
        font = run.font

        if hints.bold is not None:
            font.bold = hints.bold
            applied += 1

        if hints.italic is not None:
            font.italic = hints.italic
            applied += 1

        if hints.underline is not None:
            font.underline = hints.underline
            applied += 1

        if hints.font_size_pt is not None:
            try:
                font.size = Pt(hints.font_size_pt)
                applied += 1
            except Exception as e:
                logger.warning(f"设置字号失败: {e}")

        if hints.font_color is not None:
            applied += self._apply_color(font, hints.font_color)

        if hints.font_name is not None:
            font.name = hints.font_name
            applied += 1

        return applied

    def _apply_color(self, font, color_hex: str) -> int:
        """应用字体颜色（#RRGGBB 格式）"""
        try:
            hex_str = color_hex.lstrip("#")
            if len(hex_str) != 6:
                logger.warning(f"无效的颜色格式: {color_hex}")
                return 0
            r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
            font.color.rgb = RGBColor(r, g, b)
            return 1
        except Exception as e:
            logger.warning(f"设置颜色失败: {e}")
            return 0
