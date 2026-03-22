"""
PPTX 写回器

将修改指令 (SlideModification) 应用到原始 PPTX 文件，
在 Run 级别替换文本以保留所有格式属性。

插件集成：
  - StyleApplicator:     应用 AI 返回的样式提示 (bold, color, font_size 等)
  - AnimationApplicator: 应用 AI 返回的动画提示 (可插拔，默认关闭)

技术文档: docs/technical-spec.md §5.4, §8
"""

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
from uuid import uuid4

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu

from .models import SlideModification, TextModification, TableCellModification, SlideSelector, NewSlideContent
from .style_applicator import StyleApplicator
from .animation_applicator import AnimationApplicator

logger = logging.getLogger(__name__)


@dataclass
class ApplyResult:
    """apply() 的返回值，包含输出路径和页面索引元数据"""
    output_path: Path
    modified_indices: list[int] = field(default_factory=list)
    new_slide_indices: list[int] = field(default_factory=list)


class PPTXWriter:
    """PPTX 写回器：原始 PPTX + 修改指令 → 新 PPTX 文件

    通过构造函数注入 StyleApplicator 和 AnimationApplicator，
    实现样式/动画处理的可插拔。
    """

    def __init__(
        self,
        style_applicator: Optional[StyleApplicator] = None,
        animation_applicator: Optional[AnimationApplicator] = None,
    ):
        self._style = style_applicator or StyleApplicator()
        self._anim = animation_applicator or AnimationApplicator(enabled=False)

    def apply(
        self,
        source_path: Path,
        modifications: list[SlideModification],
        output_dir: Path,
    ) -> ApplyResult:
        """
        应用修改并产出**干净的结果 PPTX**，只包含 AI 处理过的页面。

        策略：单 Presentation 就地修改 + 删除未处理页面。
        避免双实例 deepcopy 导致 python-pptx 序列化时丢失文本修改。
        """
        prs = Presentation(str(source_path))

        total_applied = 0
        total_skipped = 0
        total_styled = 0
        total_animated = 0

        mod_slide_indices: set[int] = set()
        new_slide_mods: list[SlideModification] = []

        sorted_mods = sorted(modifications, key=lambda m: (0 if m.is_new_slide else 1))

        for slide_mod in sorted_mods:
            if slide_mod.is_new_slide and slide_mod.new_slide_content:
                new_slide_mods.append(slide_mod)
                continue

            if slide_mod.slide_index >= len(prs.slides):
                logger.warning(
                    f"slide_index {slide_mod.slide_index} 超出范围 "
                    f"(共 {len(prs.slides)} 页), 跳过"
                )
                continue

            slide = prs.slides[slide_mod.slide_index]
            shapes = list(slide.shapes)
            mod_slide_indices.add(slide_mod.slide_index)

            for text_mod in slide_mod.text_modifications:
                applied = self._apply_text_modification(shapes, text_mod, slide_mod.slide_index)
                if applied:
                    total_applied += 1
                    if text_mod.style_hints:
                        shape = shapes[text_mod.shape_index]
                        n = self._style.apply(shape, text_mod.style_hints)
                        total_styled += n
                else:
                    total_skipped += 1

            for table_mod in slide_mod.table_modifications:
                applied = self._apply_table_modification(shapes, table_mod, slide_mod.slide_index)
                if applied:
                    total_applied += 1
                else:
                    total_skipped += 1

            if slide_mod.animation_hints:
                n = self._anim.apply(slide, shapes, slide_mod.animation_hints)
                total_animated += n

        for new_mod in new_slide_mods:
            if new_mod.new_slide_content:
                self._create_new_slide(prs, new_mod.new_slide_content)
                mod_slide_indices.add(len(prs.slides) - 1)
                new_slide = prs.slides[-1]
                if new_mod.animation_hints:
                    shapes = list(new_slide.shapes)
                    self._anim.apply(new_slide, shapes, new_mod.animation_hints)
                    total_animated += len(new_mod.animation_hints)

        indices_to_remove = sorted(
            [i for i in range(len(prs.slides)) if i not in mod_slide_indices],
            reverse=True,
        )
        for idx in indices_to_remove:
            rId = prs.slides._sldIdLst[idx].get(qn("r:id"))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

        result_page_count = len(prs.slides)

        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"ai_result_{uuid4().hex[:8]}.pptx"
        prs.save(str(output_path))

        all_indices = list(range(result_page_count))
        logger.info(
            f"写回完成: {output_path.name}, "
            f"结果 {result_page_count} 页, "
            f"文本 {total_applied} 项(跳过 {total_skipped}), "
            f"样式 {total_styled} 项, 动画 {total_animated} 项"
        )
        return ApplyResult(
            output_path=output_path,
            modified_indices=all_indices,
            new_slide_indices=[],
        )

    def _apply_text_modification(
        self, shapes: list, mod: TextModification, slide_idx: int
    ) -> bool:
        """应用单个文本修改"""
        if mod.shape_index >= len(shapes):
            logger.warning(
                f"第{slide_idx}页 shape_index {mod.shape_index} 超出范围 "
                f"(共 {len(shapes)} 个 shape), 跳过"
            )
            return False

        shape = shapes[mod.shape_index]

        if not shape.has_text_frame:
            logger.warning(
                f"第{slide_idx}页 shape_index {mod.shape_index} 不是文本框, 跳过"
            )
            return False

        current_text = shape.text_frame.text.strip()
        expected = mod.original_text.strip()
        if current_text != expected:
            similarity = self._text_similarity(current_text, expected)
            if similarity < 0.5:
                logger.warning(
                    f"第{slide_idx}页 shape_index {mod.shape_index} 文本不匹配 "
                    f"(相似度 {similarity:.0%}), 跳过. "
                    f"期望: {expected[:50]}... 实际: {current_text[:50]}..."
                )
                return False
            logger.info(
                f"第{slide_idx}页 shape_index {mod.shape_index} 文本部分匹配 "
                f"(相似度 {similarity:.0%}), 继续写入"
            )

        self._replace_text_preserve_format(shape, mod.new_text)

        after_text = shape.text_frame.text.strip()[:60]
        logger.info(
            f"第{slide_idx}页 shape_{mod.shape_index} 替换完成: "
            f"'{expected[:40]}' → '{after_text}'"
        )
        return True

    def _apply_table_modification(
        self, shapes: list, mod: TableCellModification, slide_idx: int
    ) -> bool:
        """应用单个表格单元格修改"""
        if mod.shape_index >= len(shapes):
            logger.warning(
                f"第{slide_idx}页 shape_index {mod.shape_index} 超出范围, 跳过"
            )
            return False

        shape = shapes[mod.shape_index]

        if not shape.has_table:
            logger.warning(
                f"第{slide_idx}页 shape_index {mod.shape_index} 不是表格, 跳过"
            )
            return False

        table = shape.table
        if mod.row >= len(table.rows) or mod.col >= len(table.columns):
            logger.warning(
                f"第{slide_idx}页 表格 shape_{mod.shape_index} "
                f"行列索引 ({mod.row},{mod.col}) 超出范围 "
                f"({len(table.rows)}行×{len(table.columns)}列), 跳过"
            )
            return False

        cell = table.cell(mod.row, mod.col)
        self._replace_cell_text(cell, mod.new_text)
        logger.debug(
            f"第{slide_idx}页 表格 shape_{mod.shape_index} "
            f"({mod.row},{mod.col}): '{mod.original_text[:20]}' → '{mod.new_text[:20]}'"
        )
        return True

    def _replace_text_preserve_format(self, shape, new_text: str):
        """
        Run 级别文本替换 — 保留格式的核心算法。

        策略:
          1. 将新文本按换行符分割为段落
          2. 对于每个新段落:
             - 如果对应位置有旧段落: 替换第一个 Run 文本，清空后续 Run
             - 如果没有旧段落: 追加新段落，继承最后一段格式
          3. 多余的旧段落: 清空文本（保留段落结构避免 XML 异常）
        """
        tf = shape.text_frame
        new_paras = new_text.split("\n")
        old_paras = list(tf.paragraphs)

        for i, new_para_text in enumerate(new_paras):
            if i < len(old_paras):
                self._replace_paragraph_text(old_paras[i], new_para_text)
            else:
                new_para = tf.add_paragraph()
                new_para.text = new_para_text
                if old_paras:
                    self._copy_paragraph_format(old_paras[-1], new_para)

        for i in range(len(new_paras), len(old_paras)):
            self._clear_paragraph(old_paras[i])

        self._enable_auto_fit(tf)

    @staticmethod
    def _enable_auto_fit(text_frame):
        """启用文本框自动缩放，防止文字超出边界"""
        try:
            body_pr = text_frame._txBody.find(qn("a:bodyPr"))
            if body_pr is None:
                return
            for child in list(body_pr):
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag in ("noAutofit", "spAutoFit"):
                    body_pr.remove(child)
            from lxml import etree
            etree.SubElement(body_pr, qn("a:normAutofit"), attrib={"fontScale": "100000"})
        except Exception:
            pass

    def _replace_paragraph_text(self, para, new_text: str):
        """替换段落文本，保留第一个 Run 的格式"""
        runs = list(para.runs)
        if runs:
            runs[0].text = new_text
            for run in runs[1:]:
                run.text = ""
        else:
            para.text = new_text

    def _clear_paragraph(self, para):
        """清空段落文本（保留段落节点）"""
        for run in para.runs:
            run.text = ""

    def _replace_cell_text(self, cell, new_text: str):
        """替换表格单元格文本，按换行分段，保留格式"""
        tf = cell.text_frame
        new_paras = new_text.split("\n")
        old_paras = list(tf.paragraphs)

        for i, new_para_text in enumerate(new_paras):
            if i < len(old_paras):
                self._replace_paragraph_text(old_paras[i], new_para_text)
            else:
                new_para = tf.add_paragraph()
                new_para.text = new_para_text
                if old_paras:
                    self._copy_paragraph_format(old_paras[-1], new_para)

        for i in range(len(new_paras), len(old_paras)):
            self._clear_paragraph(old_paras[i])

        self._enable_auto_fit(tf)

    def _copy_paragraph_format(self, source_para, target_para):
        """复制段落格式"""
        try:
            if source_para.alignment is not None:
                target_para.alignment = source_para.alignment
        except Exception:
            pass

        try:
            if source_para.level is not None:
                target_para.level = source_para.level
        except Exception:
            pass

        try:
            src_runs = list(source_para.runs)
            tgt_runs = list(target_para.runs)
            if src_runs and tgt_runs:
                sf = src_runs[0].font
                tf = tgt_runs[0].font
                if sf.name:
                    tf.name = sf.name
                if sf.size:
                    tf.size = sf.size
                if sf.bold is not None:
                    tf.bold = sf.bold
                if sf.italic is not None:
                    tf.italic = sf.italic
                try:
                    if sf.color and sf.color.rgb:
                        tf.color.rgb = sf.color.rgb
                except Exception:
                    pass
        except Exception:
            pass

    def _create_new_slide(self, prs, content: NewSlideContent):
        """根据 AI 返回的 NewSlideContent 创建全新页面"""
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN

        layout_map = {
            "blank": 6,
            "title_and_content": 1,
            "title_only": 5,
        }
        layout_idx = layout_map.get(content.layout_hint, 1)
        try:
            layout = prs.slide_layouts[layout_idx]
        except IndexError:
            layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        if content.title:
            title_set = False
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 0:
                    shape.text = content.title
                    for para in shape.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.LEFT
                        for run in para.runs:
                            run.font.size = Pt(32)
                            run.font.bold = True
                    self._enable_auto_fit(shape.text_frame)
                    title_set = True
                    break
            if not title_set:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content.title
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(32)
                    run.font.bold = True

        if content.body_texts:
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break

            if not body_shape or not body_shape.has_text_frame:
                body_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.8), Inches(9), Inches(5.2)
                )

            tf = body_shape.text_frame
            tf.word_wrap = True

            first_para = True
            for block_idx, block in enumerate(content.body_texts):
                lines = block.split("\n")
                for line_idx, line in enumerate(lines):
                    if first_para:
                        para = tf.paragraphs[0]
                        first_para = False
                    else:
                        para = tf.add_paragraph()

                    para.text = line
                    para.alignment = PP_ALIGN.LEFT
                    is_heading = line_idx == 0 and len(lines) > 1 and ":" in line
                    for run in para.runs:
                        if is_heading:
                            run.font.size = Pt(22)
                            run.font.bold = True
                        else:
                            run.font.size = Pt(18)
                    if is_heading:
                        para.space_before = Pt(16) if block_idx > 0 else Pt(4)
                        para.space_after = Pt(4)
                    else:
                        para.space_after = Pt(6)
                        para.space_before = Pt(2)

            self._enable_auto_fit(tf)

        logger.info(f"创建新页面: title='{content.title[:30]}', body={len(content.body_texts)}段")

    def compose(
        self,
        source_files: dict[str, Path],
        selections: list[SlideSelector],
        output_dir: Path,
    ) -> Path:
        """
        从多个 PPTX 中选择页面，组合为新 PPTX。

        Args:
            source_files: {"ppt_a": Path, "ppt_b": Path, ...}
            selections: 按顺序选择的页面列表
            output_dir: 输出目录

        Returns:
            新 PPTX 文件路径
        """
        from copy import deepcopy
        from lxml import etree

        if not selections:
            raise ValueError("至少需要选择一页")

        first_source = selections[0].source
        first_path = source_files.get(first_source)
        if not first_path:
            raise ValueError(f"找不到源文件: {first_source}")

        base_prs = Presentation(str(first_path))

        presentations = {}
        for key, path in source_files.items():
            presentations[key] = Presentation(str(path))

        existing_slides = list(base_prs.slides)
        for slide in existing_slides:
            rId = base_prs.slides._sldIdLst[-1].get(qn("r:id"))
            base_prs.part.drop_rel(rId)
            base_prs.slides._sldIdLst.remove(base_prs.slides._sldIdLst[-1])

        for sel in selections:
            src_prs = presentations.get(sel.source)
            if not src_prs:
                logger.warning(f"源 {sel.source} 不存在, 跳过")
                continue
            if sel.slide_index >= len(src_prs.slides):
                logger.warning(
                    f"源 {sel.source} 页码 {sel.slide_index} 超出范围, 跳过"
                )
                continue

            src_slide = src_prs.slides[sel.slide_index]
            self._copy_slide(base_prs, src_slide, src_prs)

        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"composed_{uuid4().hex[:8]}.pptx"
        base_prs.save(str(output_path))

        logger.info(f"组合完成: {output_path.name}, {len(selections)} 页")
        return output_path

    def _copy_slide(self, dest_prs, src_slide, src_prs):
        """将一页幻灯片从源复制到目标演示文稿，尽量匹配源布局"""
        from copy import deepcopy
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        src_layout_name = src_slide.slide_layout.name if src_slide.slide_layout else None
        slide_layout = None
        if src_layout_name:
            for layout in dest_prs.slide_layouts:
                if layout.name == src_layout_name:
                    slide_layout = layout
                    break
        if slide_layout is None:
            slide_layout = dest_prs.slide_layouts[min(6, len(dest_prs.slide_layouts) - 1)]

        new_slide = dest_prs.slides.add_slide(slide_layout)

        for shape in new_slide.shapes:
            sp = shape._element
            sp.getparent().remove(sp)

        for shape in src_slide.shapes:
            el = deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)

        if src_slide.background and src_slide.background._element is not None:
            try:
                bg_elem = deepcopy(src_slide.background._element)
                existing_bg = new_slide._element.find(qn("p:bg"))
                if existing_bg is not None:
                    new_slide._element.remove(existing_bg)
                cs_sld = new_slide._element
                cs_sld.insert(0, bg_elem)
            except Exception:
                pass

        for rel in src_slide.part.rels.values():
            if rel.is_external:
                new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                try:
                    new_slide.part.rels.get_or_add(
                        rel.rId, rel.reltype, rel.target_part
                    )
                except Exception as e:
                    logger.debug(f"复制关系失败 {rel.reltype}: {e}")

    def _text_similarity(self, a: str, b: str) -> float:
        """简单的文本相似度（字符级 Jaccard）"""
        if not a and not b:
            return 1.0
        if not a or not b:
            return 0.0
        set_a = set(a)
        set_b = set(b)
        intersection = set_a & set_b
        union = set_a | set_b
        return len(intersection) / len(union) if union else 0.0
