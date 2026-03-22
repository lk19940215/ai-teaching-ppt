"""
课文内容生成 PPT API

端点:
  POST /generate/outline    输入课文内容，AI 生成幻灯片大纲
  POST /generate/ppt        根据大纲生成 PPTX 文件
"""

import uuid
import logging
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..config import settings
from ..ai.llm_client import LLMClient

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/generate", tags=["Generate"])


class OutlineRequest(BaseModel):
    content: str = Field(..., min_length=10, description="课文内容")
    subject: Optional[str] = None
    grade: Optional[str] = None
    slide_count: Optional[int] = Field(None, ge=3, le=30)
    provider: str = "deepseek"
    api_key: str
    base_url: Optional[str] = None
    model: Optional[str] = None
    temperature: float = 0.3
    max_tokens: int = 4000


class SlideOutline(BaseModel):
    title: str
    content: list[str]
    notes: str = ""


class OutlineResponse(BaseModel):
    title: str
    slides: list[SlideOutline]


class GeneratePptRequest(BaseModel):
    title: str
    slides: list[SlideOutline]
    subject: Optional[str] = None
    grade: Optional[str] = None


class GeneratePptResponse(BaseModel):
    success: bool
    download_url: str = ""
    file_name: str = ""
    slide_count: int = 0
    error: str = ""


OUTLINE_SYSTEM_PROMPT = """你是一位具有 10 年经验的教学设计师，精通 ADDIE 教学设计模型和布鲁姆教育目标分类学。
请根据用户提供的课文内容，按照标准教学流程设计一份 PPT 课件大纲。

## 教学流程结构（必须包含以下环节）

1. **标题页**（自动生成，不需要在 slides 中）
2. **学习目标**：用"能够..."句式列出 3-4 个具体、可测量的学习目标
3. **课堂导入**（Warm-up）：通过提问、图片、情境等激发兴趣，2-3 分钟
4. **新知呈现**（Presentation）：分步呈现核心知识点，每页聚焦一个主题
5. **操练巩固**（Practice）：设计互动练习（对话练习、填空、判断、小组活动等）
6. **拓展应用**（Production）：设计一个综合性输出活动（角色扮演、讨论、写作等）
7. **课堂小结**（Summary）：回顾重点，布置作业

## 设计原则

- 每页 2-5 个要点，语言简洁、适合投影展示
- 知识点由浅入深、循序渐进
- 注重互动性：包含提问、讨论、活动指令
- 符合目标学段学生的认知水平
- 如果是语言类学科，注意听说读写的均衡
- 备注(notes)写给教师的教学提示，如时间分配、教学策略、过渡语等

## 输出格式（严格 JSON）

{
  "title": "整份课件的标题",
  "slides": [
    {
      "title": "幻灯片标题",
      "content": ["要点1", "要点2", "要点3"],
      "notes": "教师备注：教学策略和时间建议"
    }
  ]
}"""


@router.post("/outline", response_model=OutlineResponse)
async def generate_outline(req: OutlineRequest):
    """根据课文内容生成 PPT 大纲"""
    try:
        llm = LLMClient(
            provider=req.provider,
            api_key=req.api_key,
            base_url=req.base_url,
            model=req.model,
            temperature=req.temperature,
            max_tokens=req.max_tokens,
        )

        user_prompt = f"请根据以下课文内容生成教学 PPT 大纲。\n\n"
        if req.subject:
            user_prompt += f"学科：{req.subject}\n"
        if req.grade:
            user_prompt += f"年级：{req.grade}\n"
        if req.slide_count:
            user_prompt += f"建议幻灯片数量：约 {req.slide_count} 页\n"
        user_prompt += f"\n课文内容：\n{req.content}"

        messages = [
            {"role": "system", "content": OUTLINE_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt},
        ]

        data = llm.chat_json(messages)
        title = data.get("title", "教学课件")
        slides_data = data.get("slides", [])

        if not slides_data:
            raise HTTPException(400, "AI 未能生成有效的大纲")

        slides = []
        for s in slides_data:
            slides.append(SlideOutline(
                title=s.get("title", ""),
                content=s.get("content", []),
                notes=s.get("notes", ""),
            ))

        logger.info(f"[generate_outline] 成功生成 {len(slides)} 页大纲: {title}")
        return OutlineResponse(title=title, slides=slides)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[generate_outline] 失败: {e}")
        raise HTTPException(500, f"生成大纲失败: {str(e)}")


@router.post("/ppt", response_model=GeneratePptResponse)
async def generate_ppt(req: GeneratePptRequest):
    """根据大纲生成 PPTX 文件"""
    try:
        gen_id = uuid.uuid4().hex[:8]
        output_dir = Path(settings.UPLOAD_DIR) / "generated"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"generated_{gen_id}.pptx"

        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 10 inches
        prs.slide_height = Emu(6858000)   # 7.5 inches

        _add_title_slide(prs, req.title, req.subject, req.grade)
        for slide_data in req.slides:
            _add_content_slide(prs, slide_data)

        prs.save(str(output_path))

        download_url = f"/api/v1/generate/download/{gen_id}"
        file_name = f"{req.title}.pptx"

        logger.info(f"[generate_ppt] 成功生成 {len(req.slides)+1} 页 PPT: {output_path}")
        return GeneratePptResponse(
            success=True,
            download_url=download_url,
            file_name=file_name,
            slide_count=len(req.slides) + 1,
        )

    except Exception as e:
        logger.error(f"[generate_ppt] 失败: {e}")
        return GeneratePptResponse(success=False, error=str(e))


@router.get("/download/{gen_id}")
async def download_generated(gen_id: str):
    """下载生成的 PPT"""
    output_path = Path(settings.UPLOAD_DIR) / "generated" / f"generated_{gen_id}.pptx"
    if not output_path.exists():
        raise HTTPException(404, "文件不存在")
    return FileResponse(
        path=str(output_path),
        filename=f"generated_{gen_id}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ---- PPT 构建辅助函数 ----

THEME_PRIMARY = RGBColor(0x33, 0x56, 0x9B)    # 深蓝
THEME_ACCENT = RGBColor(0x4F, 0x81, 0xBD)     # 蓝
THEME_BG = RGBColor(0xF2, 0xF4, 0xF8)         # 浅灰蓝
THEME_TEXT = RGBColor(0x33, 0x33, 0x33)        # 深灰
THEME_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_title_slide(prs: Presentation, title: str, subject: str | None, grade: str | None):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = THEME_PRIMARY

    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = THEME_WHITE
    p.alignment = PP_ALIGN.CENTER

    if subject or grade:
        sub_text = " | ".join(filter(None, [subject, grade]))
        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(0.6)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = sub_text
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xCC, 0xD5, 0xE5)
        p2.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, slide_data: SlideOutline):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    header_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(1.2)  # MSO_SHAPE.RECTANGLE
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = THEME_PRIMARY
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.2), Inches(8.4), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = THEME_WHITE
    p.alignment = PP_ALIGN.LEFT

    if slide_data.content:
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(8.4), Inches(5)
        )
        tf_content = content_box.text_frame
        tf_content.word_wrap = True

        for i, point in enumerate(slide_data.content):
            if i == 0:
                p = tf_content.paragraphs[0]
            else:
                p = tf_content.add_paragraph()
            p.text = f"  {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = THEME_TEXT
            p.space_after = Pt(12)
            p.level = 0

            bullet = p._pPr
            if bullet is None:
                from pptx.oxml.ns import qn
                from lxml import etree
                p._p.insert(0, etree.SubElement(p._p, qn('a:pPr')))
                bullet = p._pPr

    if slide_data.notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data.notes
