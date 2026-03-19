from fastapi import APIRouter, HTTPException, Body, UploadFile, File, Form, Query, Request, BackgroundTasks
from starlette.datastructures import UploadFile as StarletteUploadFile
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from pathlib import Path
from typing import Optional, AsyncGenerator, List, Dict, Any
import logging
import traceback as traceback_module
import uuid
import aiofiles
import asyncio
import json
import io
import tempfile
import hashlib
import time
import psutil
import os
import zipfile

from ..services.ppt_generator import get_ppt_generator, generate_single_slide_pptx
from ..services.ppt_to_image import convert_single_slide_to_image
from ..config import settings

router = APIRouter(prefix=settings.API_V1_STR, tags=["ppt"])

logger = logging.getLogger(__name__)

# PPT 解析缓存
_ppt_parse_cache: Dict[str, Dict[str, Any]] = {}
_PPT_CACHE_MAX_SIZE = 100
_PPT_CACHE_TTL = 1800
_MEMORY_WARNING_THRESHOLD = 0.85
_MEMORY_CRITICAL_THRESHOLD = 0.95

def _get_memory_usage() -> float:
    """获取当前进程内存使用率"""
    try:
        process = psutil.Process(os.getpid())
        mem_info = process.memory_info()
        system_mem = psutil.virtual_memory().total
        return mem_info.rss / system_mem
    except Exception:
        return 0.0

def _compute_file_hash(content: bytes) -> str:
    """计算文件内容的 SHA256 哈希值"""
    return hashlib.sha256(content).hexdigest()

def safe_join_list(items: List[Any], separator: str = '\n') -> str:
    """
    安全连接列表元素，处理混合类型（str/dict）。

    Args:
        items: 可能包含字符串或字典的列表
        separator: 连接符，默认换行符

    Returns:
        连接后的字符串

    Examples:
        >>> safe_join_list(['a', 'b'])
        'a\\nb'
        >>> safe_join_list([{'text': 'a'}, {'point': 'b'}])
        'a\\nb'
        >>> safe_join_list([])
        ''
    """
    if not items:
        return ''

    result = []
    for item in items:
        if isinstance(item, str):
            result.append(item)
        elif isinstance(item, dict):
            # 尝试提取常见的文本字段：text, content, point, title, value
            text = (item.get('text') or
                    item.get('content') or
                    item.get('point') or
                    item.get('title') or
                    item.get('value') or
                    str(item))
            result.append(text)
        else:
            result.append(str(item))

    return separator.join(result)


def test_json_serialization(data: Any, context: str = "data") -> Dict[str, Any]:
    """
    测试 JSON 序列化并返回详细的调试信息。

    Args:
        data: 要测试的数据
        context: 上下文描述（用于日志）

    Returns:
        {
            "success": bool,
            "error": str | None,
            "problematic_fields": List[str],
            "json_length": int | None,
            "duration_ms": float
        }
    """
    start_time = time.time()
    result = {
        "success": False,
        "error": None,
        "problematic_fields": [],
        "json_length": None,
        "duration_ms": 0
    }

    try:
        # 尝试整体序列化
        json_str = json.dumps(data, ensure_ascii=False, default=str)
        result["success"] = True
        result["json_length"] = len(json_str)
        logger.info(f"[JSON序列化] {context}: 成功，长度 {len(json_str)} 字符")
    except TypeError as e:
        result["error"] = str(e)
        logger.error(f"[JSON序列化] {context}: 失败 - {e}")

        # 逐字段测试，定位问题字段
        if isinstance(data, dict):
            for key, value in data.items():
                try:
                    json.dumps({key: value}, ensure_ascii=False, default=str)
                except TypeError:
                    result["problematic_fields"].append(f"{key}: {type(value).__name__}")
                    logger.error(f"[JSON序列化] {context}: 字段 '{key}' 无法序列化，类型={type(value).__name__}, 值={repr(value)[:200]}")

                    # 如果值也是字典，继续深入测试
                    if isinstance(value, dict):
                        for sub_key, sub_value in value.items():
                            try:
                                json.dumps({sub_key: sub_value}, ensure_ascii=False, default=str)
                            except TypeError:
                                result["problematic_fields"].append(f"{key}.{sub_key}: {type(sub_value).__name__}")
                                logger.error(f"[JSON序列化] {context}: 子字段 '{key}.{sub_key}' 无法序列化，类型={type(sub_value).__name__}")
        elif isinstance(data, list):
            for i, item in enumerate(data):
                try:
                    json.dumps(item, ensure_ascii=False, default=str)
                except TypeError:
                    result["problematic_fields"].append(f"[{i}]: {type(item).__name__}")
                    logger.error(f"[JSON序列化] {context}: 列表项 [{i}] 无法序列化，类型={type(item).__name__}")
    except Exception as e:
        result["error"] = f"Unexpected error: {e}"
        logger.error(f"[JSON序列化] {context}: 意外错误 - {e}")

    result["duration_ms"] = round((time.time() - start_time) * 1000, 2)
    return result


def _cleanup_cache():
    """清理过期缓存并执行 LRU 淘汰"""
    current_time = time.time()
    keys_to_remove = []
    for file_hash, cache_entry in list(_ppt_parse_cache.items()):
        if current_time - cache_entry["timestamp"] > _PPT_CACHE_TTL:
            keys_to_remove.append(file_hash)
    if len(_ppt_parse_cache) - len(keys_to_remove) > _PPT_CACHE_MAX_SIZE:
        sorted_entries = sorted(
            [(k, v) for k, v in _ppt_parse_cache.items() if k not in keys_to_remove],
            key=lambda x: x[1]["timestamp"]
        )
        keys_to_remove.extend([k for k, _ in sorted_entries[:_PPT_CACHE_MAX_SIZE // 2]])
    for key in keys_to_remove:
        del _ppt_parse_cache[key]

def _get_from_cache(file_hash: str) -> Optional[Dict[str, Any]]:
    """从缓存中获取解析结果"""
    if file_hash not in _ppt_parse_cache:
        return None
    cache_entry = _ppt_parse_cache[file_hash]
    current_time = time.time()
    if current_time - cache_entry["timestamp"] > _PPT_CACHE_TTL:
        del _ppt_parse_cache[file_hash]
        return None
    return cache_entry["data"]

def _add_to_cache(file_hash: str, data: Dict[str, Any]):
    """添加解析结果到缓存"""
    _cleanup_cache()
    _ppt_parse_cache[file_hash] = {"data": data, "timestamp": time.time()}

def _fix_invalid_namespaces(file_path: Path) -> Path:
    """修复 PPTX 文件中无效的 XML 命名空间声明，支持所有 ns%d 变体"""
    import zipfile
    import re
    import shutil

    # 创建临时目录解压
    temp_dir = file_path.parent
    fix_dir = temp_dir / f"fixed_{uuid.uuid4().hex}"
    fix_dir.mkdir(parents=True, exist_ok=True)

    try:
        # 解压 PPTX
        with zipfile.ZipFile(file_path, 'r') as z:
            z.extractall(fix_dir)

        # 修复所有 XML 文件中的无效命名空间
        # 使用通用正则匹配所有 xmlns:ns\d+="%s" 或 xmlns:ns\d+='%s' 模式（支持单引号和双引号）
        def replace_ns(match):
            """替换命名空间，保持引号类型一致"""
            ns_num = match.group(1)
            quote = match.group(2)  # 捕获引号类型（单引号或双引号）
            return f'xmlns:ns{ns_num}={quote}http://schemas.openxmlformats.org/presentationml/2006/placeholder{quote}'

        ns_pattern = re.compile(r'xmlns:ns(\d+)=([\'"])%s\2')
        # 匹配包含 { 的无效占位符模式
        placeholder_pattern = re.compile(r'xmlns:([a-zA-Z0-9_]+)=["\']\{([^}]*)\}[^"\']*["\']')

        fixed_count = 0
        total_fixed = 0

        for xml_file in fix_dir.rglob("*.xml"):
            try:
                content = xml_file.read_text(encoding='utf-8')
                original_content = content

                # 替换模式 1: xmlns:ns%d="%s" 或 xmlns:ns%d='%s'（支持单引号和双引号）
                content = ns_pattern.sub(replace_ns, content)

                # 替换模式 2: xmlns:x="{.*}" 等包含 { 的占位符
                content = placeholder_pattern.sub(r'xmlns:\1="http://schemas.openxmlformats.org/presentationml/2006/\2"', content)

                if content != original_content:
                    xml_file.write_text(content, encoding='utf-8')
                    fixed_count += 1
                    logger.info(f"修复 XML 文件无效命名空间：{xml_file.relative_to(fix_dir)}")
                    # 统计修复的命名空间数量
                    total_fixed += len(ns_pattern.findall(original_content)) + len(placeholder_pattern.findall(original_content))

            except Exception as e:
                logger.warning(f"修复 XML 文件失败 {xml_file}: {e}")

        # 重新打包为 PPTX
        fixed_path = temp_dir / f"fixed_{uuid.uuid4().hex}_{file_path.name}"
        with zipfile.ZipFile(fixed_path, 'w', zipfile.ZIP_DEFLATED) as z:
            for root, dirs, files in os.walk(fix_dir):
                for f in files:
                    file_on_disk = Path(root) / f
                    arcname = file_on_disk.relative_to(fix_dir)
                    z.write(file_on_disk, arcname)

        logger.info(f"PPTX 命名空间修复完成：{file_path.name} -> {fixed_path.name} (修复 {fixed_count} 个文件，共 {total_fixed} 处命名空间)")
        return fixed_path

    except Exception as e:
        logger.error(f"修复 PPTX 命名空间失败：{e}")
        return file_path
    finally:
        # 清理临时目录
        try:
            shutil.rmtree(fix_dir)
        except Exception:
            pass

@router.post("/ppt/generate")
async def generate_ppt(
    content: dict = Body(..., embed=True),
    grade: str = Body("6", embed=True),
    subject: str = Body("general", embed=True),
    style: str = Body("simple", embed=True),
    file_name: Optional[str] = Body(None, embed=True)
):
    """
    生成 PPT 文件
    Args:
        content: PPT 内容数据
        grade: 年级
        subject: 学科
        style: PPT 风格
        file_name: 文件名（可选）
    Returns:
        生成的 PPT 文件下载链接
    """
    try:
        # 生成文件名
        if not file_name:
            file_name = f"{content.get('title', '教学 PPT')}_{uuid.uuid4().hex[:8]}.pptx"

        # 确保 .pptx 扩展名
        if not file_name.endswith(".pptx"):
            file_name += ".pptx"

        # 生成文件路径
        output_path = settings.UPLOAD_DIR / "generated" / file_name
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # 生成 PPT
        generator = get_ppt_generator()
        generator.generate(content, output_path, grade, style, subject)

        return JSONResponse(content={
            "success": True,
            "message": "PPT 生成成功",
            "download_url": f"/uploads/generated/{file_name}",
            "file_name": file_name
        })

    except Exception as e:
        logger.error(f"PPT 生成失败：{e}")
        raise HTTPException(status_code=500, detail=f"PPT 生成失败：{str(e)}")


@router.get("/ppt/download/{file_name}")
async def download_ppt(file_name: str):
    """
    下载 PPT 文件
    Args:
        file_name: 文件名
    Returns:
        PPT 文件下载
    """
    file_path = settings.UPLOAD_DIR / "generated" / file_name

    if not file_path.exists():
        raise HTTPException(status_code=404, detail="文件不存在")

    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=file_name
    )


@router.post("/ppt/generate-full")
async def generate_full_ppt(
    text_content: str = Body(..., embed=True),
    grade: str = Body("6", embed=True),
    subject: str = Body("math", embed=True),
    slide_count: int = Body(15, embed=True),
    chapter: Optional[str] = Body(None, embed=True),
    provider: str = Body("deepseek", embed=True),
    api_key: str = Body(..., embed=True),
    style: str = Body("simple", embed=True),
    session_id: Optional[str] = Body(None, embed=True),
    temperature: Optional[float] = Body(None, embed=True),
    max_output_tokens: Optional[int] = Body(None, embed=True),
    difficulty_level: str = Body("unified", embed=True),
):
    """
    完整生成 PPT（内容生成 + 文件生成）
    Args:
        text_content: 教学文本内容
        grade: 年级
        subject: 学科
        slide_count: 幻灯片数量
        chapter: 章节名称
        provider: LLM 服务商
        api_key: API Key
        style: PPT 风格
        temperature: 温度参数（可选，从前端传递）
        max_output_tokens: 最大输出 token 数（可选，从前端传递）
        difficulty_level: 教学层次（unified/basic/intermediate/advanced）
    Returns:
        生成结果和下载链接
    """
    try:
        # 步骤 1: 调用 LLM 生成内容
        from ..services.llm import get_llm_service
        from ..services.content_generator import get_content_generator

        # 使用前端传递的参数，或使用默认值
        llm_temperature = temperature if temperature is not None else 0.7
        llm_max_tokens = max_output_tokens if max_output_tokens is not None else 4000

        llm_service = get_llm_service(
            provider=provider,
            api_key=api_key,
            temperature=llm_temperature,
            max_tokens=llm_max_tokens
        )

        generator = get_content_generator(llm_service)

        # 英语学科使用专门的内容生成器
        if subject == "english":
            ppt_content = generator.generate_for_english(
                content=text_content,
                grade=grade,
                slide_count=slide_count,
                chapter=chapter,
                difficulty_level=difficulty_level
            )
        else:
            ppt_content = generator.generate(
                content=text_content,
                grade=grade,
                subject=subject,
                slide_count=slide_count,
                chapter=chapter,
                difficulty_level=difficulty_level
            )

        # 步骤 2: 生成 PPT 文件
        ppt_generator = get_ppt_generator()
        file_name = f"{ppt_content.get('title', '教学 PPT')}_{uuid.uuid4().hex[:8]}.pptx"
        output_path = settings.UPLOAD_DIR / "generated" / file_name
        output_path.parent.mkdir(parents=True, exist_ok=True)

        ppt_generator.generate(ppt_content, output_path, grade, style, subject)

        # 保存历史记录（如果提供了 session_id）
        history_record = None
        if session_id:
            from ..models.database import get_db_session
            from ..models.history_crud import add_generation_record

            db = get_db_session()
            try:
                history_record = add_generation_record(
                    db=db,
                    session_id=session_id,
                    title=ppt_content.get('title', '教学 PPT'),
                    grade=grade,
                    subject=subject,
                    style=style,
                    file_name=file_name,
                    file_path=str(output_path),
                    content_text={"text": text_content},
                    slide_count=slide_count,
                    chapter=chapter,
                    ppt_content=ppt_content,
                )
            finally:
                db.close()

        response_data = {
            "success": True,
            "message": "PPT 生成成功",
            "content": ppt_content,
            "download_url": f"/uploads/generated/{file_name}",
            "file_name": file_name
        }
        if history_record:
            response_data["history_id"] = history_record.id

        return JSONResponse(content=response_data)

    except ValueError as e:
        logger.error(f"参数错误：{e}")
        raise HTTPException(status_code=400, detail=str(e))
    except TimeoutError as e:
        logger.error(f"生成超时：{e}")
        raise HTTPException(status_code=408, detail="生成超时，请稍后重试")
    except Exception as e:
        logger.error(f"完整 PPT 生成失败：{e}")
        raise HTTPException(status_code=500, detail=f"PPT 生成失败：{str(e)}")


async def sse_generator(progress_queue: asyncio.Queue) -> AsyncGenerator[str, None]:
    """SSE 事件生成器（含心跳机制，防止客户端/自动化工具超时）"""
    import time
    HEARTBEAT_INTERVAL = 3.0
    try:
        while True:
            try:
                event = await asyncio.wait_for(progress_queue.get(), timeout=HEARTBEAT_INTERVAL)
                if event is None:
                    break
                yield f"data: {json.dumps(event, ensure_ascii=False)}\n\n"
            except asyncio.TimeoutError:
                heartbeat = {"type": "heartbeat", "timestamp": time.time()}
                yield f"data: {json.dumps(heartbeat)}\n\n"
    except asyncio.CancelledError:
        logger.info("SSE 连接被客户端关闭")
        raise


@router.post("/ppt/generate-stream")
async def generate_full_ppt_stream(
    text_content: str = Body(..., embed=True),
    grade: str = Body("6", embed=True),
    subject: str = Body("math", embed=True),
    slide_count: int = Body(15, embed=True),
    chapter: Optional[str] = Body(None, embed=True),
    provider: str = Body("deepseek", embed=True),
    api_key: str = Body(..., embed=True),
    style: str = Body("simple", embed=True),
    session_id: Optional[str] = Body(None, embed=True),
    temperature: Optional[float] = Body(None, embed=True),
    max_output_tokens: Optional[int] = Body(None, embed=True),
    difficulty_level: str = Body("unified", embed=True),
):
    """
    完整生成 PPT（SSE 流式响应）
    Args:
        text_content: 教学文本内容
        grade: 年级
        subject: 学科
        slide_count: 幻灯片数量
        chapter: 章节名称
        provider: LLM 服务商
        api_key: API Key
        style: PPT 风格
        temperature: 温度参数（可选，从前端传递）
        max_output_tokens: 最大输出 token 数（可选，从前端传递）
        difficulty_level: 教学层次（unified/basic/intermediate/advanced）
    Returns:
        SSE 流式响应，包含进度事件和最终结果
    """
    progress_queue: asyncio.Queue = asyncio.Queue()

    async def generate_in_background():
        """后台执行生成任务"""
        logger.info("后台生成任务启动")
        try:
            # 阶段 1: 分析内容 (10%)
            logger.info("发送第一个 SSE 事件：analyzing_content 10%")
            await progress_queue.put({
                "stage": "analyzing_content",
                "progress": 10,
                "message": "正在分析教材内容..."
            })

            # 步骤 1: 调用 LLM 生成内容
            from ..services.llm import get_llm_service
            from ..services.content_generator import get_content_generator

            llm_temperature = temperature if temperature is not None else 0.7
            llm_max_tokens = max_output_tokens if max_output_tokens is not None else 4000

            # 注意：base_url 和 model 参数不传递，让 LLMService 根据 provider 使用默认值
            # 各服务商默认值已在 llm.py 中定义：
            # - DeepSeek: https://api.deepseek.com, deepseek-chat
            # - OpenAI: https://api.openai.com/v1, gpt-4
            # - Claude: https://api.anthropic.com/v1, claude-3-opus-20240229
            # - GLM: https://open.bigmodel.cn/api/paas/v4, glm-4
            llm_service = get_llm_service(
                provider=provider,
                api_key=api_key,
                temperature=llm_temperature,
                max_tokens=llm_max_tokens
            )

            generator = get_content_generator(llm_service)

            # 阶段 2: 生成大纲 (30%)
            await progress_queue.put({
                "stage": "generating_outline",
                "progress": 30,
                "message": "正在调用 AI 生成 PPT 大纲..."
            })

            # 英语学科使用专门的内容生成器
            if subject == "english":
                ppt_content = generator.generate_for_english(
                    content=text_content,
                    grade=grade,
                    slide_count=slide_count,
                    chapter=chapter,
                    difficulty_level=difficulty_level
                )
            else:
                ppt_content = generator.generate(
                    content=text_content,
                    grade=grade,
                    subject=subject,
                    slide_count=slide_count,
                    chapter=chapter,
                    difficulty_level=difficulty_level
                )

            # 阶段 3: 构建幻灯片 (60%)
            await progress_queue.put({
                "stage": "building_slides",
                "progress": 60,
                "message": "正在构建幻灯片页面..."
            })

            # 步骤 2: 生成 PPT 文件
            ppt_generator = get_ppt_generator()
            file_name = f"{ppt_content.get('title', '教学 PPT')}_{uuid.uuid4().hex[:8]}.pptx"
            output_path = settings.UPLOAD_DIR / "generated" / file_name
            output_path.parent.mkdir(parents=True, exist_ok=True)

            ppt_generator.generate(ppt_content, output_path, grade, style, subject)

            # 阶段 4: 添加动画效果 (85%)
            await progress_queue.put({
                "stage": "adding_animations",
                "progress": 85,
                "message": "正在添加动画效果..."
            })

            # 保存历史记录（如果提供了 session_id）
            history_record = None
            if session_id:
                from ..models.database import get_db_session
                from ..models.history_crud import add_generation_record

                db = get_db_session()
                try:
                    history_record = add_generation_record(
                        db=db,
                        session_id=session_id,
                        title=ppt_content.get('title', '教学 PPT'),
                        grade=grade,
                        subject=subject,
                        style=style,
                        file_name=file_name,
                        file_path=str(output_path),
                        content_text={"text": text_content},
                        slide_count=slide_count,
                        chapter=chapter,
                        ppt_content=ppt_content,
                    )
                finally:
                    db.close()

            response_data = {
                "success": True,
                "message": "PPT 生成成功",
                "content": ppt_content,
                "download_url": f"/uploads/generated/{file_name}",
                "file_name": file_name
            }
            if history_record:
                response_data["history_id"] = history_record.id

            # 阶段 5: 完成 (100%)
            await progress_queue.put({
                "stage": "complete",
                "progress": 100,
                "message": "生成完成！",
                "result": response_data
            })

        except ValueError as e:
            logger.error(f"参数错误：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": f"参数错误：{str(e)}"
            })
        except TimeoutError as e:
            logger.error(f"生成超时：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": "生成超时，请稍后重试"
            })
        except Exception as e:
            logger.error(f"完整 PPT 生成失败：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": f"PPT 生成失败：{str(e)}"
            })
        finally:
            # 结束 SSE 流
            await progress_queue.put(None)

    # 启动后台任务
    task = asyncio.create_task(generate_in_background())

    # 让出控制权给事件循环，确保后台任务立即开始执行
    # 这样第一个 SSE 事件（analyzing_content 10%）会在返回响应前被放入队列
    await asyncio.sleep(0)

    # 返回 SSE 流
    return StreamingResponse(
        sse_generator(progress_queue),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@router.get("/ppt/generate-stream")
async def generate_full_ppt_stream_get(
    text_content: str,
    grade: str = "6",
    subject: str = "math",
    slide_count: int = 15,
    chapter: Optional[str] = None,
    provider: str = "deepseek",
    api_key: str = "",
    style: str = "simple",
    session_id: Optional[str] = None,
    temperature: Optional[float] = None,
    max_output_tokens: Optional[int] = None,
    difficulty_level: str = "unified",
):
    """
    完整生成 PPT（SSE 流式响应 - GET 版本，用于 EventSource）
    """
    logger.info(f"收到 SSE 生成请求：text_content={text_content[:50]}..., grade={grade}, subject={subject}")
    progress_queue: asyncio.Queue = asyncio.Queue()

    async def generate_in_background():
        """后台执行生成任务"""
        logger.info("后台生成任务启动")
        try:
            # 阶段 1: 分析内容 (10%)
            logger.info("发送第一个 SSE 事件：analyzing_content 10%")
            await progress_queue.put({
                "stage": "analyzing_content",
                "progress": 10,
                "message": "正在分析教材内容..."
            })

            # 步骤 1: 调用 LLM 生成内容
            from ..services.llm import get_llm_service
            from ..services.content_generator import get_content_generator

            llm_temperature = temperature if temperature is not None else 0.7
            llm_max_tokens = max_output_tokens if max_output_tokens is not None else 4000

            # 注意：base_url 和 model 参数不传递，让 LLMService 根据 provider 使用默认值
            # 各服务商默认值已在 llm.py 中定义：
            # - DeepSeek: https://api.deepseek.com, deepseek-chat
            # - OpenAI: https://api.openai.com/v1, gpt-4
            # - Claude: https://api.anthropic.com/v1, claude-3-opus-20240229
            # - GLM: https://open.bigmodel.cn/api/paas/v4, glm-4
            llm_service = get_llm_service(
                provider=provider,
                api_key=api_key,
                temperature=llm_temperature,
                max_tokens=llm_max_tokens
            )

            generator = get_content_generator(llm_service)

            # 阶段 2: 生成大纲 (30%)
            await progress_queue.put({
                "stage": "generating_outline",
                "progress": 30,
                "message": "正在调用 AI 生成 PPT 大纲..."
            })

            # 英语学科使用专门的内容生成器
            if subject == "english":
                ppt_content = generator.generate_for_english(
                    content=text_content,
                    grade=grade,
                    slide_count=slide_count,
                    chapter=chapter,
                    difficulty_level=difficulty_level
                )
            else:
                ppt_content = generator.generate(
                    content=text_content,
                    grade=grade,
                    subject=subject,
                    slide_count=slide_count,
                    chapter=chapter,
                    difficulty_level=difficulty_level
                )

            # 阶段 3: 构建幻灯片 (60%)
            await progress_queue.put({
                "stage": "building_slides",
                "progress": 60,
                "message": "正在构建幻灯片页面..."
            })

            # 步骤 2: 生成 PPT 文件
            ppt_generator = get_ppt_generator()
            file_name = f"{ppt_content.get('title', '教学 PPT')}_{uuid.uuid4().hex[:8]}.pptx"
            output_path = settings.UPLOAD_DIR / "generated" / file_name
            output_path.parent.mkdir(parents=True, exist_ok=True)

            ppt_generator.generate(ppt_content, output_path, grade, style, subject)

            # 阶段 4: 添加动画效果 (85%)
            await progress_queue.put({
                "stage": "adding_animations",
                "progress": 85,
                "message": "正在添加动画效果..."
            })

            # 保存历史记录（如果提供了 session_id）
            history_record = None
            if session_id:
                from ..models.database import get_db_session
                from ..models.history_crud import add_generation_record

                db = get_db_session()
                try:
                    history_record = add_generation_record(
                        db=db,
                        session_id=session_id,
                        title=ppt_content.get('title', '教学 PPT'),
                        grade=grade,
                        subject=subject,
                        style=style,
                        file_name=file_name,
                        file_path=str(output_path),
                        content_text={"text": text_content},
                        slide_count=slide_count,
                        chapter=chapter,
                        ppt_content=ppt_content,
                    )
                finally:
                    db.close()

            response_data = {
                "success": True,
                "message": "PPT 生成成功",
                "content": ppt_content,
                "download_url": f"/uploads/generated/{file_name}",
                "file_name": file_name
            }
            if history_record:
                response_data["history_id"] = history_record.id

            # 阶段 5: 完成 (100%)
            await progress_queue.put({
                "stage": "complete",
                "progress": 100,
                "message": "生成完成！",
                "result": response_data
            })

        except ValueError as e:
            logger.error(f"参数错误：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": f"参数错误：{str(e)}"
            })
        except TimeoutError as e:
            logger.error(f"生成超时：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": "生成超时，请稍后重试"
            })
        except Exception as e:
            logger.error(f"完整 PPT 生成失败：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": f"PPT 生成失败：{str(e)}"
            })
        finally:
            # 结束 SSE 流
            await progress_queue.put(None)

    # 启动后台任务
    task = asyncio.create_task(generate_in_background())

    # 让出控制权给事件循环，确保后台任务立即开始执行
    # 这样第一个 SSE 事件（analyzing_content 10%）会在返回响应前被放入队列
    await asyncio.sleep(0)

    # 返回 SSE 流
    return StreamingResponse(
        sse_generator(progress_queue),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@router.post("/ppt/merge")
async def merge_ppts(
    files: list[UploadFile] = File(..., description="要合并的多个 PPT 文件"),
    title: str = Form("合并课件", description="合并后课件的标题"),
):
    """
    合并多个 PPT 文件为一个新课件
    Args:
        files: 要合并的 PPTX 文件列表（至少 2 个）
        title: 合并后课件的标题
    Returns:
        合并后的 PPT 文件下载链接
    """
    try:
        if len(files) < 2:
            raise HTTPException(status_code=400, detail="至少需要上传 2 个 PPT 文件")

        if len(files) > 10:
            raise HTTPException(status_code=400, detail="最多支持合并 10 个 PPT 文件")

        # 验证文件类型并保存临时文件
        from pathlib import Path
        import tempfile

        temp_dir = Path(tempfile.gettempdir()) / "ppt_merge"
        temp_dir.mkdir(parents=True, exist_ok=True)

        ppt_paths = []
        try:
            for file in files:
                if not file.filename.lower().endswith(".pptx"):
                    raise HTTPException(
                        status_code=400,
                        detail=f"不支持的文件格式：{file.filename}，仅支持 .pptx 格式"
                    )

                # 保存临时文件
                temp_path = temp_dir / f"{uuid.uuid4().hex}_{file.filename}"
                async with aiofiles.open(temp_path, "wb") as f:
                    content = await file.read()
                    await f.write(content)
                ppt_paths.append(temp_path)

            # 生成输出文件名
            output_file_name = f"merged_{uuid.uuid4().hex[:8]}.pptx"
            output_path = settings.UPLOAD_DIR / "generated" / output_file_name

            # 调用合并服务
            from ..services.ppt_generator import get_ppt_generator
            generator = get_ppt_generator()
            generator.merge_ppts(ppt_paths, output_path, title)

            return JSONResponse(content={
                "success": True,
                "message": f"成功合并 {len(files)} 个 PPT 文件",
                "download_url": f"/uploads/generated/{output_file_name}",
                "file_name": output_file_name,
                "merged_count": len(files)
            })

        finally:
            # 清理临时文件
            for temp_path in ppt_paths:
                try:
                    if temp_path.exists():
                        temp_path.unlink()
                except Exception as e:
                    logger.warning(f"清理临时文件失败：{temp_path}, {e}")

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PPT 合并失败：{e}")
        raise HTTPException(status_code=500, detail=f"PPT 合并失败：{str(e)}")


@router.post("/ppt/parse")
async def parse_ppt(
    file: UploadFile = File(..., description="要解析的 PPTX 文件"),
    extract_enhanced: bool = Form(False, description="是否提取增强元数据（图片、表格、样式等）"),
    max_image_size: int = Form(512, description="图片最大尺寸（宽度/高度），用于控制Base64大小"),
    timeout: int = Form(30, description="解析超时时间（秒），默认 30 秒"),
    use_cache: bool = Form(True, description="是否启用缓存（默认启用）"),
):
    """
    解析 PPTX 文件，提取每页的内容用于前端预览
    Args:
        file: PPTX 文件
        extract_enhanced: 是否提取增强元数据（图片、表格、样式等）
        max_image_size: 图片最大尺寸（宽度/高度），用于控制Base64大小
    Returns:
        JSON 格式：
        基础模式：{ pages: [{ index, title, content, shapes }] }
        增强模式：{ pages: [{ index, title, content: [{type, text, font, position}], shapes, layout }] }
    """
    try:
        start_time = time.time()

        # 内存检查
        memory_usage = _get_memory_usage()
        logger.info(f"PPT 解析 - 内存使用率: {memory_usage:.1%} (警告阈值: {_MEMORY_WARNING_THRESHOLD:.0%}, 临界阈值: {_MEMORY_CRITICAL_THRESHOLD:.0%})")

        if memory_usage > _MEMORY_CRITICAL_THRESHOLD:
            logger.error(f"服务器内存紧张，拒绝解析请求: {memory_usage:.1%}")
            raise HTTPException(status_code=503, detail=f"服务器内存紧张 ({memory_usage:.1%})")
        if memory_usage > _MEMORY_WARNING_THRESHOLD:
            logger.warning(f"内存使用率较高：{memory_usage:.1%}，启用降级模式")
            extract_enhanced = False
        
        # 读取文件内容
        content_bytes = await file.read()
        file_size = len(content_bytes)
        
        # 文件大小检查
        if file_size > 50 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="文件大小不能超过 50MB")
        
        # 检查缓存
        file_hash = _compute_file_hash(content_bytes)
        cache_key = f"{file_hash}_{'enhanced' if extract_enhanced else 'basic'}_{max_image_size}"
        
        if use_cache:
            cached_result = _get_from_cache(cache_key)
            if cached_result:
                logger.info(f"缓存命中：{file.filename}")
                return JSONResponse(content={**cached_result, "from_cache": True})
        
        # 验证文件类型
        if not file.filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="仅支持 .pptx 格式文件")

        # 保存到临时文件
        temp_dir = Path(tempfile.gettempdir()) / "ppt_parse"
        temp_dir.mkdir(parents=True, exist_ok=True)
        temp_path = temp_dir / f"{uuid.uuid4().hex}_{file.filename}"

        try:
            async with aiofiles.open(temp_path, "wb") as f:
                await f.write(content_bytes)

            # 修复无效的 XML 命名空间（某些 PPTX 文件包含 xmlns:ns2="%s" 等无效 URI）
            def _fix_invalid_namespaces(file_path: Path) -> Path:
                """修复 PPTX 文件中无效的 XML 命名空间声明，支持所有 ns%d 变体"""
                import zipfile
                import re
                import shutil

                # 创建临时目录解压（直接执行修复，不移除快速预检）
                fix_dir = temp_dir / f"fixed_{uuid.uuid4().hex}"
                fix_dir.mkdir(parents=True, exist_ok=True)

                try:
                    # 解压 PPTX
                    with zipfile.ZipFile(file_path, 'r') as z:
                        z.extractall(fix_dir)

                    # 修复所有 XML 文件中的无效命名空间
                    # 使用通用正则匹配所有 xmlns:ns\d+="%s" 或 xmlns:ns\d+='%s' 模式（支持单引号和双引号）
                    def replace_ns(match):
                        """替换命名空间，保持引号类型一致"""
                        ns_num = match.group(1)
                        quote = match.group(2)  # 捕获引号类型（单引号或双引号）
                        return f'xmlns:ns{ns_num}={quote}http://schemas.openxmlformats.org/presentationml/2006/placeholder{quote}'

                    ns_pattern = re.compile(r'xmlns:ns(\d+)=([\'"])%s\2')
                    # 匹配包含 { 的无效占位符模式
                    placeholder_pattern = re.compile(r'xmlns:([a-zA-Z0-9_]+)=["\']\{([^}]*)\}[^"\']*["\']')

                    fixed_count = 0
                    total_fixed = 0

                    for xml_file in fix_dir.rglob("*.xml"):
                        try:
                            content = xml_file.read_text(encoding='utf-8')
                            original_content = content

                            # 替换模式1: xmlns:ns%d="%s" 或 xmlns:ns%d='%s'（支持单引号和双引号）
                            content = ns_pattern.sub(replace_ns, content)

                            # 替换模式2: xmlns:x="{.*}" 等包含 { 的占位符
                            content = placeholder_pattern.sub(r'xmlns:\1="http://schemas.openxmlformats.org/presentationml/2006/\2"', content)

                            if content != original_content:
                                xml_file.write_text(content, encoding='utf-8')
                                fixed_count += 1
                                logger.info(f"修复 XML 文件无效命名空间：{xml_file.relative_to(fix_dir)}")
                                # 统计修复的命名空间数量
                                total_fixed += len(ns_pattern.findall(original_content)) + len(placeholder_pattern.findall(original_content))

                        except Exception as e:
                            logger.warning(f"修复 XML 文件失败 {xml_file}: {e}")

                    # 重新打包为 PPTX
                    fixed_path = temp_dir / f"fixed_{uuid.uuid4().hex}_{file_path.name}"
                    with zipfile.ZipFile(fixed_path, 'w', zipfile.ZIP_DEFLATED) as z:
                        for root, dirs, files in os.walk(fix_dir):
                            for f in files:
                                file_on_disk = Path(root) / f
                                arcname = file_on_disk.relative_to(fix_dir)
                                z.write(file_on_disk, arcname)

                    logger.info(f"PPTX 命名空间修复完成：{file_path.name} -> {fixed_path.name} (修复 {fixed_count} 个文件, 共 {total_fixed} 处命名空间)")
                    return fixed_path

                except Exception as e:
                    logger.error(f"修复 PPTX 命名空间失败：{e}")
                    return file_path
                finally:
                    # 清理临时目录
                    try:
                        shutil.rmtree(fix_dir)
                    except Exception:
                        pass

            # 应用命名空间修复
            logger.info(f"=== XML Namespace Fix Enabled ===")
            temp_path = _fix_invalid_namespaces(temp_path)

            # 解析 PPTX
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            try:
                from pptx.enum.text import MSO_AUTO_SHAPE_TYPE
            except ImportError:
                MSO_AUTO_SHAPE_TYPE = None  # 旧版本 python-pptx 可能没有此枚举
            from pptx.dml.color import RGBColor
            import base64
            from io import BytesIO
            from PIL import Image

            prs = Presentation(temp_path)
            pages: List[Dict[str, Any]] = []

            # 提取单个文本运行的样式信息
            def extract_text_run_style(run):
                """提取文本运行的样式信息"""
                font_info = {}
                if hasattr(run, 'font'):
                    font = run.font
                    if font.name:
                        font_info['name'] = font.name
                    if font.size:
                        font_info['size'] = font.size.pt  # 转换为磅值
                    if font.bold is not None:
                        font_info['bold'] = font.bold
                    if font.italic is not None:
                        font_info['italic'] = font.italic
                    if font.underline is not None:
                        font_info['underline'] = font.underline
                    if font.color and hasattr(font.color, 'rgb') and font.color.rgb:
                        # 转换 RGBColor 为十六进制字符串
                        rgb = font.color.rgb
                        font_info['color'] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                return font_info

            # 压缩并转换图片为 Base64
            def image_to_base64(image_bytes: bytes, max_size: int = 512) -> str:
                """将图片压缩并转换为 Base64"""
                try:
                    img = Image.open(BytesIO(image_bytes))

                    # 如果图片过大，进行缩放
                    if img.width > max_size or img.height > max_size:
                        img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)

                    # 转换为 RGB（如果需要）
                    if img.mode in ('RGBA', 'P', 'LA'):
                        img = img.convert('RGB')

                    # 保存为 JPEG 并压缩
                    output = BytesIO()
                    img.save(output, format='JPEG', quality=75, optimize=True, progressive=True)
                    output.seek(0)

                    # 转换为 Base64
                    base64_str = base64.b64encode(output.getvalue()).decode('utf-8')
                    return f"data:image/jpeg;base64,{base64_str}"
                except Exception as e:
                    logger.warning(f"图片压缩失败：{e}")
                    return None

            # 解析单个形状
            def parse_shape(shape, slide_width: int, slide_height: int):
                """解析单个形状，提取详细信息"""
                shape_info = {
                    'type': 'unknown',
                    'name': shape.name if hasattr(shape, 'name') else '',
                    'position': {
                        'x': shape.left.pt if hasattr(shape, 'left') and shape.left else 0,
                        'y': shape.top.pt if hasattr(shape, 'top') and shape.top else 0,
                        'width': shape.width.pt if hasattr(shape, 'width') and shape.width else 0,
                        'height': shape.height.pt if hasattr(shape, 'height') and shape.height else 0
                    }
                }

                # 提取形状类型
                if hasattr(shape, 'shape_type'):
                    shape_type = shape.shape_type
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        shape_info['type'] = 'picture'
                    elif shape_type == MSO_SHAPE_TYPE.TABLE:
                        shape_info['type'] = 'table'
                    elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        shape_info['type'] = 'text_box'
                    elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        shape_info['type'] = 'auto_shape'
                        if hasattr(shape, 'auto_shape_type'):
                            shape_info['auto_shape_type'] = str(shape.auto_shape_type)
                    elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        shape_info['type'] = 'placeholder'
                    else:
                        shape_info['type'] = str(shape_type)

                # 提取文本内容（增强模式）
                if extract_enhanced and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    paragraphs = []
                    for paragraph in shape.text_frame.paragraphs:
                        runs = []
                        for run in paragraph.runs:
                            run_info = {
                                'text': run.text,
                                'font': extract_text_run_style(run)
                            }
                            runs.append(run_info)

                        if runs:
                            paragraphs.append({
                                'runs': runs,
                                'alignment': str(paragraph.alignment) if hasattr(paragraph, 'alignment') else 'left'
                            })

                    if paragraphs:
                        shape_info['text_content'] = paragraphs

                # 提取图片（增强模式）
                if extract_enhanced and hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        base64_str = image_to_base64(image_bytes, max_image_size)
                        if base64_str:
                            shape_info['image_base64'] = base64_str
                            shape_info['image_format'] = image.ext
                    except Exception as e:
                        logger.warning(f"提取图片失败：{e}")

                # 提取表格（增强模式）
                if extract_enhanced and hasattr(shape, 'has_table') and shape.has_table:
                    try:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                cell_text = ''
                                if hasattr(cell, 'text_frame'):
                                    for paragraph in cell.text_frame.paragraphs:
                                        cell_text += paragraph.text + '\n'
                                row_data.append(cell_text.strip())
                            table_data.append(row_data)
                        shape_info['table_data'] = table_data
                    except Exception as e:
                        logger.warning(f"提取表格失败：{e}")

                # 计算相对位置（百分比）
                if slide_width > 0 and slide_height > 0:
                    shape_info['position_relative'] = {
                        'x': round((shape_info['position']['x'] / slide_width) * 100, 2),
                        'y': round((shape_info['position']['y'] / slide_height) * 100, 2),
                        'width': round((shape_info['position']['width'] / slide_width) * 100, 2),
                        'height': round((shape_info['position']['height'] / slide_height) * 100, 2)
                    }

                return shape_info

            # 复杂元素检测（图表、SmartArt 等）
            COMPLEX_SHAPE_TYPES = {
                MSO_SHAPE_TYPE.CHART,      # 图表
                MSO_SHAPE_TYPE.DIAGRAM,    # SmartArt
                MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,  # 嵌入对象
                MSO_SHAPE_TYPE.MEDIA,      # 音视频
            }

            complex_element_slides: List[int] = []

            for slide_idx, slide in enumerate(prs.slides):
                # 获取页面尺寸
                slide_width = prs.slide_width.pt if hasattr(prs, 'slide_width') else 960
                slide_height = prs.slide_height.pt if hasattr(prs, 'slide_height') else 540

                page_data: Dict[str, Any] = {
                    'index': slide_idx + 1,  # 从 1 开始计数
                    'title': '',
                    'content': [],
                    'shapes': [],
                    'layout': {
                        'width': slide_width,
                        'height': slide_height
                    }
                }

                # 检测复杂元素
                has_complex_elements = False
                complex_element_types: List[str] = []

                for shape in slide.shapes:
                    if hasattr(shape, 'shape_type') and shape.shape_type in COMPLEX_SHAPE_TYPES:
                        has_complex_elements = True
                        complex_element_types.append(str(shape.shape_type))

                    shape_info = parse_shape(shape, slide_width, slide_height)
                    page_data['shapes'].append(shape_info)

                    # 提取标题和内容（基础模式保持兼容）
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text_content = ""
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                text_content += para_text + "\n"

                        if text_content.strip():
                            # 第一个文本框通常作为标题（通过形状名称判断）
                            if not page_data['title'] and ('Title' in shape_info['name'] or shape_info['type'] == 'placeholder'):
                                page_data['title'] = text_content.strip()
                            else:
                                # 基础模式：仅文本
                                if not extract_enhanced:
                                    page_data['content'].append(text_content.strip())
                                # 增强模式：包含样式和位置
                                else:
                                    page_data['content'].append({
                                        'type': 'text',
                                        'text': text_content.strip(),
                                        'position': shape_info['position'],
                                        'font': shape_info.get('text_content', [{}])[0].get('runs', [{}])[0].get('font', {}) if shape_info.get('text_content') else {}
                                    })

                # 添加复杂元素标记
                if has_complex_elements:
                    page_data['has_complex_elements'] = True
                    page_data['complex_element_types'] = list(set(complex_element_types))  # 去重
                    complex_element_slides.append(slide_idx + 1)  # 记录页码（从 1 开始）
                    logger.warning(f"第 {slide_idx + 1} 页包含复杂元素：{complex_element_types}")

                pages.append(page_data)

            result = {
                'success': True,
                'file_name': file.filename,
                'file_size': file_size,
                'total_pages': len(pages),
                'enhanced': extract_enhanced,
                'pages': pages,
                'complex_element_slides': complex_element_slides,  # 包含复杂元素的页码列表
                'parse_time_ms': int((time.time() - start_time) * 1000)
            }
            
            # 缓存结果
            if use_cache:
                _add_to_cache(cache_key, result)
            
            return JSONResponse(content={
                **result,
                "total_time_ms": int((time.time() - start_time) * 1000),
                "memory_usage": _get_memory_usage()
            })

        finally:
            # 清理临时文件
            try:
                if temp_path.exists():
                    temp_path.unlink()
            except Exception as e:
                logger.warning(f"清理临时文件失败：{temp_path}, {e}")

    except HTTPException:
        raise
    except zipfile.BadZipFile as e:
        # 文件格式错误（不是有效的 ZIP/PPTX 文件）
        logger.warning(f"文件格式错误：{file.filename} - {e}")
        raise HTTPException(
            status_code=400,
            detail=f"文件格式错误：'{file.filename}' 不是有效的 PPTX 文件，请确保文件完整且未损坏"
        )
    except Exception as e:
        import traceback
        logger.error(f"PPT 解析失败：{e}")
        logger.error(f"Traceback:\n{traceback.format_exc()}")
        logger.error(f"文件名: {file.filename if 'file' in locals() else 'N/A'}, 大小: {file_size if 'file_size' in locals() else 'N/A'}")
        raise HTTPException(status_code=500, detail=f"PPT 解析失败：{str(e)}")


@router.post("/ppt/smart-merge")
async def smart_merge_ppt(
    file_a: UploadFile = File(..., description="PPT 文件 A"),
    file_b: UploadFile = File(..., description="PPT 文件 B"),
    page_prompts: str = Form("{}", description="页面级提示语 JSON"),
    global_prompt: str = Form("", description="全局合并提示语"),
    api_key: str = Form(..., description="LLM API Key"),
    provider: str = Form("deepseek", description="LLM 服务商"),
    title: str = Form("智能合并课件", description="合并后课件标题"),
    temperature: Optional[float] = Form(0.3, description="LLM 温度参数"),
    max_tokens: Optional[int] = Form(2000, description="LLM 最大输出 token 数"),
):
    """
    智能合并两个 PPT 文件：根据页面级提示语调用 LLM 生成合并策略

    Args:
        file_a: PPT 文件 A
        file_b: PPT 文件 B
        page_prompts: 页面级提示语 JSON，格式：{"a_pages": {"1": "提示语 1", "2": "提示语 2"}, "b_pages": {"1": "提示语 1"}}
        global_prompt: 全局合并提示语，如"将 A 的概念讲解与 B 的例题结合"
        api_key: LLM API Key
        provider: LLM 服务商（deepseek/openai/claude/glm）
        title: 合并后课件标题
        temperature: LLM 温度参数
        max_tokens: LLM 最大输出 token 数

    Returns:
        合并后的 PPT 文件下载链接
    """
    try:
        # 验证文件类型
        if not file_a.filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="文件 A 格式错误，仅支持 .pptx 格式")
        if not file_b.filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="文件 B 格式错误，仅支持 .pptx 格式")

        # 解析页面提示语
        try:
            page_prompts_dict = json.loads(page_prompts) if page_prompts else {}
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="页面提示语 JSON 格式错误")

        # 保存临时文件
        temp_dir = Path(tempfile.gettempdir()) / "ppt_smart_merge"
        temp_dir.mkdir(parents=True, exist_ok=True)

        temp_a = temp_dir / f"{uuid.uuid4().hex}_{file_a.filename}"
        temp_b = temp_dir / f"{uuid.uuid4().hex}_{file_b.filename}"

        ppt_paths = []
        try:
            # 保存文件 A
            async with aiofiles.open(temp_a, "wb") as f:
                content = await file_a.read()
                await f.write(content)
            ppt_paths.append(temp_a)

            # 保存文件 B
            async with aiofiles.open(temp_b, "wb") as f:
                content = await file_b.read()
                await f.write(content)
            ppt_paths.append(temp_b)

            # 步骤 1: 解析两个 PPT 的内容
            from pptx import Presentation

            prs_a = Presentation(temp_a)
            prs_b = Presentation(temp_b)

            # 提取 A 的内容摘要
            a_pages_info = []
            for i, slide in enumerate(prs_a.slides):
                page_info = {"index": i + 1, "title": "", "content_summary": ""}
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        if not page_info["title"]:
                            page_info["title"] = shape.text.strip()[:50]
                        else:
                            page_info["content_summary"] += shape.text.strip()[:100] + "..."
                            break
                a_pages_info.append(page_info)

            # 提取 B 的内容摘要
            b_pages_info = []
            for i, slide in enumerate(prs_b.slides):
                page_info = {"index": i + 1, "title": "", "content_summary": ""}
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        if not page_info["title"]:
                            page_info["title"] = shape.text.strip()[:50]
                        else:
                            page_info["content_summary"] += shape.text.strip()[:100] + "..."
                            break
                b_pages_info.append(page_info)

            # 步骤 2: 构建 LLM 提示词，生成合并策略
            from ..services.llm import get_llm_service, LLMProvider

            # 构建页面提示语字符串
            a_prompts_str = ""
            for page_num, prompt in page_prompts_dict.get("a_pages", {}).items():
                page_data = a_pages_info[int(page_num) - 1] if int(page_num) <= len(a_pages_info) else {}
                a_prompts_str += f"- A 第{page_num}页 [{page_data.get('title', '无标题')}]: {prompt}\n"

            b_prompts_str = ""
            for page_num, prompt in page_prompts_dict.get("b_pages", {}).items():
                page_data = b_pages_info[int(page_num) - 1] if int(page_num) <= len(b_pages_info) else {}
                b_prompts_str += f"- B 第{page_num}页 [{page_data.get('title', '无标题')}]: {prompt}\n"

            # 系统提示词
            system_prompt = """你是一个专业的 PPT 合并助手。请根据用户提供的页面级提示语，生成 PPT 合并策略。

请严格按照以下 JSON 格式输出：
{
  "slides_to_merge": [
    {
      "from_a": [1, 2],      // 从 A 中选取的页码
      "from_b": [3, 4],      // 从 B 中选取的页码
      "action": "combine",   // 合并方式：combine(合并到一页)/append_a(追加 A)/append_b(追加 B)
      "instruction": "保留标题，正文合并"  // 合并指令
    }
  ],
  "slides_to_skip_a": [5, 6],  // 从 A 中跳过的页码
  "slides_to_skip_b": [7, 8],  // 从 B 中跳过的页码
  "global_adjustments": "统一字体和颜色"  // 全局调整说明
}

要求：
1. 所有页码从 1 开始计数
2. 封面页（第 1 页）通常不需要合并，会自动跳过
3. action 只能是 combine/append_a/append_b
4. 未被明确处理的页面默认追加到末尾"""

            # 用户提示词
            user_prompt = f"""PPT A 内容（共{len(a_pages_info)}页）：
{json.dumps(a_pages_info, ensure_ascii=False, indent=2)}

PPT B 内容（共{len(b_pages_info)}页）：
{json.dumps(b_pages_info, ensure_ascii=False, indent=2)}

页面级提示语：
A 页面：
{a_prompts_str if a_prompts_str else "无"}

B 页面：
{b_prompts_str if b_prompts_str else "无"}

全局提示语：
{global_prompt if global_prompt else "无"}

请根据以上信息，生成合并策略 JSON。"""

            # 调用 LLM
            llm_service = get_llm_service(
                provider=provider,
                api_key=api_key,
                temperature=temperature,
                max_tokens=max_tokens
            )

            strategy_response, usage_info = llm_service.chat_with_usage([
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ])

            # 解析 LLM 返回的策略
            try:
                merge_strategy = json.loads(strategy_response)
            except json.JSONDecodeError as e:
                logger.error(f"LLM 返回的策略 JSON 解析失败：{strategy_response[:500]}")
                # 尝试修复 JSON
                from ..services.llm import LLMService
                fix_service = LLMService()
                fixed_json = fix_service._fix_json(strategy_response)
                try:
                    merge_strategy = json.loads(fixed_json)
                except:
                    raise HTTPException(status_code=500, detail=f"LLM 生成的策略 JSON 格式错误：{e}")

            logger.info(f"✅ LLM 响应成功 - provider={provider}, model={llm_service.model}")
            logger.info(f"   Tokens 使用：prompt={usage_info['prompt_tokens']}, completion={usage_info['completion_tokens']}, total={usage_info['total_tokens']}")
            logger.info(f"   Request ID: {usage_info['request_id']}")
            logger.info(f"   Finish reason: {usage_info['finish_reason']}")
            logger.info(f"LLM 生成的合并策略：{json.dumps(merge_strategy, ensure_ascii=False)}")

            # 步骤 3: 执行智能合并
            output_file_name = f"smart_merged_{uuid.uuid4().hex[:8]}.pptx"
            output_path = settings.UPLOAD_DIR / "generated" / output_file_name

            from ..services.ppt_generator import get_ppt_generator
            generator = get_ppt_generator()
            generator.smart_merge_ppts(
                ppt_a_path=temp_a,
                ppt_b_path=temp_b,
                output_path=output_path,
                merge_strategy=merge_strategy,
                title=title
            )

            return JSONResponse(content={
                "success": True,
                "message": "智能合并成功",
                "download_url": f"/uploads/generated/{output_file_name}",
                "file_name": output_file_name,
                "strategy": merge_strategy,
                "merged_from": [file_a.filename, file_b.filename]
            })

        finally:
            # 清理临时文件
            for temp_path in ppt_paths:
                try:
                    if temp_path.exists():
                        temp_path.unlink()
                except Exception as e:
                    logger.warning(f"清理临时文件失败：{temp_path}, {e}")

    except HTTPException:
        raise
    except ValueError as e:
        logger.error(f"参数错误：{e}")
        raise HTTPException(status_code=400, detail=str(e))
    except TimeoutError as e:
        logger.error(f"LLM 调用超时：{e}")
        raise HTTPException(status_code=408, detail="LLM 调用超时，请稍后重试")
    except Exception as e:
        logger.error(f"智能合并失败：{e}")
        raise HTTPException(status_code=500, detail=f"智能合并失败：{str(e)}")


@router.post("/ppt/smart-merge-stream")
async def smart_merge_ppt_stream(
    file_a: UploadFile = File(..., description="PPT 文件 A"),
    file_b: UploadFile = File(..., description="PPT 文件 B"),
    page_prompts: str = Form("{}", description="页面级提示语 JSON"),
    global_prompt: str = Form("", description="全局合并提示语"),
    api_key: str = Form(..., description="LLM API Key"),
    provider: str = Form("deepseek", description="LLM 服务商"),
    title: str = Form("智能合并课件", description="合并后课件标题"),
    temperature: Optional[float] = Form(0.3, description="LLM 温度参数"),
    max_tokens: Optional[int] = Form(2000, description="LLM 最大输出 token 数"),
):
    """
    智能合并两个 PPT 文件（SSE 流式响应）：根据页面级提示语调用 LLM 生成合并策略

    Args:
        file_a: PPT 文件 A
        file_b: PPT 文件 B
        page_prompts: 页面级提示语 JSON，格式：{"a_pages": {"1": "提示语 1", "2": "提示语 2"}, "b_pages": {"1": "提示语 1"}}
        global_prompt: 全局合并提示语，如"将 A 的概念讲解与 B 的例题结合"
        api_key: LLM API Key
        provider: LLM 服务商（deepseek/openai/claude/glm）
        title: 合并后课件标题
        temperature: LLM 温度参数
        max_tokens: LLM 最大输出 token 数

    Returns:
        SSE 流式响应，包含进度事件和最终结果
    """
    progress_queue: asyncio.Queue = asyncio.Queue()

    async def merge_in_background():
        """后台执行合并任务"""
        logger.info("后台智能合并任务启动")
        try:
            # 阶段 1: 上传文件 (10%)
            await progress_queue.put({
                "stage": "uploading_files",
                "progress": 10,
                "message": "正在上传 PPT 文件..."
            })
            logger.info("智能合并阶段 1: 上传文件 (10%)")

            # 验证文件类型
            if not file_a.filename.lower().endswith(".pptx"):
                raise HTTPException(status_code=400, detail="文件 A 格式错误，仅支持 .pptx 格式")
            if not file_b.filename.lower().endswith(".pptx"):
                raise HTTPException(status_code=400, detail="文件 B 格式错误，仅支持 .pptx 格式")

            # 解析页面提示语
            try:
                page_prompts_dict = json.loads(page_prompts) if page_prompts else {}
            except json.JSONDecodeError:
                raise HTTPException(status_code=400, detail="页面提示语 JSON 格式错误")

            # 保存临时文件
            temp_dir = Path(tempfile.gettempdir()) / "ppt_smart_merge"
            temp_dir.mkdir(parents=True, exist_ok=True)

            temp_a = temp_dir / f"{uuid.uuid4().hex}_{file_a.filename}"
            temp_b = temp_dir / f"{uuid.uuid4().hex}_{file_b.filename}"

            ppt_paths = []
            try:
                async with aiofiles.open(temp_a, "wb") as f:
                    content = await file_a.read()
                    await f.write(content)
                ppt_paths.append(temp_a)

                async with aiofiles.open(temp_b, "wb") as f:
                    content = await file_b.read()
                    await f.write(content)
                ppt_paths.append(temp_b)

                # 阶段 2: 解析 PPT 内容 (25%)
                await progress_queue.put({
                    "stage": "parsing",
                    "progress": 25,
                    "message": "正在解析 PPT 内容..."
                })
                logger.info("智能合并阶段 2: 解析 PPT 内容 (25%)")

                # 应用 XML 命名空间修复（与 /ppt/parse API 相同的修复逻辑）
                logger.info(f"=== XML Namespace Fix Enabled for Merge ===")
                temp_a = _fix_invalid_namespaces(temp_a)
                temp_b = _fix_invalid_namespaces(temp_b)

                from pptx import Presentation

                prs_a = Presentation(temp_a)
                prs_b = Presentation(temp_b)

                a_pages_info = []
                for i, slide in enumerate(prs_a.slides):
                    page_info = {"index": i + 1, "title": "", "content_summary": ""}
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            if not page_info["title"]:
                                page_info["title"] = shape.text.strip()[:50]
                            else:
                                page_info["content_summary"] += shape.text.strip()[:100] + "..."
                                break
                    a_pages_info.append(page_info)

                b_pages_info = []
                for i, slide in enumerate(prs_b.slides):
                    page_info = {"index": i + 1, "title": "", "content_summary": ""}
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text.strip():
                            if not page_info["title"]:
                                page_info["title"] = shape.text.strip()[:50]
                            else:
                                page_info["content_summary"] += shape.text.strip()[:100] + "..."
                                break
                    b_pages_info.append(page_info)

                # 阶段 3: 调用 LLM 生成策略 (50%)
                await progress_queue.put({
                    "stage": "calling_llm",
                    "progress": 50,
                    "message": "正在调用 AI 生成合并策略..."
                })
                logger.info("智能合并阶段 3: 调用 LLM 生成合并策略 (50%)")
                logger.info(f"LLM 请求参数 - provider={provider}, temperature={temperature}, max_tokens={max_tokens}")
                logger.info(f"PPT A 页面数: {len(a_pages_info)}, PPT B 页面数: {len(b_pages_info)}")
                logger.info(f"页面提示语数量 - A: {len(page_prompts_dict.get('a_pages', {}))}, B: {len(page_prompts_dict.get('b_pages', {}))}")
                logger.info(f"全局提示语: {global_prompt[:50] if global_prompt else '无'}")

                from ..services.llm import get_llm_service

                # 构建结构化提示词（支持 keep/discard 格式）
                a_prompts_str = ""
                for page_num, prompt_data in page_prompts_dict.get("a_pages", {}).items():
                    page_data = a_pages_info[int(page_num) - 1] if int(page_num) <= len(a_pages_info) else {}
                    page_title = page_data.get('title', '无标题')
                    if isinstance(prompt_data, dict):
                        keep = prompt_data.get('keep', '')
                        discard = prompt_data.get('discard', '')
                        parts = []
                        if keep:
                            parts.append(f"【保留】{keep}")
                        if discard:
                            parts.append(f"【废弃】{discard}")
                        a_prompts_str += f"- A 第{page_num}页 [{page_title}]: {'; '.join(parts)}\n"
                    else:
                        a_prompts_str += f"- A 第{page_num}页 [{page_title}]: {prompt_data}\n"

                b_prompts_str = ""
                for page_num, prompt_data in page_prompts_dict.get("b_pages", {}).items():
                    page_data = b_pages_info[int(page_num) - 1] if int(page_num) <= len(b_pages_info) else {}
                    page_title = page_data.get('title', '无标题')
                    if isinstance(prompt_data, dict):
                        keep = prompt_data.get('keep', '')
                        discard = prompt_data.get('discard', '')
                        parts = []
                        if keep:
                            parts.append(f"【保留】{keep}")
                        if discard:
                            parts.append(f"【废弃】{discard}")
                        b_prompts_str += f"- B 第{page_num}页 [{page_title}]: {'; '.join(parts)}\n"
                    else:
                        b_prompts_str += f"- B 第{page_num}页 [{page_title}]: {prompt_data}\n"

                # 发送 thinking 状态，让前端知道 AI 正在工作
                await progress_queue.put({
                    "stage": "thinking",
                    "progress": 45,
                    "message": "AI 正在分析两个 PPT 的内容结构..."
                })

                system_prompt = """你是一位专业的教学课件合并专家。你的任务是分析两个 PPT 的内容，根据用户的指令生成最优的合并策略。

## 你的能力
- 理解教学内容的逻辑结构和知识点关系
- 识别重复、冗余和互补的内容
- 根据用户的"保留"和"废弃"指令精确执行操作
- 生成清晰可执行的合并策略

## 输出格式
请严格输出以下 JSON 格式（不要包含其他文字）：
```json
{
  "slides_to_merge": [
    {
      "from_a": [页码数组],
      "from_b": [页码数组],
      "action": "combine|append_a|append_b|rewrite",
      "instruction": "具体的合并指令说明"
    }
  ],
  "slides_to_skip_a": [要跳过的 A 页码],
  "slides_to_skip_b": [要跳过的 B 页码],
  "global_adjustments": "全局调整说明",
  "merge_rationale": "合并策略的整体思路说明"
}
```

## action 说明
- `combine`: 将 A 和 B 的指定页面合并为一页（内容融合）
- `append_a`: 保留 A 的指定页面原样追加
- `append_b`: 保留 B 的指定页面原样追加
- `rewrite`: 需要重新组织内容（将在 instruction 中说明如何重写）

## 规则
1. 页码从 1 开始计数
2. 用户标记【保留】的内容必须保留，标记【废弃】的内容必须移除
3. 未被用户明确标记的页面，根据内容质量和全局策略自动决定
4. 避免重复内容：如果 A 和 B 有相似页面，选择内容更丰富的版本
5. 保持教学逻辑连贯：合并后的页面顺序应符合教学流程
6. instruction 字段要具体说明如何处理每组页面的内容"""

                user_prompt = f"""## PPT A 内容（共 {len(a_pages_info)} 页）
{json.dumps(a_pages_info, ensure_ascii=False, indent=2)}

## PPT B 内容（共 {len(b_pages_info)} 页）
{json.dumps(b_pages_info, ensure_ascii=False, indent=2)}

## 用户页面级指令
### A 页面指令：
{a_prompts_str if a_prompts_str else "（无特定指令）"}

### B 页面指令：
{b_prompts_str if b_prompts_str else "（无特定指令）"}

## 用户全局合并策略
{global_prompt if global_prompt else "（无特定策略，请根据内容自动生成最优合并方案）"}

请分析以上内容，生成合并策略 JSON。"""

                logger.info("准备初始化 LLM 服务...")

                llm_service = get_llm_service(
                    provider=provider,
                    api_key=api_key,
                    temperature=temperature,
                    max_tokens=max_tokens
                )

                logger.info("LLM 服务初始化完成，开始调用 chat API...")
                logger.info(f"System prompt 长度: {len(system_prompt)}, User prompt 长度: {len(user_prompt)}")

                await progress_queue.put({
                    "stage": "thinking",
                    "progress": 50,
                    "message": "AI 正在生成合并策略，请稍候..."
                })

                import time as _time
                llm_start_time = _time.time()

                try:
                    strategy_response, usage_info = llm_service.chat_with_usage([
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ])
                    llm_elapsed = _time.time() - llm_start_time
                    logger.info(f"✅ LLM 响应成功 - provider={provider}, model={llm_service.model}")
                    logger.info(f"   Tokens 使用：prompt={usage_info['prompt_tokens']}, completion={usage_info['completion_tokens']}, total={usage_info['total_tokens']}")
                    logger.info(f"   Request ID: {usage_info['request_id']}")
                    logger.info(f"   Finish reason: {usage_info['finish_reason']}")
                    logger.info(f"   API 调用耗时：{llm_elapsed:.1f}s")
                    logger.info(f"   返回字符数：{len(strategy_response) if strategy_response else 0}")
                    logger.info(f"LLM 响应内容: {strategy_response[:500]}...")
                except Exception as e:
                    llm_elapsed = _time.time() - llm_start_time
                    logger.error(f"❌ LLM chat 失败（耗时 {llm_elapsed:.1f}s）：{type(e).__name__}: {e}", exc_info=True)
                    raise

                await progress_queue.put({
                    "stage": "calling_llm",
                    "progress": 65,
                    "message": f"AI 策略生成完成（耗时 {llm_elapsed:.1f}s），正在解析..."
                })

                try:
                    merge_strategy = json.loads(strategy_response)
                except json.JSONDecodeError as e:
                    logger.error(f"❌ LLM 返回的策略 JSON 格式解析失败")
                    logger.error(f"LLM 原始响应: {strategy_response[:500]}...")
                    from ..services.llm import LLMService
                    fix_service = LLMService()
                    # 先提取 JSON（处理 markdown 代码块），再修复格式
                    extracted_json = fix_service._extract_json_from_response(strategy_response)
                    logger.info(f"🔧 尝试提取 JSON: {extracted_json[:200]}...")
                    fixed_json = fix_service._fix_json(extracted_json)
                    logger.info(f"🔧 尝试修复 JSON: {fixed_json[:200]}...")
                    try:
                        merge_strategy = json.loads(fixed_json)
                        logger.info("JSON 修复并解析成功")
                    except Exception as fix_e:
                        logger.error(f"JSON 修复失败：{fix_e}")
                        raise HTTPException(status_code=500, detail=f"LLM 生成的策略 JSON 格式错误：{e}")

                logger.info(f"LLM 生成的合并策略：{json.dumps(merge_strategy, ensure_ascii=False)}")

                # 策略验证与修正
                if "slides_to_merge" not in merge_strategy:
                    merge_strategy["slides_to_merge"] = []
                if "slides_to_skip_a" not in merge_strategy:
                    merge_strategy["slides_to_skip_a"] = []
                if "slides_to_skip_b" not in merge_strategy:
                    merge_strategy["slides_to_skip_b"] = []
                if "global_adjustments" not in merge_strategy:
                    merge_strategy["global_adjustments"] = ""

                # 验证页码范围
                total_a = len(a_pages_info)
                total_b = len(b_pages_info)
                for item in merge_strategy["slides_to_merge"]:
                    if "from_a" in item:
                        item["from_a"] = [p for p in item["from_a"] if 1 <= p <= total_a]
                    if "from_b" in item:
                        item["from_b"] = [p for p in item["from_b"] if 1 <= p <= total_b]
                    if "action" not in item or item["action"] not in ("combine", "append_a", "append_b", "rewrite"):
                        item["action"] = "combine"

                merge_strategy["slides_to_skip_a"] = [p for p in merge_strategy["slides_to_skip_a"] if 1 <= p <= total_a]
                merge_strategy["slides_to_skip_b"] = [p for p in merge_strategy["slides_to_skip_b"] if 1 <= p <= total_b]

                rationale = merge_strategy.get("merge_rationale", "")
                if rationale:
                    logger.info(f"AI 合并思路: {rationale}")

                # 阶段 4: 执行合并 (75%)
                await progress_queue.put({
                    "stage": "merging",
                    "progress": 75,
                    "message": "正在执行智能合并..."
                })
                logger.info("智能合并阶段 4: 执行智能合并 (75%)")
                logger.info(f"合并策略: {json.dumps(merge_strategy, ensure_ascii=False)}")

                output_file_name = f"smart_merged_{uuid.uuid4().hex[:8]}.pptx"
                output_path = settings.UPLOAD_DIR / "generated" / output_file_name

                from ..services.ppt_generator import get_ppt_generator
                generator = get_ppt_generator()
                generator.smart_merge_ppts(
                    ppt_a_path=temp_a,
                    ppt_b_path=temp_b,
                    output_path=output_path,
                    merge_strategy=merge_strategy,
                    title=title
                )

                logger.info(f"✅ 文件合并完成: {output_file_name}")

                response_data = {
                    "success": True,
                    "message": "智能合并成功",
                    "download_url": f"/uploads/generated/{output_file_name}",
                    "file_name": output_file_name,
                    "strategy": merge_strategy,
                    "merged_from": [file_a.filename, file_b.filename]
                }

                # 阶段 5: 完成 (100%)
                await progress_queue.put({
                    "stage": "complete",
                    "progress": 100,
                    "message": "合并完成！",
                    "result": response_data
                })
                logger.info("✅ 智能合并完成 (100%)")

            finally:
                for temp_path in ppt_paths:
                    try:
                        if temp_path.exists():
                            temp_path.unlink()
                    except Exception as e:
                        logger.warning(f"清理临时文件失败：{temp_path}, {e}")

        except ValueError as e:
            logger.error(f"参数错误：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": f"参数错误：{str(e)}"
            })
        except TimeoutError as e:
            logger.error(f"LLM 调用超时：{e}")
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": "LLM 调用超时，请稍后重试"
            })
        except Exception as e:
            logger.error(f"智能合并失败：{type(e).__name__} - {e}")
            logger.error(f"Exception details: hasattr detail={hasattr(e, 'detail')}, detail value={getattr(e, 'detail', 'N/A')}")
            logger.error(f"str(e)={str(e)}, repr(e)={repr(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            error_msg = getattr(e, 'detail', None) if hasattr(e, 'detail') else str(e)
            await progress_queue.put({
                "stage": "error",
                "progress": 0,
                "message": f"智能合并失败：{error_msg or type(e).__name__}"
            })
        finally:
            await progress_queue.put(None)

    task = asyncio.create_task(merge_in_background())
    await asyncio.sleep(0)

    return StreamingResponse(
        sse_generator(progress_queue),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@router.post("/ppt/parse-structured")
async def parse_ppt_structured(
    file: UploadFile = File(..., description="要解析的 PPTX 文件"),
    max_image_size: int = Form(512, description="图片最大尺寸（用于压缩）"),
):
    """
    解析 PPTX 文件为中间结构（供 AI 消费）

    返回结构化的中间数据，包含：
    - 每页的元素列表（文本、图片、表格）
    - 元素的位置、样式
    - 教学语义（页面类型、教学角色）

    设计文档: .claude-coder/plans/ppt-merge-technical-design.md#2-数据结构设计

    Returns:
        JSON 格式的中间结构
    """
    import tempfile

    try:
        # 验证文件类型
        if not file.filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="仅支持 .pptx 格式文件")

        # 保存到临时文件
        content_bytes = await file.read()
        temp_dir = Path(tempfile.gettempdir()) / "ppt_parse_structured"
        temp_dir.mkdir(parents=True, exist_ok=True)
        temp_path = temp_dir / f"{uuid.uuid4().hex}_{file.filename}"

        async with aiofiles.open(temp_path, "wb") as f:
            await f.write(content_bytes)

        # 调用解析服务
        from ..services.ppt_content_parser import parse_pptx_to_structure

        result = parse_pptx_to_structure(temp_path, max_image_size=max_image_size)

        # 清理临时文件
        try:
            temp_path.unlink()
        except Exception:
            pass

        logger.info(f"结构化解析完成: {file.filename}, {len(result.get('slides', []))} 页")

        return JSONResponse(content=result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"结构化解析失败: {e}")
        raise HTTPException(status_code=500, detail=f"解析失败: {str(e)}")


@router.get("/ppt/libreoffice-status")
async def get_libreoffice_status():
    """
    检查 LibreOffice 安装状态

    Returns:
        installed: 是否已安装
        path: 安装路径（已安装时）
        guide: 安装指引（未安装时）
    """
    from ..services.ppt_to_image import check_libreoffice

    result = check_libreoffice()
    return JSONResponse(content=result)


@router.post("/ppt/convert-to-images")
async def convert_ppt_to_images(
    file: UploadFile = File(..., description="要转换的 PPTX 文件"),
    resolution: str = Form("high", description="图片分辨率：high/medium/low"),
    pages: str = Form("", description="指定页码，逗号分隔（可选，默认全部）"),
):
    """
    将 PPTX 每页转换为 PNG 图片

    使用 LibreOffice headless 模式实现 100% 真实渲染。

    Args:
        file: PPTX 文件
        resolution: 图片分辨率 (high=1920x1080, medium=1280x720, low=640x360)
        pages: 指定页码，如 "0,1,2"（可选）

    Returns:
        images: 图片列表 [{page, url, width, height}]

    设计文档: .claude-coder/plans/ppt-merge-technical-design.md#5-版本化管理设计
    """
    import tempfile

    try:
        # 验证文件类型
        if not file.filename.lower().endswith(".pptx"):
            raise HTTPException(status_code=400, detail="仅支持 .pptx 格式文件")

        # 解析页码
        pages_list = None
        if pages.strip():
            try:
                pages_list = [int(p.strip()) for p in pages.split(",")]
            except ValueError:
                raise HTTPException(status_code=400, detail="页码格式错误，应为逗号分隔的数字")

        # 保存到临时文件
        content_bytes = await file.read()
        temp_dir = Path(tempfile.gettempdir()) / "ppt_to_image"
        temp_dir.mkdir(parents=True, exist_ok=True)
        temp_path = temp_dir / f"{uuid.uuid4().hex}_{file.filename}"

        async with aiofiles.open(temp_path, "wb") as f:
            await f.write(content_bytes)

        # 输出目录
        output_dir = settings.UPLOAD_DIR / "versions"
        output_dir.mkdir(parents=True, exist_ok=True)

        # 调用转换服务
        from ..services.ppt_to_image import convert_pptx_to_images

        result = convert_pptx_to_images(
            pptx_path=temp_path,
            output_dir=output_dir,
            resolution=resolution,
            pages=pages_list
        )

        # 清理临时文件
        try:
            temp_path.unlink()
        except Exception:
            pass

        if not result["success"]:
            # 检查是否是 LibreOffice 未安装
            if "LibreOffice 未安装" in (result.get("error") or ""):
                raise HTTPException(
                    status_code=503,
                    detail={
                        "error": "LibreOffice 未安装",
                        "guide": result.get("error", "")
                    }
                )
            raise HTTPException(status_code=500, detail=result.get("error", "转换失败"))

        logger.info(f"PPT 转图片完成: {file.filename}, {len(result['images'])} 页")

        return JSONResponse(content={
            "success": True,
            "images": result["images"],
            "total": len(result["images"])
        })

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PPT 转图片失败: {e}")
        raise HTTPException(status_code=500, detail=f"转换失败: {str(e)}")



# ==================== AI 内容融合引擎 API ====================
# 修复：移除错误的 sse_generator 导入（它在 ppt.py 中定义，不在 ppt_generator.py 中）

@router.post('/ppt/ai-merge')
async def ai_merge_ppts(
    file_a: UploadFile = File(..., description='PPT A 文件'),
    file_b: UploadFile = File(..., description='PPT B 文件'),
    merge_type: str = Form('full', description='合并类型'),
    selected_pages_a: str = Form('', description='PPT A 选中页码'),
    selected_pages_b: str = Form('', description='PPT B 选中页码'),
    single_page_action: str = Form('polish', description='单页处理动作'),
    single_page_index: int = Form(0, description='单页处理的页码'),
    source_doc: str = Form('A', description='单页处理的源文档'),
    custom_prompt: str = Form('', description='自定义合并提示语'),
    provider: str = Form('deepseek', description='LLM 服务商'),
    api_key: str = Form(..., description='API Key'),
    temperature: float = Form(0.3, description='温度参数'),
    max_tokens: int = Form(3000, description='最大输出 token'),
):
    import tempfile
    import asyncio
    progress_queue = asyncio.Queue()
    async def merge_in_background():
        temp_a = None
        temp_b = None
        try:
            logger.info(f'开始 AI 内容融合：type={merge_type}')
            import sys
            sys.stdout.reconfigure(encoding='utf-8')
            sys.stderr.reconfigure(encoding='utf-8')
            from ..services.ppt_content_parser import parse_pptx_to_structure
            from ..services.content_merger import get_content_merger

            # DEBUG: 测试导入是否正常
            logger.info('DEBUG: Imports completed successfully')

            doc_a = {"slides": []}
            doc_b = {"slides": []}

            # feat-168 fix: 单页处理模式只解析源文件，避免解析空的 placeholder 文件
            need_parse_a = merge_type != 'single' or source_doc == 'A'
            need_parse_b = merge_type != 'single' or source_doc == 'B'

            if need_parse_a:
                await progress_queue.put({"stage": "analysis", "progress": 10, "message": "正在解析 PPT A 内容..."})
                content_a = await file_a.read()
                if content_a:  # 只有当文件非空时才解析
                    temp_dir_a = Path(tempfile.gettempdir()) / 'ppt_ai_merge'
                    temp_dir_a.mkdir(parents=True, exist_ok=True)
                    temp_a = temp_dir_a / f'{uuid.uuid4().hex}_{file_a.filename}'
                    async with aiofiles.open(temp_a, 'wb') as f:
                        await f.write(content_a)
                    doc_a = parse_pptx_to_structure(temp_a)
                    logger.info(f'PPT A 解析完成：{len(doc_a.get("slides", []))} 页')

            if need_parse_b:
                await progress_queue.put({"stage": "analysis", "progress": 25, "message": "正在解析 PPT B 内容..."})
                content_b = await file_b.read()
                if content_b:  # 只有当文件非空时才解析
                    temp_dir_b = Path(tempfile.gettempdir()) / 'ppt_ai_merge'
                    temp_dir_b.mkdir(parents=True, exist_ok=True)
                    temp_b = temp_dir_b / f'{uuid.uuid4().hex}_{file_b.filename}'
                    async with aiofiles.open(temp_b, 'wb') as f:
                        await f.write(content_b)
                    doc_b = parse_pptx_to_structure(temp_b)
                    logger.info(f'PPT B 解析完成：{len(doc_b.get("slides", []))} 页')

            await progress_queue.put({"stage": "thinking", "progress": 50, "message": "正在调用 AI 生成合并策略..."})
            logger.info('DEBUG: About to create merger')
            merger = get_content_merger(provider=provider, api_key=api_key, temperature=temperature, max_tokens=max_tokens)
            logger.info('DEBUG: Merger created successfully')
            if merge_type == 'full':
                logger.info('DEBUG: Calling merge_full')
                plan = merger.merge_full(doc_a, doc_b, custom_prompt)
                logger.info('DEBUG: merge_full completed')
                result = {'merge_type': 'full', 'plan': {'merge_strategy': plan.merge_strategy, 'summary': plan.summary, 'knowledge_points': plan.knowledge_points, 'slide_plan': [{'action': p.action.value, 'source': p.source, 'slide_index': p.slide_index, 'sources': p.sources, 'new_content': p.new_content, 'instruction': p.instruction, 'reason': p.reason} for p in plan.slide_plan]}}
            elif merge_type == 'single':
                logger.info('DEBUG: Entering single merge mode')
                source_doc_data = doc_a if source_doc == 'A' else doc_b
                slides = source_doc_data.get('slides', [])
                if single_page_index < 0 or single_page_index >= len(slides):
                    raise ValueError(f'页码超出范围：{single_page_index}')

                # DEBUG: 打印 slide 数据类型
                slide = slides[single_page_index]
                tc = slide.get('teaching_content', {})
                mp = tc.get('main_points', [])
                logger.info(f'DEBUG: slide teaching_content.main_points type={type(mp)}, value={mp[:2] if mp else []}')

                logger.info('DEBUG: About to call process_single_page')
                try:
                    single_result = merger.process_single_page(slide, single_page_action, custom_prompt)
                    logger.info('DEBUG: process_single_page returned successfully')
                except Exception as e:
                    logger.error(f'DEBUG: process_single_page failed with: {e}')
                    import traceback
                    logger.error(f'DEBUG: Traceback: {traceback.format_exc()}')
                    raise

                # DEBUG: 打印结果类型
                nc = single_result.get('new_content', {})
                nc_mp = nc.get('main_points', [])
                logger.info(f'DEBUG: new_content.main_points type={type(nc_mp)}, first_elem_type={type(nc_mp[0]) if nc_mp else None}')

                # feat-191: 包装为 MergePlan 格式，与 full/partial 模式保持一致
                action_desc = {
                    'polish': '润色',
                    'expand': '扩展',
                    'rewrite': '改写',
                    'extract': '知识点提取'
                }.get(single_page_action, single_page_action)
                # 从 changes 中提取 reason 字段组成原因说明
                # 修复：确保 reason 值为字符串类型，避免 join 失败
                changes_list = single_result.get('changes', [])
                if changes_list and isinstance(changes_list[0], dict):
                    reason_parts = []
                    for c in changes_list:
                        reason_val = c.get('reason', '')
                        if isinstance(reason_val, str) and reason_val:
                            reason_parts.append(reason_val)
                        elif reason_val:
                            reason_parts.append(str(reason_val))
                    reason_text = '; '.join(reason_parts[:3]) if reason_parts else f'AI{action_desc}处理'
                elif changes_list and isinstance(changes_list[0], str):
                    reason_text = '; '.join(str(c) for c in changes_list)
                else:
                    reason_text = f'AI{action_desc}处理'
                result = {
                    'merge_type': 'single',
                    'plan': {
                        'merge_strategy': f'单页{action_desc}操作',
                        'summary': f'对第 {single_page_index + 1} 页执行 {action_desc}操作',
                        'knowledge_points': [],
                        'slide_plan': [{
                            'action': single_result.get('action', single_page_action),
                            'source': source_doc,
                            'slide_index': single_page_index,
                            'sources': None,
                            'new_content': single_result.get('new_content'),
                            'instruction': custom_prompt or f'执行{action_desc}操作',
                            'reason': reason_text
                        }]
                    }
                }
            elif merge_type == 'partial':
                pages_a_idx = [int(p.strip()) for p in selected_pages_a.split(',') if p.strip()] if selected_pages_a else []
                pages_b_idx = [int(p.strip()) for p in selected_pages_b.split(',') if p.strip()] if selected_pages_b else []
                slides_a, slides_b = doc_a.get('slides', []), doc_b.get('slides', [])
                pages_a = [slides_a[i] for i in pages_a_idx if 0 <= i < len(slides_a)]
                pages_b = [slides_b[i] for i in pages_b_idx if 0 <= i < len(slides_b)]
                if not pages_a and not pages_b:
                    raise ValueError('至少需要选择一个页面')
                result = merger.merge_pages(pages_a, pages_b, custom_prompt)
                result['merge_type'] = 'partial'
            else:
                raise ValueError(f'不支持的合并类型：{merge_type}')
            logger.info('DEBUG: About to send result to queue')
            logger.info(f'DEBUG: Result keys: {result.keys()}')

            # 增强：使用 test_json_serialization 进行详细的序列化测试
            serialization_result = test_json_serialization(result, context="ai_merge_result")
            if not serialization_result["success"]:
                logger.error(f'[JSON序列化] ai-merge 失败字段: {serialization_result["problematic_fields"]}')
                raise ValueError(f"JSON 序列化失败: {serialization_result['error']}")
            logger.info(f'[JSON序列化] ai-merge 成功，耗时 {serialization_result["duration_ms"]}ms')

            await progress_queue.put({"stage": "complete", "progress": 100, "message": "AI 内容融合完成！", "result": result})
        except ValueError as e:
            logger.error(f'参数错误：{e}')
            await progress_queue.put({"stage": "error", "progress": 0, "message": f'参数错误：{str(e)}'})
        except Exception as e:
            import traceback
            tb_str = traceback.format_exc()
            print(f'DEBUG ERROR: AI 融合失败：{e}')
            print(f'DEBUG Traceback: {tb_str}')
            logger.error(f'AI 融合失败：{e}')
            logger.error(f'Traceback: {traceback_module.format_exc()}')
            # 在错误消息中包含 traceback 的关键行以便调试
            # 提取错误位置
            tb_lines = [line.strip() for line in tb_str.split('\n') if line.strip() and ('File' in line or 'Error' in line or 'in ' in line)][:5]
            error_detail = ' | '.join(tb_lines)
            await progress_queue.put({"stage": "error", "progress": 0, "message": f'AI 融合失败：{str(e)} [{error_detail}]'})
        finally:
            for temp_path in [temp_a, temp_b]:
                try:
                    if temp_path and temp_path.exists(): temp_path.unlink()
                except: pass
            await progress_queue.put(None)
    task = asyncio.create_task(merge_in_background())
    await asyncio.sleep(0)
    return StreamingResponse(sse_generator(progress_queue), media_type='text/event-stream', headers={'Cache-Control': 'no-cache', 'Connection': 'keep-alive', 'X-Accel-Buffering': 'no'})


# ==================== feat-162: 增量调整 API ====================

@router.post('/ppt/regenerate-slide')
async def regenerate_slide(
    file_a: UploadFile = File(..., description='PPT A 文件'),
    file_b: UploadFile = File(..., description='PPT B 文件'),
    slide_index: int = Form(..., description='要重新生成的页码（在方案中的索引）'),
    prompt: str = Form('', description='重新生成的提示语'),
    provider: str = Form('deepseek', description='LLM 服务商'),
    api_key: str = Form(..., description='API Key'),
    temperature: float = Form(0.3, description='温度参数'),
    max_tokens: int = Form(2000, description='最大输出 token'),
    current_action: str = Form('keep', description='当前页的动作类型'),
    current_content: str = Form('', description='当前页的内容'),
    source_doc: str = Form('A', description='来源文档'),
    source_slide_index: int = Form(0, description='源页码'),
):
    """
    feat-162: 重新生成合并方案中的某一页

    Request:
    - file_a, file_b: PPT 文件
    - slide_index: 在方案中的页码索引（0-indexed）
    - prompt: 用户自定义的重新生成指令
    - current_action: 当前页的动作类型（keep/merge/create等）
    - current_content: 当前页的内容

    Response:
    {
        "new_slide": {
            "action": "keep",
            "new_content": "重新生成的内容",
            "reason": "基于用户指令重新生成"
        }
    }
    """
    import tempfile
    from ..services.ppt_content_parser import parse_pptx_to_structure
    from ..services.content_merger import get_content_merger

    temp_a = None
    temp_b = None

    try:
        logger.info(f'feat-162: 开始重新生成第 {slide_index} 页，prompt={prompt}')

        # 解析 PPT 文件
        content_a = await file_a.read()
        temp_dir_a = Path(tempfile.gettempdir()) / 'ppt_regenerate'
        temp_dir_a.mkdir(parents=True, exist_ok=True)
        temp_a = temp_dir_a / f'{uuid.uuid4().hex}_{file_a.filename}'
        async with aiofiles.open(temp_a, 'wb') as f:
            await f.write(content_a)
        doc_a = parse_pptx_to_structure(temp_a)

        content_b = await file_b.read()
        temp_dir_b = Path(tempfile.gettempdir()) / 'ppt_regenerate'
        temp_dir_b.mkdir(parents=True, exist_ok=True)
        temp_b = temp_dir_b / f'{uuid.uuid4().hex}_{file_b.filename}'
        async with aiofiles.open(temp_b, 'wb') as f:
            await f.write(content_b)
        doc_b = parse_pptx_to_structure(temp_b)

        # 获取源文档的对应页面
        source_doc_data = doc_a if source_doc == 'A' else doc_b
        slides = source_doc_data.get('slides', [])

        # 构建 AI 提示语
        if prompt:
            regenerate_instruction = f"""
用户希望重新生成此页面，具体要求：
{prompt}

请根据上述要求，生成新的页面内容。
"""
        else:
            regenerate_instruction = f"""
请重新生成此页面，优化内容表达和教学效果。
当前动作类型：{current_action}
当前内容：{current_content or '（无）'}
"""

        # 调用 LLM 生成新内容
        merger = get_content_merger(provider=provider, api_key=api_key, temperature=temperature, max_tokens=max_tokens)

        # 根据当前动作类型决定如何处理
        if current_action in ['keep', 'polish', 'expand', 'rewrite', 'extract']:
            # 单页处理模式
            if source_slide_index < len(slides):
                source_slide = slides[source_slide_index]
                result = merger.process_single_page(source_slide, current_action if current_action != 'keep' else 'polish', prompt)
                new_slide = {
                    'action': result.get('action', current_action),
                    'new_content': result.get('new_content', {}).get('title', '') + '\n' + safe_join_list(result.get('new_content', {}).get('main_points', [])),
                    'reason': f'基于用户指令重新生成：{prompt[:50]}...' if prompt else 'AI 优化重新生成'
                }
            else:
                raise ValueError(f'源页码超出范围：{source_slide_index}')
        else:
            # 创建新内容模式
            result = merger.process_single_page(
                {'elements': [], 'teaching_content': {'title': '新页面'}},
                'create',
                prompt or f'为合并方案创建第 {slide_index + 1} 页内容'
            )
            new_slide = {
                'action': 'create',
                'new_content': result.get('new_content', {}).get('title', '') + '\n' + safe_join_list(result.get('new_content', {}).get('main_points', [])),
                'reason': '用户创建的新页面'
            }

        logger.info(f'feat-162: 重新生成成功，new_slide={new_slide}')

        # 增强：测试 JSON 序列化
        response_data = {'success': True, 'new_slide': new_slide}
        serialization_result = test_json_serialization(response_data, context="regenerate_slide")
        if not serialization_result["success"]:
            logger.error(f'[JSON序列化] regenerate_slide 失败字段: {serialization_result["problematic_fields"]}')

        return JSONResponse(content=response_data)

    except ValueError as e:
        logger.error(f'feat-162: 参数错误：{e}')
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f'feat-162: 重新生成失败：{e}')
        raise HTTPException(status_code=500, detail=f'重新生成失败：{str(e)}')
    finally:
        for temp_path in [temp_a, temp_b]:
            try:
                if temp_path and temp_path.exists():
                    temp_path.unlink()
            except:
                pass


# ==================== 版本化管理 API ====================

@router.post("/session/create")
async def create_ppt_session(request: Request):
    """
    创建 PPT 合并会话

    Request: FormData with ppt_a and ppt_b files

    Response:
    {
        "session_id": "abc12345",
        "documents": {
            "ppt_a": {"source_file": "...", "total_slides": 5},
            "ppt_b": {"source_file": "...", "total_slides": 8}
        }
    }
    """
    import tempfile
    from ..services.version_manager import get_version_manager

    # 从 FormData 获取上传的文件
    form = await request.form()

    # 调试日志：打印所有 form 字段
    logger.info(f"收到 session/create 请求，form 字段：{list(form.keys())}")
    for key, value in form.items():
        logger.info(f"  - {key}: type={type(value).__name__}, value={value if not isinstance(value, (UploadFile, StarletteUploadFile)) else f'UploadFile({value.filename})'}")

    # 提取所有 .pptx 文件
    files = {}
    for key, value in form.items():
        # 检查是否是 UploadFile 类型（FastAPI 或 Starlette）
        if isinstance(value, (UploadFile, StarletteUploadFile)):
            filename = value.filename or ""
            if filename.lower().endswith(('.pptx', '.ppt')):
                files[key] = value
                logger.info(f"  -> 识别为 PPT 文件：{key} = {filename}")
        # 兼容处理：检查是否有 filename 属性
        elif hasattr(value, 'filename') and hasattr(value, 'file'):
            filename = getattr(value, 'filename', '') or ""
            if filename.lower().endswith(('.pptx', '.ppt')):
                files[key] = value
                logger.info(f"  -> 识别为 PPT 文件（兼容模式）：{key} = {filename}")

    logger.info(f"创建 PPT 合并会话，识别到的 PPT 文件数：{len(files)}, 文件 keys：{list(files.keys())}")

    if len(files) < 1:
        raise HTTPException(status_code=400, detail="至少需要上传一个文件")

    # 保存上传文件到临时目录
    temp_files = {}
    temp_dir = Path(tempfile.gettempdir()) / 'ppt_sessions'
    temp_dir.mkdir(parents=True, exist_ok=True)

    try:
        for doc_id, file in files.items():
            if not file.filename.lower().endswith(('.pptx', '.ppt')):
                raise HTTPException(status_code=400, detail=f"文件格式错误：{file.filename}，仅支持 .pptx 格式")

            temp_path = temp_dir / f"{uuid.uuid4().hex}_{file.filename}"
            content = await file.read()
            async with aiofiles.open(temp_path, 'wb') as f:
                await f.write(content)
            temp_files[doc_id] = str(temp_path)

        # 创建会话
        version_manager = get_version_manager()
        session = version_manager.create_session(temp_files)

        # 返回简化信息
        documents_info = {}
        for doc_id, doc_state in session.documents.items():
            documents_info[doc_id] = {
                "source_file": doc_state.source_file,
                "total_slides": len(doc_state.slides)
            }

        logger.info(f"会话创建成功：{session.session_id}")

        return JSONResponse(content={
            "session_id": session.session_id,
            "documents": documents_info
        })

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"创建会话失败：{e}")
        raise HTTPException(status_code=500, detail=f"创建会话失败：{str(e)}")


@router.post("/version/create")
async def create_version(
    session_id: str = Body(..., description="会话 ID"),
    document_id: str = Body(..., description="文档 ID (ppt_a, ppt_b)"),
    slide_index: int = Body(..., description="页码（0-indexed）"),
    operation: str = Body(..., description="操作类型 (ai_polish, ai_expand, ai_rewrite, ai_extract)"),
    prompt: Optional[str] = Body(None, description="AI 提示语"),
    new_pptx: Optional[str] = Body(None, description="新生成的单页 PPTX 路径"),
    content_snapshot: Optional[Dict[str, Any]] = Body(None, description="AI 修改的内容快照（中间结构数据）"),
    generate_preview: bool = Body(True, description="是否生成预览图片，默认 true")
):
    """
    创建新版本

    Request:
    {
        "session_id": "abc12345",
        "document_id": "ppt_a",
        "slide_index": 1,
        "operation": "ai_polish",
        "prompt": "用更通俗的语言解释",
        "new_pptx": "/path/to/new_slide.pptx" (可选),
        "content_snapshot": {...} (可选，AI 修改的中间结构数据),
        "generate_preview": true (可选，是否生成预览图片)
    }

    Response:
    {
        "version": "v2",
        "image_url": "http://localhost:8000/public/versions/abc12345/ppt_a_slide1_v2.png",
        "created_at": "10:05:00",
        "has_preview": true
    }
    """
    from ..services.version_manager import get_version_manager

    logger.info(f"创建版本：{document_id} slide {slide_index} ({operation}), generate_preview={generate_preview}")

    try:
        version_manager = get_version_manager()

        # 验证会话存在
        session = version_manager.get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail=f"会话不存在：{session_id}")

        # 确定要使用的 PPTX 路径
        new_pptx_path = Path(new_pptx) if new_pptx else None

        # 如果需要生成预览且提供了 content_snapshot，但未提供 new_pptx，
        # 则自动从 content_snapshot 生成临时 PPTX
        if generate_preview and content_snapshot and not new_pptx_path:
            from ..services.ppt_generator import get_ppt_generator

            # 创建临时目录存放生成的 PPTX
            temp_dir = settings.PUBLIC_DIR / "versions" / session_id / "temp"
            temp_dir.mkdir(parents=True, exist_ok=True)

            # 生成临时 PPTX 文件名
            temp_pptx_path = temp_dir / f"{document_id}_slide{slide_index}_temp.pptx"

            logger.info(f"从 content_snapshot 生成临时 PPTX: {temp_pptx_path}")

            # 调用生成器创建单页 PPTX
            ppt_generator = get_ppt_generator()
            ppt_generator.generate_single_slide_pptx(
                content_snapshot=content_snapshot,
                output_path=temp_pptx_path
            )

            new_pptx_path = temp_pptx_path

        # 创建新版本
        version = version_manager.create_version(
            session_id=session_id,
            document_id=document_id,
            slide_index=slide_index,
            operation=operation,
            prompt=prompt,
            new_pptx=new_pptx_path,
            content_snapshot=content_snapshot
        )

        return JSONResponse(content={
            "version": version.version,
            "image_url": version.image_url,
            "created_at": version.created_at,
            "operation": version.operation,
            "prompt": version.prompt,
            "has_preview": bool(version.image_url)
        })

    except HTTPException:
        raise
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"创建版本失败：{e}")
        raise HTTPException(status_code=500, detail=f"创建版本失败：{str(e)}")


@router.post("/version/restore")
async def restore_version(
    session_id: str = Body(..., description="会话 ID"),
    document_id: str = Body(..., description="文档 ID"),
    slide_index: int = Body(..., description="页码"),
    target_version: str = Body(..., description="目标版本号 (v1, v2, ...)")
):
    """
    恢复历史版本

    Request:
    {
        "session_id": "abc12345",
        "document_id": "ppt_a",
        "slide_index": 1,
        "target_version": "v1"
    }

    Response:
    {
        "success": true,
        "current_version": "v1"
    }
    """
    from ..services.version_manager import get_version_manager

    logger.info(f"恢复版本：{document_id} slide {slide_index} -> {target_version}")

    try:
        version_manager = get_version_manager()

        # 验证会话存在
        session = version_manager.get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail=f"会话不存在：{session_id}")

        # 恢复版本
        version_manager.restore_version(
            session_id=session_id,
            document_id=document_id,
            slide_index=slide_index,
            target_version=target_version
        )

        return JSONResponse(content={
            "success": True,
            "current_version": target_version
        })

    except HTTPException:
        raise
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"恢复版本失败：{e}")
        raise HTTPException(status_code=500, detail=f"恢复版本失败：{str(e)}")


@router.post("/slide/toggle")
async def toggle_slide_status(
    session_id: str = Body(..., description="会话 ID"),
    document_id: str = Body(..., description="文档 ID"),
    slide_index: int = Body(..., description="页码"),
    action: str = Body(..., description="操作：delete 或 restore")
):
    """
    删除/恢复页面

    Request:
    {
        "session_id": "abc12345",
        "document_id": "ppt_a",
        "slide_index": 2,
        "action": "delete"
    }

    Response:
    {
        "success": true,
        "status": "deleted"
    }
    """
    from ..services.version_manager import get_version_manager, SlideStatus

    logger.info(f"切换页面状态：{document_id} slide {slide_index} ({action})")

    if action not in ["delete", "restore"]:
        raise HTTPException(status_code=400, detail=f"无效操作：{action}，仅支持 delete/restore")

    try:
        version_manager = get_version_manager()

        # 验证会话存在
        session = version_manager.get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail=f"会话不存在：{session_id}")

        # 切换状态
        status = version_manager.toggle_slide(
            session_id=session_id,
            document_id=document_id,
            slide_index=slide_index,
            action=action
        )

        return JSONResponse(content={
            "success": True,
            "status": status.value
        })

    except HTTPException:
        raise
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"切换页面状态失败：{e}")
        raise HTTPException(status_code=500, detail=f"切换页面状态失败：{str(e)}")


@router.post("/ppt/generate-final")
async def generate_final_ppt(
    session_id: str = Body(..., description="会话 ID"),
    title: str = Body(..., description="最终 PPT 标题"),
    grade: str = Body("6", description="年级"),
    subject: str = Body("general", description="学科"),
    style: str = Body("simple", description="风格"),
    final_selection: Optional[List[str]] = Body(None, description="最终选择的版本ID列表，如 ['ppt_a_0_v1', 'ppt_b_1_v1']"),
    content_snapshots: Optional[Dict[str, Any]] = Body(None, description="前端附带的内容快照（用于 merged/AI 处理版本）")
):
    """
    基于当前选择的版本生成最终 PPT

    流程：
    1. 获取会话数据
    2. 如果提供了 final_selection，只包含指定的版本
    3. 否则收集所有活跃页面的当前版本
    4. 对于 v1 版本：从源 PPTX 直接复制页面
    5. 对于 v2+ 版本：从 content_snapshot 内容快照重建页面
    6. 生成统一样式的最终 PPT

    Request:
    {
        "session_id": "abc12345",
        "title": "合并后的课件标题",
        "grade": "6",
        "subject": "math",
        "style": "simple",
        "final_selection": ["ppt_a_0_v1", "ppt_b_1_v1", "merged_0_v1"]
    }

    Response:
    {
        "success": true,
        "download_url": "/uploads/generated/final_xxx.pptx",
        "file_name": "final_xxx.pptx",
        "total_slides": 10
    }
    """
    from ..services.version_manager import get_version_manager, SlideStatus
    from ..services.ppt_generator import generate_ppt_from_versions
    from pptx import Presentation
    from io import BytesIO
    import re

    logger.info(f"生成最终 PPT: session_id={session_id}, title={title}, final_selection={final_selection}")

    try:
        version_manager = get_version_manager()
        session = version_manager.get_session(session_id)

        if not session:
            raise HTTPException(status_code=404, detail=f"会话不存在：{session_id}")

        merged_slides = []  # List of (source_pptx, slide_index, content_snapshot)

        if final_selection:
            # 使用 final_selection 指定的版本
            logger.info(f"使用 final_selection: {len(final_selection)} 个版本")

            for version_id in final_selection:
                # 解析 version_id: "ppt_a_0_v1" -> doc_id="ppt_a", slide_index=0, version="v1"
                # 或 "merged_0_v1" -> doc_id="merged", slide_index=0, version="v1"
                match = re.match(r'^(ppt_a|ppt_b|merged)_(\d+)_v(\d+)$', version_id)
                if not match:
                    logger.warning(f"无法解析 version_id: {version_id}，跳过")
                    continue

                doc_id = match.group(1)
                slide_index = int(match.group(2))
                version_name = f"v{match.group(3)}"

                # 处理 merged 类型（融合生成的幻灯片）
                if doc_id == 'merged':
                    # 优先使用前端传来的 content_snapshots
                    frontend_snapshot = (content_snapshots or {}).get(version_id)
                    if frontend_snapshot:
                        merged_slides.append({
                            "source_pptx": None,
                            "slide_index": slide_index,
                            "content_snapshot": frontend_snapshot,
                            "version": version_name
                        })
                    elif 'ppt_a' in session.documents:
                        doc_state = session.documents['ppt_a']
                        source_pptx = doc_state.source_file
                        if Path(source_pptx).exists():
                            merged_slides.append({
                                "source_pptx": source_pptx,
                                "slide_index": 0,
                                "content_snapshot": None,
                                "version": version_name
                            })
                    continue

                # 处理普通幻灯片
                if doc_id not in session.documents:
                    logger.warning(f"文档不存在: {doc_id}，跳过")
                    continue

                doc_state = session.documents[doc_id]
                source_pptx = doc_state.source_file

                if not Path(source_pptx).exists():
                    logger.warning(f"源文件不存在：{source_pptx}，跳过")
                    continue

                if slide_index not in doc_state.slides:
                    logger.warning(f"幻灯片不存在: {doc_id}[{slide_index}]，跳过")
                    continue

                slide_state = doc_state.slides[slide_index]

                # 查找指定版本的 content_snapshot
                content_snapshot = None
                for version in slide_state.versions:
                    if version.version == version_name:
                        if version.content_snapshot:
                            content_snapshot = version.content_snapshot
                        break

                # 如果后端没有 content_snapshot，尝试使用前端传来的
                if not content_snapshot and content_snapshots and version_id in content_snapshots:
                    content_snapshot = content_snapshots[version_id]

                merged_slides.append({
                    "source_pptx": source_pptx,
                    "slide_index": slide_index,
                    "content_snapshot": content_snapshot,
                    "version": version_name
                })
        else:
            # 兼容旧逻辑：收集所有活跃页面
            logger.info("使用默认逻辑：收集所有活跃页面")

            for doc_id, doc_state in session.documents.items():
                source_pptx = doc_state.source_file
                # 验证文件存在
                if not Path(source_pptx).exists():
                    logger.warning(f"源文件不存在：{source_pptx}，跳过")
                    continue

                for slide_index, slide_state in doc_state.slides.items():
                    if slide_state.status == SlideStatus.DELETED or not slide_state.current_version:
                        continue  # 跳过已删除页面

                    # 查找当前版本的 content_snapshot
                    content_snapshot = None
                    for version in slide_state.versions:
                        if version.version == slide_state.current_version:
                            if version.content_snapshot:
                                content_snapshot = version.content_snapshot
                            break

                    merged_slides.append({
                        "source_pptx": source_pptx,
                        "slide_index": slide_index,
                        "content_snapshot": content_snapshot,
                        "version": slide_state.current_version
                    })

        if not merged_slides:
            raise HTTPException(status_code=400, detail="没有可生成的页面")

        logger.info(f"收集到 {len(merged_slides)} 个页面用于生成最终 PPT")

        # 生成最终 PPT
        output_path = generate_ppt_from_versions(
            merged_slides,
            title=title,
            grade=grade,
            subject=subject,
            style=style
        )

        file_name = output_path.name

        logger.info(f"最终 PPT 生成成功：{file_name}")

        response_data = {
            "success": True,
            "download_url": f"/uploads/generated/{file_name}",
            "file_name": file_name,
            "total_slides": len(merged_slides)
        }

        # 增强：测试 JSON 序列化
        serialization_result = test_json_serialization(response_data, context="generate_final_ppt")
        if not serialization_result["success"]:
            logger.error(f'[JSON序列化] generate_final_ppt 失败字段: {serialization_result["problematic_fields"]}')

        return JSONResponse(content=response_data)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"生成最终 PPT 失败：{e}")
        raise HTTPException(status_code=500, detail=f"生成最终 PPT 失败：{str(e)}")


@router.get("/session/{session_id}")
async def get_session_info(session_id: str):
    """
    获取会话详情（包含所有文档的完整版本历史）

    Response:
    {
        "session_id": "abc12345",
        "documents": {
            "ppt_a": {
                "source_file": "...",
                "slides": {
                    "0": {
                        "current_version": "v1",
                        "status": "active",
                        "versions": [
                            {"version": "v1", "image_url": "...", "operation": "原始上传"}
                        ]
                    }
                }
            }
        }
    }
    """
    from ..services.version_manager import get_version_manager

    try:
        version_manager = get_version_manager()
        session = version_manager.get_session(session_id)

        if not session:
            raise HTTPException(status_code=404, detail=f"会话不存在：{session_id}")

        # 转换为可序列化的格式
        def serialize_slide_state(slide_state):
            return {
                "current_version": slide_state.current_version,
                "status": slide_state.status.value,
                "versions": [
                    {
                        "version": v.version,
                        "image_url": v.image_url,
                        "created_at": v.created_at,
                        "operation": v.operation,
                        "prompt": v.prompt
                    }
                    for v in slide_state.versions
                ]
            }

        documents = {}
        for doc_id, doc_state in session.documents.items():
            documents[doc_id] = {
                "source_file": doc_state.source_file,
                "slides": {
                    str(idx): serialize_slide_state(slide)
                    for idx, slide in doc_state.slides.items()
                }
            }

        response_data = {
            "session_id": session.session_id,
            "documents": documents
        }

        # 增强：测试 JSON 序列化
        serialization_result = test_json_serialization(response_data, context="get_session_info")
        if not serialization_result["success"]:
            logger.error(f'[JSON序列化] get_session_info 失败字段: {serialization_result["problematic_fields"]}')

        return JSONResponse(content=response_data)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"获取会话信息失败：{e}")
        raise HTTPException(status_code=500, detail=f"获取会话信息失败：{str(e)}")


@router.get("/session/{session_id}/history")
async def get_slide_version_history(
    session_id: str,
    document_id: str = Query(..., description="文档 ID"),
    slide_index: int = Query(..., description="页码")
):
    """
    获取指定页面的版本历史

    Response:
    {
        "versions": [
            {"version": "v1", "image_url": "...", "operation": "原始上传"},
            {"version": "v2", "image_url": "...", "operation": "ai_polish", "prompt": "..."}
        ]
    }
    """
    from ..services.version_manager import get_version_manager

    try:
        version_manager = get_version_manager()
        versions = version_manager.get_version_history(session_id, document_id, slide_index)

        if not versions:
            # 检查是否是因为会话/页面不存在
            session = version_manager.get_session(session_id)
            if not session:
                raise HTTPException(status_code=404, detail=f"会话不存在：{session_id}")

        return JSONResponse(content={
            "versions": [
                {
                    "version": v.version,
                    "image_url": v.image_url,
                    "created_at": v.created_at,
                    "operation": v.operation,
                    "prompt": v.prompt
                }
                for v in versions
            ]
        })

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"获取版本历史失败：{e}")
        raise HTTPException(status_code=500, detail=f"获取版本历史失败：{str(e)}")


# ==================== feat-212: 单页 AI 处理 + 实时 PPT 图片预览 ====================

@router.post("/ppt/ai-merge-single")
async def ai_merge_single_page(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(..., description="PPT 文件"),
    page_index: int = Form(..., description="要处理的页码（0-indexed）"),
    action: str = Form(..., description="操作类型：polish/expand/rewrite/extract"),
    custom_prompt: Optional[str] = Form(None, description="自定义提示语"),
    provider: str = Form("deepseek", description="LLM 服务商"),
    api_key: str = Form(..., description="API Key"),
    temperature: float = Form(0.3, description="温度参数"),
    max_tokens: int = Form(2000, description="最大输出 token"),
):
    """
    feat-212: 单页 AI 处理 + 实时生成 PPT 图片预览

    流程：
    1. 解析原始 PPT 获取页面结构
    2. 调用 LLM 进行 AI 处理
    3. 生成单页 PPTX 文件
    4. 调用 LibreOffice 转换为 PNG
    5. 返回处理结果和图片 URL

    Request:
    - file: PPT 文件
    - page_index: 要处理的页码（0-indexed）
    - action: 操作类型（polish/expand/rewrite/extract）
    - custom_prompt: 可选的自定义提示语
    - provider: LLM 服务商（默认 deepseek）
    - api_key: API Key

    Response:
    {
        "success": true,
        "content": { "title": "...", "main_points": [...], ... },
        "preview_url": "http://localhost:8000/public/images/xxx/page_0.png",
        "image_path": "/path/to/image.png",
        "degraded": false
    }
    """
    temp_pptx = None
    output_pptx = None

    try:
        logger.info(f"单页处理请求: page={page_index}, action={action}, file={file.filename}")

        # 1. 保存上传文件
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="上传文件为空")

        temp_dir = Path(tempfile.gettempdir()) / 'ppt_ai_merge_single'
        temp_dir.mkdir(parents=True, exist_ok=True)

        temp_pptx = temp_dir / f'{uuid.uuid4().hex}_{file.filename}'
        async with aiofiles.open(temp_pptx, 'wb') as f:
            await f.write(content)

        logger.info(f"文件已保存: {temp_pptx}")

        # 2. 解析 PPT 获取原始内容
        from ..services.ppt_content_parser import parse_pptx_to_structure

        doc = parse_pptx_to_structure(temp_pptx)
        slides = doc.get("slides", [])

        if page_index < 0 or page_index >= len(slides):
            raise HTTPException(
                status_code=400,
                detail=f"页码超出范围：{page_index}，有效范围 0-{len(slides) - 1}"
            )

        original_slide = slides[page_index]
        logger.info(f"PPT 解析完成: {len(slides)} 页，处理第 {page_index} 页")

        # 3. 调用 LLM 进行处理
        from ..services.content_merger import get_content_merger

        merger = get_content_merger(
            provider=provider,
            api_key=api_key,
            temperature=temperature,
            max_tokens=max_tokens
        )

        logger.info(f"开始 AI 处理: action={action}")
        result = merger.process_single_page(original_slide, action, custom_prompt)
        logger.info(f"AI 处理完成")

        # 提取处理后的内容
        new_content = result.get("new_content", {})

        # 确保 new_content 包含必要字段
        if not new_content:
            new_content = {
                "title": original_slide.get("teaching_content", {}).get("title", ""),
                "main_points": [],
                "additional_content": ""
            }

        # 类型安全：确保 main_points 是字符串列表
        main_points = new_content.get("main_points", [])
        if not isinstance(main_points, list):
            main_points = []
        else:
            safe_points = []
            for p in main_points:
                if isinstance(p, str):
                    safe_points.append(p)
                elif isinstance(p, dict):
                    text = (p.get('text') or p.get('content') or p.get('point') or str(p))
                    safe_points.append(text)
                else:
                    safe_points.append(str(p))
            main_points = safe_points

        new_content["main_points"] = main_points

        # 4. 生成单页 PPTX
        output_pptx = temp_dir / f'slide_{page_index}_{action}_{uuid.uuid4().hex}.pptx'

        logger.info(f"生成单页 PPTX: {output_pptx}")
        generate_single_slide_pptx(
            content=new_content,
            output_path=output_pptx,
            slide_index=page_index,
            layout_type="auto"
        )

        # 5. 转换为 PNG
        session_id = uuid.uuid4().hex
        image_output_dir = settings.PUBLIC_DIR / "images" / session_id
        image_output_dir.mkdir(parents=True, exist_ok=True)

        logger.info(f"转换为 PNG: {image_output_dir}")
        image_result = convert_single_slide_to_image(
            pptx_path=output_pptx,
            output_dir=image_output_dir,
            page_index=0,  # 单页 PPTX 只有一页
            resolution="high"
        )

        # 6. 构建返回结果
        response = {
            "success": True,
            "content": new_content,
            "preview_url": image_result.get("url") if image_result and image_result.get("success") else None,
            "image_path": image_result.get("path") if image_result else None,
            "degraded": image_result.get("degraded", False) if image_result else True,
            "error": image_result.get("error") if image_result and not image_result.get("success") else None
        }

        # 7. 测试 JSON 序列化
        serialization_result = test_json_serialization(response, context="ai_merge_single_response")
        if not serialization_result["success"]:
            logger.error(f"[JSON序列化] ai-merge-single 失败字段: {serialization_result['problematic_fields']}")
            raise ValueError(f"JSON 序列化失败: {serialization_result['error']}")

        logger.info(f"单页处理完成: preview_url={response.get('preview_url')}")

        # 8. 注册后台清理任务（延迟清理，确保文件可被访问）
        def cleanup_temp_files():
            import time
            time.sleep(30)  # 延迟 30 秒清理，给前端足够时间下载图片
            for temp_path in [temp_pptx, output_pptx]:
                try:
                    if temp_path and temp_path.exists():
                        temp_path.unlink()
                        logger.debug(f"已清理临时文件: {temp_path}")
                except Exception as e:
                    logger.warning(f"清理临时文件失败: {temp_path}, error={e}")

        background_tasks.add_task(cleanup_temp_files)

        return JSONResponse(content=response)

    except HTTPException:
        raise
    except ValueError as e:
        logger.error(f"参数错误: {e}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"单页处理失败: {e}")
        logger.error(f"Traceback: {traceback_module.format_exc()}")
        raise HTTPException(status_code=500, detail=f"单页处理失败: {str(e)}")
