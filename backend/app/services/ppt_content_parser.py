# -*- coding: utf-8 -*-
"""
PPT 内容结构解析服务
将 PPTX 文件解析为中间数据结构供 AI 消费

设计文档: .claude-coder/plans/ppt-merge-technical-design.md#2-数据结构设计
"""

import base64
import uuid
import logging
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from ..models.ppt_structure import (
    DocumentData, SlideData, ElementData, Position, Style, Paragraph,
    TeachingContent, ElementType, SlideType, TeachingRole, to_dict
)
from .teaching_semantic_extractor import TeachingSemanticExtractor

logger = logging.getLogger(__name__)


class PptContentParser:
    """PPT 内容结构解析器"""

    # 复杂元素类型（无法 AI 合并）
    COMPLEX_SHAPE_TYPES = {
        MSO_SHAPE_TYPE.CHART,      # 图表
        MSO_SHAPE_TYPE.DIAGRAM,    # SmartArt
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,  # 嵌入对象
        MSO_SHAPE_TYPE.MEDIA,      # 音视频
    }

    # 占位符文本列表（PPT 模板中的示例文字）
    PLACEHOLDER_TEXTS = [
        # 中文占位符
        "添加标题",
        "添加标题内容",
        "添加副标题",
        "单击此处添加标题",
        "单击此处添加副标题",
        "单击此处编辑标题",
        "单击此处编辑副标题",
        "请输入标题",
        "请输入副标题",
        "点击此处添加文本",
        "在此处输入文本",
        "双击此处添加文本",
        # 英文占位符
        "click to add title",
        "click to add text",
        "click to add subtitle",
        "add title",
        "add subtitle",
        "add text",
        "enter title",
        "enter subtitle",
        "enter text",
        "title",
        "subtitle",
        "text",
    ]

    def __init__(self, max_image_size: int = 512):
        """
        初始化解析器
        Args:
            max_image_size: 图片最大尺寸（用于压缩）
        """
        self.max_image_size = max_image_size
        # 初始化教学语义提取器
        self.semantic_extractor = TeachingSemanticExtractor()

    def _is_placeholder_text(self, text: str) -> bool:
        """
        判断文本是否为占位符
        Args:
            text: 待检测文本
        Returns:
            是否为占位符文本
        """
        # 标准化：移除空格、转小写
        normalized = text.strip().lower().replace(" ", "").replace("\t", "")

        # 检查是否匹配任何占位符模式
        for placeholder in self.PLACEHOLDER_TEXTS:
            normalized_placeholder = placeholder.lower().replace(" ", "").replace("\t", "")
            # 完全匹配或包含关系（占位符文本通常很短）
            if normalized == normalized_placeholder or normalized_placeholder in normalized:
                return True

        return False

    def parse(self, file_path: Path) -> DocumentData:
        """
        解析 PPTX 文件
        Args:
            file_path: PPTX 文件路径
        Returns:
            DocumentData 中间结构
        """
        logger.info(f"开始解析 PPT: {file_path.name}")

        prs = Presentation(file_path)
        slides_data: List[SlideData] = []
        complex_element_slides: List[int] = []

        # 幻灯片尺寸
        slide_width = prs.slide_width / 914400  # EMU 转英寸
        slide_height = prs.slide_height / 914400

        # 遍历每页
        total_slides = len(prs.slides)
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = self._parse_slide(
                slide, slide_idx, slide_width, slide_height, total_slides
            )
            slides_data.append(slide_data)

            # 检测复杂元素
            if self._has_complex_elements(slide):
                complex_element_slides.append(slide_idx)
                logger.warning(f"第 {slide_idx + 1} 页包含复杂元素（图表/SmartArt）")

        # 构建文档数据
        doc_id = str(uuid.uuid4())
        document = DocumentData(
            document_id=doc_id,
            source_file=file_path.name,
            slides=slides_data,
            total_slides=len(slides_data),
            complex_elements_detected=len(complex_element_slides) > 0,
            complex_element_slides=complex_element_slides
        )

        logger.info(f"解析完成: {len(slides_data)} 页, 复杂元素页面: {complex_element_slides}")
        return document

    def _parse_slide(
        self, slide, slide_idx: int, width: float, height: float, total_slides: int
    ) -> SlideData:
        """解析单页幻灯片"""
        elements: List[ElementData] = []
        has_images = False
        has_tables = False
        title = None

        for shape_idx, shape in enumerate(slide.shapes):
            element = self._parse_shape(shape, shape_idx, width, height)
            if element:
                elements.append(element)

                # 统计
                if element.type == ElementType.IMAGE:
                    has_images = True
                elif element.type == ElementType.TABLE:
                    has_tables = True

                # 提取标题
                if element.type == ElementType.TITLE and not title:
                    title = element.text

        # 使用教学语义提取器进行语义分析
        teaching_content = self.semantic_extractor.extract_teaching_content(
            slide_idx=slide_idx,
            total_slides=total_slides,
            elements=elements,
            title=title
        )

        # 从语义提取结果获取 slide_type 和 teaching_role
        semantic_result = self.semantic_extractor.extract(
            slide_idx=slide_idx,
            total_slides=total_slides,
            elements=elements,
            title=title
        )

        # 更新 teaching_content 的 has_images 和 has_tables
        teaching_content.has_images = has_images
        teaching_content.has_tables = has_tables

        return SlideData(
            slide_index=slide_idx,
            slide_type=semantic_result.slide_type,
            teaching_role=semantic_result.teaching_role,
            elements=elements,
            teaching_content=teaching_content,
            layout_width=width,
            layout_height=height
        )

    def _parse_shape(self, shape, shape_idx: int, width: float, height: float) -> Optional[ElementData]:
        """解析单个 Shape"""
        element_id = f"elem_{shape_idx:03d}"

        # 计算位置百分比
        pos = Position(
            x_pct=round(shape.left / 914400 / width * 100, 2),
            y_pct=round(shape.top / 914400 / height * 100, 2),
            width_pct=round(shape.width / 914400 / width * 100, 2),
            height_pct=round(shape.height / 914400 / height * 100, 2)
        )

        # 根据类型解析
        shape_type = shape.shape_type

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._parse_picture(shape, element_id, pos)

        elif shape.has_table:
            return self._parse_table(shape, element_id, pos)

        elif shape.has_text_frame:
            return self._parse_text(shape, element_id, pos)

        else:
            # 其他类型（形状、占位符等）
            return ElementData(
                element_id=element_id,
                type=ElementType.SHAPE,
                position=pos,
                raw_shape_type=str(shape_type)
            )

    def _parse_picture(self, shape, element_id: str, pos: Position) -> ElementData:
        """解析图片"""
        try:
            image = shape.image
            image_bytes = image.blob

            # 压缩图片
            compressed_base64 = self._compress_image(image_bytes)

            return ElementData(
                element_id=element_id,
                type=ElementType.IMAGE,
                position=pos,
                image_base64=compressed_base64,
                image_format=image.ext,
                raw_shape_type=str(MSO_SHAPE_TYPE.PICTURE)
            )
        except Exception as e:
            logger.warning(f"图片解析失败: {e}")
            return ElementData(
                element_id=element_id,
                type=ElementType.IMAGE,
                position=pos,
                raw_shape_type=str(MSO_SHAPE_TYPE.PICTURE)
            )

    def _parse_table(self, shape, element_id: str, pos: Position) -> ElementData:
        """解析表格"""
        try:
            table = shape.table
            rows_data: List[List[str]] = []
            headers: List[str] = []

            for row_idx, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_text.append(cell_text)

                rows_data.append(row_text)

                # 第一行作为表头
                if row_idx == 0:
                    headers = row_text

            return ElementData(
                element_id=element_id,
                type=ElementType.TABLE,
                position=pos,
                table_data=rows_data,
                table_headers=headers,
                raw_shape_type=str(MSO_SHAPE_TYPE.TABLE)
            )
        except Exception as e:
            logger.warning(f"表格解析失败: {e}")
            return ElementData(
                element_id=element_id,
                type=ElementType.TABLE,
                position=pos,
                raw_shape_type=str(MSO_SHAPE_TYPE.TABLE)
            )

    def _parse_text(self, shape, element_id: str, pos: Position) -> Optional[ElementData]:
        """解析文本"""
        paragraphs: List[Paragraph] = []
        all_text: List[str] = []
        style: Optional[Style] = None

        for para in shape.text_frame.paragraphs:
            para_text = para.text.strip()
            if not para_text:
                continue

            all_text.append(para_text)

            # 提取样式（从第一个 run）
            para_style = None
            for run in para.runs:
                if run.text:
                    para_style = self._extract_style(run.font)
                    break

            paragraphs.append(Paragraph(
                text=para_text,
                style=para_style
            ))

            # 记录第一个样式作为整体样式
            if style is None:
                style = para_style

        # 判断文本类型
        full_text = "\n".join(all_text)

        # 优先检查占位符：如果是占位符文本，返回 None 过滤掉
        if self._is_placeholder_text(full_text):
            logger.debug(f"检测到占位符文本，已过滤: {full_text[:50]}...")
            return None

        element_type = self._detect_text_type(shape, full_text, pos)

        return ElementData(
            element_id=element_id,
            type=element_type,
            position=pos,
            text=full_text,
            paragraphs=paragraphs,
            style=style,
            raw_shape_type=str(shape.shape_type)
        )

    def _extract_style(self, font) -> Optional[Style]:
        """提取字体样式"""
        try:
            style = Style()

            if font.name:
                style.font_name = font.name
            if font.size:
                style.font_size = font.size.pt
            if font.bold is not None:
                style.bold = font.bold
            if font.italic is not None:
                style.italic = font.italic
            if font.underline is not None:
                style.underline = font.underline
            if font.color and hasattr(font.color, 'rgb') and font.color.rgb:
                rgb = font.color.rgb
                style.color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"

            return style
        except Exception:
            return None

    def _detect_text_type(self, shape, text: str, pos: Position) -> ElementType:
        """判断文本类型"""
        # 优先检查占位符（双重保障，_parse_text 已过滤，此处作为防御性检查）
        if self._is_placeholder_text(text):
            return ElementType.PLACEHOLDER

        # 优先判断列表
        if text.strip().startswith(('•', '-', '·', '◆', '►', '1.', '2.', '3.', '4.', '5.')):
            return ElementType.LIST_ITEM

        # 获取字体大小
        max_font_size = 0
        avg_font_size = 0
        font_count = 0

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        size_pt = run.font.size.pt
                        max_font_size = max(max_font_size, size_pt)
                        avg_font_size += size_pt
                        font_count += 1

        if font_count > 0:
            avg_font_size /= font_count

        # 根据位置和字体大小判断
        # 标题通常：位置靠上、字体大、文本短
        if pos.y_pct < 25:  # 顶部 25%
            if max_font_size > 28 or avg_font_size > 24:
                return ElementType.TITLE
            elif max_font_size > 20 or avg_font_size > 18:
                return ElementType.SUBTITLE

        # 文本较短且字体大，可能是标题
        if len(text) < 50 and max_font_size > 24:
            return ElementType.TITLE

        # 文本居中且较短，可能是副标题
        if pos.x_pct > 20 and pos.width_pct < 60 and len(text) < 100:
            if avg_font_size > 18:
                return ElementType.SUBTITLE

        return ElementType.TEXT_BODY

    def _compress_image(self, image_bytes: bytes) -> str:
        """压缩图片并返回 Base64"""
        try:
            from PIL import Image

            img = Image.open(BytesIO(image_bytes))

            # 缩放
            if img.width > self.max_image_size or img.height > self.max_image_size:
                img.thumbnail((self.max_image_size, self.max_image_size), Image.Resampling.LANCZOS)

            # 转换格式
            if img.mode in ('RGBA', 'P', 'LA'):
                img = img.convert('RGB')

            output = BytesIO()
            img.save(output, format='JPEG', quality=75)
            output.seek(0)

            return base64.b64encode(output.read()).decode('utf-8')
        except Exception as e:
            logger.warning(f"图片压缩失败: {e}")
            return base64.b64encode(image_bytes).decode('utf-8')

    def _has_complex_elements(self, slide) -> bool:
        """检测页面是否包含复杂元素"""
        for shape in slide.shapes:
            if shape.shape_type in self.COMPLEX_SHAPE_TYPES:
                return True
        return False

    def _detect_slide_type(self, slide_idx: int, elements: List[ElementData]) -> SlideType:
        """判断页面类型"""
        if slide_idx == 0:
            return SlideType.TITLE_SLIDE

        # 根据元素判断
        has_title = any(e.type == ElementType.TITLE for e in elements)

        if not has_title:
            return SlideType.SECTION_SLIDE

        return SlideType.CONTENT_SLIDE

    def _detect_teaching_role(self, slide_idx: int, elements: List[ElementData]) -> TeachingRole:
        """判断教学角色"""
        if slide_idx == 0:
            return TeachingRole.COVER

        # 根据文本内容判断
        text_content = " ".join([e.text for e in elements if e.text])

        # 简单关键词判断
        if any(kw in text_content for kw in ['练习', '习题', '试一试']):
            return TeachingRole.EXERCISE
        elif any(kw in text_content for kw in ['例题', '例', '解']):
            return TeachingRole.EXAMPLE
        elif any(kw in text_content for kw in ['总结', '小结', '回顾']):
            return TeachingRole.SUMMARY
        elif any(kw in text_content for kw in ['作业', '课后']):
            return TeachingRole.HOMEWORK
        elif any(kw in text_content for kw in ['目录', 'Contents']):
            return TeachingRole.OUTLINE

        return TeachingRole.CONCEPT

    def _extract_teaching_content(self, elements: List[ElementData], title: Optional[str]) -> TeachingContent:
        """提取教学语义内容"""
        content = TeachingContent(title=title)

        for element in elements:
            if element.type == ElementType.TEXT_BODY and element.text:
                # 简单提取：非标题文本作为要点
                content.main_points.append(element.text[:100])

            elif element.type == ElementType.LIST_ITEM and element.text:
                content.knowledge_points.append(element.text)

        return content


def parse_pptx_to_structure(file_path: Path, max_image_size: int = 512) -> Dict[str, Any]:
    """
    解析 PPTX 文件为中间结构（便捷函数）
    Args:
        file_path: PPTX 文件路径
        max_image_size: 图片最大尺寸
    Returns:
        中间结构字典
    """
    parser = PptContentParser(max_image_size=max_image_size)
    document = parser.parse(file_path)
    return to_dict(document)