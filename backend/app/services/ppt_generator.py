from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
from typing import Dict, Any, List, Optional
from io import BytesIO
import logging
import jieba

# 动画效果注入模块（条件导入）
try:
    from .pptx_animator import PPTXAnimator
    ANIMATION_AVAILABLE = True
except ImportError:
    ANIMATION_AVAILABLE = False

# 互动触发器注入模块（条件导入）
try:
    from .pptx_interactive import InteractivePPTGenerator
    INTERACTIVE_AVAILABLE = True
except ImportError:
    INTERACTIVE_AVAILABLE = False

logger = logging.getLogger(__name__)

# 预定义颜色，供类属性使用
_COLOR_RED = RGBColor(255, 107, 107)
_COLOR_ORANGE = RGBColor(255, 159, 67)
_COLOR_GOLD = RGBColor(255, 215, 0)
_COLOR_DARK_ORANGE = RGBColor(255, 140, 0)
_COLOR_TURQUOISE = RGBColor(64, 224, 208)
COLOR_TURQUOISE_2 = RGBColor(78, 205, 196)
COLOR_TURQUOISE_3 = RGBColor(102, 205, 170)
COLOR_BLUE_1 = RGBColor(52, 152, 219)
COLOR_BLUE_2 = RGBColor(68, 138, 255)
COLOR_BLUE_3 = RGBColor(63, 81, 181)
COLOR_BLUE_4 = RGBColor(106, 90, 205)
COLOR_BLUE_5 = RGBColor(41, 128, 185)
COLOR_GREEN_1 = RGBColor(67, 160, 71)
COLOR_GREEN_2 = RGBColor(39, 174, 96)
COLOR_GREEN_3 = RGBColor(102, 187, 106)
COLOR_BLUE_6 = RGBColor(33, 150, 243)
COLOR_GREEN_4 = RGBColor(76, 175, 80)

# 学科主题色映射
SUBJECT_COLORS = {
    "math": {"primary": COLOR_BLUE_5, "secondary": COLOR_BLUE_1, "name": "数学蓝"},
    "english": {"primary": COLOR_GREEN_2, "secondary": COLOR_GREEN_3, "name": "英语绿"},
    "chinese": {"primary": _COLOR_RED, "secondary": _COLOR_GOLD, "name": "语文红"},
    "physics": {"primary": _COLOR_ORANGE, "secondary": COLOR_BLUE_5, "name": "物理橙"},
    "chemistry": {"primary": COLOR_BLUE_3, "secondary": COLOR_GREEN_1, "name": "化学蓝"},
    "biology": {"primary": COLOR_GREEN_1, "secondary": COLOR_TURQUOISE_2, "name": "生物绿"},
    "history": RGBColor(128, 0, 128),
    "politics": RGBColor(186, 85, 211),
    "geography": RGBColor(34, 139, 34),
    "general": {"primary": COLOR_BLUE_6, "secondary": COLOR_TURQUOISE_2, "name": "通用蓝"},
}


def add_pinyin_text(paragraph, text: str, font_size: int, is_title: bool = False):
    """
    添加带拼音的文本（简化版：仅对常见字添加拼音标注）
    实际项目中建议使用 pypinyin 库进行完整拼音转换
    Args:
        paragraph: 段落对象
        text: 要添加的文本
        font_size: 字号
        is_title: 是否标题
    """
    # 简化实现：直接添加文本，拼音标注需要额外的 pypinyin 库支持
    # 在实际部署时，可以安装 pypinyin: pip install pypinyin
    # 然后使用：from pypinyin import lazy_pinyin
    paragraph.text = text
    if is_title:
        paragraph.font.size = Pt(font_size + 4)
        paragraph.font.bold = True
    else:
        paragraph.font.size = Pt(font_size)


class PPTPageType:
    """PPT 页面类型"""
    COVER = "封面页"
    TOC = "目录页"
    CONTENT = "知识点讲解页"
    INTERACTION = "互动问答页"
    EXERCISE = "课堂练习页"
    SUMMARY = "总结回顾页"

    # 新增页面类型（优化版：增加页面类型多样性）
    DIAGRAM = "图示页"  # 流程图、结构图、概念图
    TABLE = "表格页"  # 对比表、分类表、数据表
    COMPARISON = "对比分析页"  # 对比两种或多种概念

    # 英语学科特殊页面类型
    VOCABULARY = "单词学习页"
    GRAMMAR = "语法讲解页"
    DIALOGUE = "情景对话页"
    ANALYSIS = "课文分析页"

    # 英语学科深度增强页面类型（feat-033）
    WORD_ROOT = "词根词缀页"  # 词根/前缀/后缀分析
    WORD_STRUCTURE = "构词法图解页"  # 单词构成树状图
    ETYMOLOGY = "词源故事页"  # 单词来源和文化背景
    COLLOCATION = "固定搭配页"  # 短语/搭配/惯用语
    VOCAB_NETWORK = "词汇网络页"  # 同义词/反义词/上义词/下义词
    COLLOCATION_MATRIX = "搭配矩阵页"  # 搭配表格
    SYNTAX_TREE = "语法树页"  # 句子成分分析
    SYNTAX_STRUCTURE = "句法结构页"  # 树状图展示
    CLAUSE_ANALYSIS = "从句分析页"  # 各类从句讲解
    SPECIAL_SENTENCE = "特殊句型页"  # 被动/倒装/强调/虚拟
    TENSE_TIMELINE = "时态时间轴页"  # 时态可视化
    TENSE_OVERVIEW = "时态概览页"  # 16 种时态总览
    TENSE_COMPARISON = "时态对比页"  # 易混时态对比

    # 数学学科特殊页面类型
    CONCEPT = "概念引入页"  # 直观→抽象概念形成
    FORMULA_DERIVATION = "公式推导页"  # 逐步标注推导过程
    EXAMPLE_PROBLEM = "例题讲解页"  # 四步法（审题→分析→解题→反思）
    VARIATION_PRACTICE = "变式训练页"  # 梯度练习
    COMMON_MISTAKES = "易错警示页"  # 错误 vs 正确对比
    MISTAKE_ANALYSIS = "错题分析页"  # feat-041: 常见错误 vs 正确方法对比分析

    # 语文教学习页面类型（feat-029）
    CHARACTER_LEARNING = "生字学习页"  # 语文专属：字 + 拼音 + 组词 + 造句
    POEM_APPRECIATION = "古诗鉴赏页"  # 语文专属：原文 + 注释 + 赏析
    READING_COMPREHENSION = "阅读理解页"  # 语文专属：段落分析 + 主题提炼
    WRITING_GUIDANCE = "作文指导页"  # 语文专属：审题 + 结构 + 素材
    PARAGRAPH_ANALYSIS = "段落分析页"  # 语文专属：层次划分 + 关键句
    RHETORIC_DEVICE = "修辞手法页"  # 语文专属：比喻/拟人/排比等

    # 先行组织者与脚手架策略（feat-037）
    BRIDGE = "概念桥接页"  # 连接新旧知识：已知 X→新知 Y→关系说明

    # 互动课堂页面类型（feat-017）
    QUIZ = "互动选择题页"  # ABCD 选项选择题
    CLICK_TO_REVEAL = "点击显示答案页"  # 点击显示隐藏内容
    MATCHING = "拖拽匹配游戏页"  # 单词 - 释义配对
    QUIZ_GAME = "随堂测验页"  # 多题测验
    FILL_BLANK = "填空题页"  # 填空练习

    # 理科实验页面类型（feat-031 物理/化学/生物）
    EXPERIMENT_DESIGN = "实验设计页"  # 探究目的、器材、步骤
    EXPERIMENT_OBSERVATION = "现象观察页"  # 实验现象记录
    DATA_ANALYSIS = "数据分析页"  # 数据处理、规律发现
    LAW_SUMMARY = "规律归纳页"  # 从实验提炼规律
    APPLICATION_TRANSFER = "应用迁移页"  # 联系实际应用
    EXPERIMENT_STEPS = "实验步骤页"  # 仪器、药品、操作流程
    REACTION_PHENOMENA = "反应现象页"  # 颜色变化、气体、沉淀
    CHEMICAL_EQUATION = "化学方程式页"  # 方程式书写、配平
    MICROSCOPIC_EXPLANATION = "微观解释页"  # 分子/原子/离子模型
    PERIODIC_TREND = "元素周期律页"  # 周期表位置、性质递变
    STRUCTURE_OBSERVATION = "结构观察页"  # 细胞/组织/器官结构
    FUNCTION_ANALYSIS = "功能分析页"  # 结构与功能适应
    PROCESS_DESCRIPTION = "过程描述页"  # 生命活动过程
    SYSTEM_THINKING = "系统思维页"  # 层次关系、整体协调
    EXPERIMENT_INQUIRY = "实验探究页"  # 观察实验、对照实验

    # 文科教学页面类型（feat-032 历史/政治/地理）
    TIMELINE = "时间轴页"  # 历史专属：时序关系、发展阶段
    BACKGROUND_ANALYSIS = "背景分析页"  # 历史专属：时代背景
    CAUSE_ANALYSIS = "因果分析页"  # 历史专属：原因探究
    HISTORICAL_MATERIAL = "史料解读页"  # 历史专属：原始史料分析
    PROCESS_NARRATIVE = "过程叙述页"  # 历史专属：事件经过
    IMPACT_EVALUATION = "影响评价页"  # 历史专属：历史影响
    MULTI_PERSPECTIVE = "多视角评价页"  # 历史专属：多元评价
    HISTORICAL_REVELATION = "历史启示页"  # 历史专属：以史为鉴
    CONCEPT_DEFINITION = "概念界定页"  # 政治专属：概念内涵外延
    PRINCIPLE_EXPLANATION = "原理阐述页"  # 政治专属：理论原理
    CASE_ANALYSIS = "案例分析页"  # 政治专属：理论联系实际
    POLICY_INTERPRETATION = "政策解读页"  # 政治专属：政策分析
    VALUE_GUIDANCE = "价值引领页"  # 政治专属：价值观教育
    DIALECTICAL_ANALYSIS = "辩证分析页"  # 政治专属：辩证思维
    LOCATION = "位置定位页"  # 地理专属：经纬度、海陆位置
    NATURAL_ENVIRONMENT = "自然环境页"  # 地理专属：地形/气候/水文
    HUMAN_FEATURES = "人文特征页"  # 地理专属：人口/城市/产业
    REGIONAL_COMPARISON = "区域比较页"  # 地理专属：区域差异分析
    HUMAN_LAND_RELATION = "人地关系页"  # 地理专属：人地协调
    MAP_ANALYSIS = "地图分析页"  # 地理专属：地图阅读

    # 问题链与元认知提示系统（feat-038）
    QUESTION_CHAIN = "问题链递进页"  # 递进式问题序列：是什么→为什么→怎么用→如果→还有
    METACOGNITIVE_BUBBLE = "思考气泡页"  # 元认知提示：云朵形状思考气泡 + 特殊配色
    REFLECTION_CHECKPOINT = "反思 checkpoint 页"  # 学习 checkpoints，促进自我监控

    # 情境创设自动生成（feat-039）
    SCENARIO_IMPORT = "情境导入页"  # 贴近学生生活的情境导入：情境描述 + 引导问题 + 视觉建议


class PPTStyle:
    """PPT 样式配置（优化版：细化年级自适应规则）"""

    # 风格配置（三种 PPT 风格）
    STYLE_CONFIG = {
        "fun": {
            "name": "活泼趣味",
            "description": "鲜艳暖色配色，活泼生动，适合低年级",
            "color_scheme": "warm",  # 暖色系
            "font_family": "rounded",  # 圆体
            "layout_density": "spacious",  # 宽松布局，大量留白
            "decoration_level": "high",  # 高装饰性
            "border_radius": 12,  # 大圆角
            "use_gradients": True,  # 使用渐变
            "animation_style": "bouncy",  # 活泼动画
        },
        "simple": {
            "name": "简约清晰",
            "description": "灰蓝冷色配色，简洁专业，适合高年级",
            "color_scheme": "cool",  # 冷色系
            "font_family": "sans-serif",  # 黑体
            "layout_density": "compact",  # 紧凑布局
            "decoration_level": "low",  # 低装饰性
            "border_radius": 4,  # 小圆角
            "use_gradients": False,  # 不使用渐变
            "animation_style": "fade",  # 淡入动画
        },
        "theme": {
            "name": "学科主题",
            "description": "根据学科自动选择主题色",
            "color_scheme": "dynamic",  # 动态配色
            "font_family": "sans-serif",  # 黑体
            "layout_density": "balanced",  # 平衡布局
            "decoration_level": "medium",  # 中等装饰性
            "border_radius": 8,  # 中圆角
            "use_gradients": False,
            "animation_style": "slide",  # 滑动动画
        },
    }

    # 年级组配置（三个学段）
    GRADE_GROUPS = {
        "elementary_low": {
            "name": "小学低年级（1-3 年级）",
            "grades": ["1", "2", "3"],
            "description": "特大字号、鲜艳配色、大量留白、趣味化表达",
            "teaching_strategy": "直观形象教学为主，多采用图片、动画、游戏等形式， attention span 短，每页内容少而精",
            "interaction_style": "简单问答、跟读、指认、小游戏",
        },
        "elementary_high": {
            "name": "小学高年级（4-6 年级）",
            "grades": ["4", "5", "6"],
            "description": "较大字号、明快配色、图文并茂、互动游戏式问答",
            "teaching_strategy": "逐步过渡到抽象思维，增加逻辑推理内容，培养自主学习能力",
            "interaction_style": "小组讨论、角色扮演、竞赛游戏、思维导图",
        },
        "middle": {
            "name": "初中（7-9 年级）",
            "grades": ["7", "8", "9"],
            "description": "标准字号、稳重配色、知识体系化、重点标注",
            "teaching_strategy": "系统化知识讲解，注重知识间的联系，培养批判性思维和综合运用能力",
            "interaction_style": "探究性问题、辩论、案例分析、项目式学习",
        },
        "high_school": {
            "name": "高中（10-12 年级）",
            "grades": ["10", "11", "12"],
            "description": "学术风格、严谨配色、高信息密度、高考导向",
            "teaching_strategy": "深度知识讲解，强调学科核心素养，注重知识迁移和综合运用，强化应试能力",
            "interaction_style": "学术研讨、高考真题分析、变式训练、专题探究",
        },
    }

    # 年级样式配置（细化到每个年级）
    GRADE_STYLES = {
        # 小学低年级（1-3 年级）：特大字号、鲜艳配色
        "1": {
            "font_size": 36,  # 增大字号
            "title_size": 48,
            "primary": _COLOR_RED,
            "secondary": COLOR_TURQUOISE_2,
            "group": "elementary_low",
            "need_pinyin": True,  # 语文需要拼音标注
            "max_content_lines": 3,  # 每页最多内容行数
            "use_icons": True,  # 使用图标辅助理解
            "animation_suggestion": "使用活泼的进入动画",
        },
        "2": {
            "font_size": 34,
            "title_size": 44,
            "primary": _COLOR_ORANGE,
            "secondary": COLOR_TURQUOISE_3,
            "group": "elementary_low",
            "need_pinyin": True,
            "max_content_lines": 4,
            "use_icons": True,
            "animation_suggestion": "使用活泼的进入动画",
        },
        "3": {
            "font_size": 30,
            "title_size": 40,
            "primary": _COLOR_GOLD,
            "secondary": _COLOR_DARK_ORANGE,
            "group": "elementary_low",
            "need_pinyin": True,
            "max_content_lines": 4,
            "use_icons": True,
            "animation_suggestion": "使用适中的进入动画",
        },
        # 小学高年级（4-6 年级）：较大字号、明快配色
        "4": {
            "font_size": 26,
            "title_size": 36,
            "primary": _COLOR_TURQUOISE,
            "secondary": COLOR_TURQUOISE_2,
            "group": "elementary_high",
            "need_pinyin": False,
            "max_content_lines": 5,
            "use_icons": True,
            "animation_suggestion": "使用简洁的进入动画",
        },
        "5": {
            "font_size": 24,
            "title_size": 34,
            "primary": COLOR_BLUE_1,
            "secondary": COLOR_BLUE_2,
            "group": "elementary_high",
            "need_pinyin": False,
            "max_content_lines": 6,
            "use_icons": False,
            "animation_suggestion": "使用简洁的进入动画",
        },
        "6": {
            "font_size": 22,
            "title_size": 32,
            "primary": COLOR_BLUE_3,
            "secondary": COLOR_BLUE_4,
            "group": "elementary_high",
            "need_pinyin": False,
            "max_content_lines": 6,
            "use_icons": False,
            "animation_suggestion": "使用淡入动画",
        },
        # 初中（7-9 年级）：标准字号、稳重配色
        "7": {
            "font_size": 20,
            "title_size": 30,
            "primary": COLOR_BLUE_5,
            "secondary": COLOR_GREEN_1,
            "group": "middle",
            "need_pinyin": False,
            "max_content_lines": 7,
            "use_icons": False,
            "animation_suggestion": "使用淡入或擦除动画",
        },
        "8": {
            "font_size": 18,
            "title_size": 28,
            "primary": COLOR_GREEN_2,
            "secondary": COLOR_GREEN_3,
            "group": "middle",
            "need_pinyin": False,
            "max_content_lines": 8,
            "use_icons": False,
            "animation_suggestion": "使用简洁的切换动画",
        },
        "9": {
            "font_size": 18,
            "title_size": 26,
            "primary": COLOR_BLUE_6,
            "secondary": COLOR_GREEN_4,
            "group": "middle",
            "need_pinyin": False,
            "max_content_lines": 8,
            "use_icons": False,
            "animation_suggestion": "使用无动画或简单切换",
        },
        # 高中（10-12 年级）：学术风格、严谨配色
        "10": {
            "font_size": 18,
            "title_size": 26,
            "primary": COLOR_BLUE_3,
            "secondary": COLOR_BLUE_5,
            "group": "high_school",
            "need_pinyin": False,
            "max_content_lines": 8,
            "use_icons": False,
            "animation_suggestion": "使用专业克制的切换动画",
        },
        "11": {
            "font_size": 16,
            "title_size": 24,
            "primary": COLOR_GREEN_2,
            "secondary": COLOR_BLUE_3,
            "group": "high_school",
            "need_pinyin": False,
            "max_content_lines": 9,
            "use_icons": False,
            "animation_suggestion": "使用简洁动画或无动画",
        },
        "12": {
            "font_size": 16,
            "title_size": 24,
            "primary": COLOR_BLUE_6,
            "secondary": COLOR_GREEN_1,
            "group": "high_school",
            "need_pinyin": False,
            "max_content_lines": 10,
            "use_icons": False,
            "animation_suggestion": "使用无动画，专注内容呈现",
        },
    }

    @classmethod
    def get_style(cls, grade: str) -> Dict[str, Any]:
        """获取年级对应的样式"""
        return cls.GRADE_STYLES.get(grade, cls.GRADE_STYLES["6"])

    @classmethod
    def get_style_config(cls, style: str = "simple") -> Dict[str, Any]:
        """获取 PPT 风格配置"""
        return cls.STYLE_CONFIG.get(style, cls.STYLE_CONFIG["simple"])

    @classmethod
    def get_subject_color(cls, subject: str) -> Dict[str, Any]:
        """获取学科主题色"""
        color_info = SUBJECT_COLORS.get(subject, SUBJECT_COLORS["general"])
        # 处理某些学科直接定义了颜色值的情况
        if isinstance(color_info, RGBColor):
            return {"primary": color_info, "secondary": SUBJECT_COLORS["general"]["secondary"], "name": subject}
        return color_info

    @classmethod
    def get_grade_group(cls, grade: str) -> str:
        """获取年级所属组别"""
        style = cls.get_style(grade)
        return style.get("group", "elementary_high")

    @classmethod
    def need_pinyin(cls, grade: str, subject: str = "chinese") -> bool:
        """判断是否需要拼音标注（仅小学低年级语文）"""
        if subject != "chinese":
            return False
        style = cls.get_style(grade)
        return style.get("need_pinyin", False)


class PPTGenerator:
    """PPT 文件生成器"""

    def __init__(self):
        """初始化 PPT 生成器"""
        pass

    def _setup_text_frame(self, text_frame, vertical_anchor=MSO_VERTICAL_ANCHOR.TOP, add_margins=True):
        """
        统一设置文本框属性，防止文本溢出

        Args:
            text_frame: 文本框对象
            vertical_anchor: 垂直对齐方式（默认顶部对齐）
            add_margins: 是否添加边距（默认添加）
        """
        text_frame.word_wrap = True  # 自动换行
        text_frame.vertical_anchor = vertical_anchor

        if add_margins:
            # 设置边距（防止文字贴边）
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)

        # 注意：不设置 auto_size，因为它可能导致 XML 命名空间错误
        # 文本溢出问题主要通过 word_wrap 和 margins 解决

    def _copy_slide_from_pptx(self, target_prs: Presentation, source_pptx_path: str, slide_index: int):
        """
        从源 PPTX 复制页面到目标 PPT

        Args:
            target_prs: 目标 PPT（Presentation 对象）
            source_pptx_path: 源 PPTX 文件路径
            slide_index: 源页码（0-indexed）
        """
        from pptx import Presentation as Prs

        source_prs = Prs(source_pptx_path)
        if slide_index >= len(source_prs.slides):
            raise ValueError(f"页码超出范围：{slide_index} >= {len(source_prs.slides)}")

        source_slide = source_prs.slides[slide_index]

        # 使用 blank layout 创建新页面
        blank_layout = target_prs.slide_layouts[6]  # 6 = Blank
        target_slide = target_prs.slides.add_slide(blank_layout)

        # 复制源页面的所有形状
        for shape in source_slide.shapes:
            if shape.shape_type == 14:  # MSO_SHAPE_TYPE.PICTURE
                # 复制图片
                try:
                    img_bytes = shape.image.blob
                    # 计算位置
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    target_slide.shapes.add_picture(
                        BytesIO(img_bytes), left, top, width, height
                    )
                except Exception as e:
                    logger.warning(f"复制图片失败：{e}")
            elif hasattr(shape, "text") and shape.text:
                # 复制文本框
                try:
                    # 创建新文本框
                    new_shape = target_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    frame = new_shape.text_frame
                    frame.clear()

                    # 复制段落
                    for i, para in enumerate(shape.text_frame.paragraphs):
                        if i == 0:
                            frame.paragraphs[0].text = para.text
                        else:
                            frame.add_paragraph().text = para.text
                except Exception as e:
                    logger.warning(f"复制文本框失败：{e}")

        logger.info(f"成功复制页面 {slide_index} 从 {source_pptx_path}")

    def _add_slide_from_snapshot(self, target_prs: Presentation, snapshot: Dict[str, Any], font_size: int):
        """
        从 content_snapshot 重建页面

        Args:
            target_prs: 目标 PPT
            snapshot: AI 修改的内容快照（中间结构数据）
            font_size: 正文字号
        """
        # 使用 blank layout
        blank_layout = target_prs.slide_layouts[6]
        slide = target_prs.slides.add_slide(blank_layout)

        # 适配嵌套格式：content_snapshot 可能是 {action, polished_content: {...}} 形式
        action = snapshot.get("action", "")
        if action and action in ("polish", "expand", "rewrite", "extract"):
            title, content_items = self._extract_content_from_snapshot(snapshot, action)
            elements = [{"type": "text_body", "content": item} if isinstance(item, str) else item for item in content_items]
            main_points = []
        else:
            title = snapshot.get("title", "")
            elements = snapshot.get("elements", [])
            main_points = snapshot.get("main_points", [])

        y_position = Inches(0.5)

        # 添加标题
        if title:
            title_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(font_size + 4)
            title_frame.paragraphs[0].font.bold = True
            y_position = Inches(1.5)

        # 添加要点
        for point in main_points:
            para_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.8))
            para_frame = para_shape.text_frame
            para_frame.text = f"• {point}"
            para_frame.paragraphs[0].font.size = Pt(font_size)
            y_position += Inches(0.7)

        # 添加其他元素
        for elem in elements:
            elem_type = elem.get("type", "text")
            content = elem.get("content", "")

            if elem_type in ("text_body", "list_item", "title"):
                elem_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.6))
                elem_frame = elem_shape.text_frame
                elem_frame.text = content
                elem_frame.paragraphs[0].font.size = Pt(font_size)
                y_position += Inches(0.6)

        logger.info(f"成功从 snapshot 重建页面")

    def generate_single_slide_pptx(
        self,
        content_snapshot: Dict[str, Any],
        output_path: Path,
        font_size: int = 24,
        primary_color: RGBColor = COLOR_BLUE_5
    ) -> Path:
        """
        从 AI 返回的 content_snapshot 生成单页 PPTX

        Args:
            content_snapshot: AI 修改的内容快照，格式根据 action 类型不同：
                - polish: {action, polished_content: {title, main_points}}
                - expand: {action, expanded_content: {title, expanded_points, original_points}}
                - rewrite: {action, rewritten_content: {title, main_content}}
                - extract: {action, extracted_knowledge: {core_concepts, formulas, methods}}
            output_path: 输出文件路径
            font_size: 正文字号
            primary_color: 主色调

        Returns:
            生成的 PPTX 文件路径
        """
        action = content_snapshot.get("action", "polish")

        # 根据 action 类型提取内容
        title, content_items = self._extract_content_from_snapshot(content_snapshot, action)

        # 创建单页 PPT
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 使用 blank layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        y_position = Inches(0.5)

        # 添加标题
        if title:
            title_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(font_size + 6)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = primary_color
            y_position = Inches(1.5)

        # 添加内容项
        for item in content_items:
            if isinstance(item, dict):
                # 结构化内容项
                item_type = item.get("type", "text")
                item_content = item.get("content", "")

                if item_type == "heading":
                    # 小标题
                    heading_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.6))
                    heading_frame = heading_shape.text_frame
                    heading_frame.text = item_content
                    heading_frame.paragraphs[0].font.size = Pt(font_size + 2)
                    heading_frame.paragraphs[0].font.bold = True
                    y_position += Inches(0.6)
                elif item_type == "formula":
                    # 公式（特殊样式）
                    formula_shape = slide.shapes.add_textbox(Inches(0.8), y_position, Inches(8.4), Inches(0.6))
                    formula_frame = formula_shape.text_frame
                    formula_frame.text = f"📐 {item_content}"
                    formula_frame.paragraphs[0].font.size = Pt(font_size)
                    formula_frame.paragraphs[0].font.italic = True
                    y_position += Inches(0.6)
                else:
                    # 普通文本
                    text_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.6))
                    text_frame = text_shape.text_frame
                    text_frame.text = f"• {item_content}"
                    text_frame.paragraphs[0].font.size = Pt(font_size)
                    y_position += Inches(0.5)
            else:
                # 简单字符串内容
                text_shape = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.6))
                text_frame = text_shape.text_frame
                text_frame.text = f"• {item}"
                text_frame.paragraphs[0].font.size = Pt(font_size)
                y_position += Inches(0.5)

        # 保存文件
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

        logger.info(f"成功生成单页 PPTX: {output_path}")
        return output_path

    def _extract_content_from_snapshot(
        self,
        content_snapshot: Dict[str, Any],
        action: str
    ) -> tuple:
        """
        从 content_snapshot 提取标题和内容项

        Args:
            content_snapshot: AI 返回的内容快照
            action: 动作类型

        Returns:
            (title, content_items) 元组
        """
        title = ""
        content_items = []

        if action == "polish":
            # 润色结果
            polished = content_snapshot.get("polished_content", {})
            title = polished.get("title", "")
            main_points = polished.get("main_points", [])
            content_items = main_points

        elif action == "expand":
            # 扩展结果
            expanded = content_snapshot.get("expanded_content", {})
            title = expanded.get("title", "")

            # 合并原有要点和扩展要点
            original = expanded.get("original_points", [])
            expanded_points = expanded.get("expanded_points", [])
            new_examples = expanded.get("new_examples", [])
            additional = expanded.get("additional_content", "")

            # 添加原有要点
            if original:
                content_items.append({"type": "heading", "content": "原有内容"})
                content_items.extend(original)

            # 添加扩展要点
            if expanded_points:
                content_items.append({"type": "heading", "content": "扩展内容"})
                content_items.extend(expanded_points)

            # 添加新例子
            if new_examples:
                content_items.append({"type": "heading", "content": "新增示例"})
                content_items.extend(new_examples)

            # 添加额外内容
            if additional:
                content_items.append({"type": "heading", "content": "补充说明"})
                content_items.append(additional)

        elif action == "rewrite":
            # 改写结果
            rewritten = content_snapshot.get("rewritten_content", {})
            title = rewritten.get("title", "")
            main_content = rewritten.get("main_content", "")
            style_features = rewritten.get("style_features", [])

            if main_content:
                # 将主要内容按段落分割
                paragraphs = main_content.split("\n\n") if "\n\n" in main_content else [main_content]
                content_items.extend(paragraphs)

            if style_features:
                content_items.append({"type": "heading", "content": "风格特点"})
                content_items.extend(style_features)

        elif action == "extract":
            # 提取结果
            extracted = content_snapshot.get("extracted_knowledge", {})
            title = "知识点提取"

            # 核心概念
            concepts = extracted.get("core_concepts", [])
            if concepts:
                content_items.append({"type": "heading", "content": "核心概念"})
                for concept in concepts:
                    if isinstance(concept, dict):
                        concept_name = concept.get("concept", "")
                        concept_def = concept.get("definition", "")
                        content_items.append(f"{concept_name}：{concept_def}")
                    else:
                        content_items.append(concept)

            # 公式
            formulas = extracted.get("formulas", [])
            if formulas:
                content_items.append({"type": "heading", "content": "重要公式"})
                for formula in formulas:
                    if isinstance(formula, dict):
                        formula_name = formula.get("name", "")
                        formula_expr = formula.get("formula", "")
                        content_items.append({"type": "formula", "content": f"{formula_name}: {formula_expr}"})
                    else:
                        content_items.append({"type": "formula", "content": formula})

            # 方法
            methods = extracted.get("methods", [])
            if methods:
                content_items.append({"type": "heading", "content": "解题方法"})
                for method in methods:
                    if isinstance(method, dict):
                        method_name = method.get("name", "")
                        steps = method.get("steps", [])
                        content_items.append(f"{method_name}：{' → '.join(steps)}")
                    else:
                        content_items.append(method)

            # 易错点
            mistakes = extracted.get("common_mistakes", [])
            if mistakes:
                content_items.append({"type": "heading", "content": "易错提醒"})
                for mistake in mistakes:
                    if isinstance(mistake, dict):
                        wrong = mistake.get("mistake", "")
                        correct = mistake.get("correction", "")
                        content_items.append(f"❌ {wrong} → ✅ {correct}")
                    else:
                        content_items.append(mistake)

        else:
            # 未知 action，尝试通用提取
            title = content_snapshot.get("title", "")
            content_items = content_snapshot.get("main_points", content_snapshot.get("content", []))

        return title, content_items

    def generate(
        self,
        content: Dict[str, Any],
        output_path: Path,
        grade: str = "6",
        style: str = "simple",
        subject: str = "general"
    ) -> Path:
        """
        生成 PPT 文件
        Args:
            content: PPT 内容数据
            output_path: 输出文件路径
            grade: 年级
            style: PPT 风格
            subject: 学科（用于拼音标注等特性）
        Returns:
            生成的 PPT 文件路径
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # 获取年级样式配置
            grade_style = PPTStyle.get_style(grade)
            title_size = grade_style["title_size"]
            content_size = grade_style["font_size"]
            grade_group = grade_style.get("group", "elementary_high")
            use_pinyin = PPTStyle.need_pinyin(grade, subject)

            # 获取 PPT 风格配置
            ppt_style_config = PPTStyle.get_style_config(style)

            # 确定主色和辅助色：根据风格决定
            if style == "theme" and subject in SUBJECT_COLORS:
                # 学科主题模式：使用学科颜色
                subject_color = PPTStyle.get_subject_color(subject)
                primary_color = subject_color["primary"]
                secondary_color = subject_color["secondary"]
            else:
                # 其他模式：使用年级默认颜色
                primary_color = grade_style["primary"]
                secondary_color = grade_style["secondary"]

                # 活泼风格：调整为更鲜艳的暖色
                if style == "fun":
                    primary_color = _COLOR_RED if grade_group == "elementary_low" else _COLOR_ORANGE
                    secondary_color = COLOR_TURQUOISE_2 if grade_group == "elementary_low" else COLOR_TURQUOISE_3
                # 简约风格：使用冷色调
                elif style == "simple":
                    primary_color = COLOR_BLUE_5 if grade_group == "middle" else COLOR_BLUE_1
                    secondary_color = COLOR_GREEN_1 if grade_group == "middle" else COLOR_BLUE_2

            # 生成封面页
            self._add_cover_slide(
                prs,
                content.get("title", "教学 PPT"),
                title_size,
                primary_color,
                use_pinyin
            )

            # 生成目录页
            if any(s.get("page_type") == PPTPageType.TOC for s in content.get("slides", [])):
                self._add_toc_slide(
                    prs,
                    content.get("slides", []),
                    title_size,
                    content_size,
                    secondary_color
                )

            # 生成内容页
            for slide_data in content.get("slides", []):
                page_type = slide_data.get("page_type", PPTPageType.CONTENT)

                if page_type == PPTPageType.CONTENT:
                    self._add_content_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.INTERACTION:
                    self._add_interaction_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                elif page_type == PPTPageType.EXERCISE:
                    self._add_exercise_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                elif page_type == PPTPageType.SUMMARY:
                    self._add_summary_slide(
                        prs,
                        content,
                        title_size,
                        content_size,
                        primary_color
                    )
                # 新增页面类型支持
                elif page_type == PPTPageType.DIAGRAM:
                    self._add_diagram_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.TABLE:
                    self._add_table_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                elif page_type == PPTPageType.COMPARISON:
                    self._add_mistake_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 英语学科专属页面类型
                elif page_type == PPTPageType.VOCABULARY:
                    self._add_vocabulary_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.GRAMMAR:
                    self._add_grammar_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                elif page_type == PPTPageType.DIALOGUE:
                    self._add_dialogue_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.ANALYSIS:
                    self._add_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                # 英语学科深度增强页面类型（feat-033）
                elif page_type == PPTPageType.WORD_ROOT:
                    self._add_word_root_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.WORD_STRUCTURE:
                    self._add_word_structure_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.ETYMOLOGY:
                    self._add_etymology_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                elif page_type == PPTPageType.COLLOCATION:
                    self._add_collocation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.VOCAB_NETWORK:
                    self._add_vocab_network_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.COLLOCATION_MATRIX:
                    self._add_collocation_matrix_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.SYNTAX_TREE:
                    self._add_syntax_tree_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.SYNTAX_STRUCTURE:
                    self._add_syntax_structure_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.CLAUSE_ANALYSIS:
                    self._add_clause_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                elif page_type == PPTPageType.SPECIAL_SENTENCE:
                    self._add_special_sentence_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.TENSE_TIMELINE:
                    self._add_tense_timeline_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.TENSE_OVERVIEW:
                    self._add_tense_overview_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.TENSE_COMPARISON:
                    self._add_tense_comparison_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 语文学科专属页面类型（feat-029）
                elif page_type == PPTPageType.CHARACTER_LEARNING:
                    # 生字学习页：字 + 拼音 + 组词 + 造句
                    self._add_character_learning_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.POEM_APPRECIATION:
                    # 古诗鉴赏页：原文 + 注释 + 赏析
                    self._add_poem_appreciation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.READING_COMPREHENSION:
                    # 阅读理解页：段落分析 + 主题提炼
                    self._add_reading_comprehension_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.WRITING_GUIDANCE:
                    # 作文指导页：审题 + 结构 + 素材
                    self._add_writing_guidance_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.PARAGRAPH_ANALYSIS:
                    # 段落分析页：层次划分 + 关键句
                    self._add_paragraph_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.RHETORIC_DEVICE:
                    # 修辞手法页：比喻/拟人/排比等
                    self._add_rhetoric_device_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 数学学科专属页面类型
                elif page_type == PPTPageType.CONCEPT:
                    # 概念引入页：复用内容页逻辑，增加概念字段处理
                    self._add_content_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.FORMULA_DERIVATION:
                    # 公式推导页：复用内容页逻辑，公式步骤在 content 中展示
                    self._add_content_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.EXAMPLE_PROBLEM:
                    # 例题讲解页：复用内容页逻辑
                    self._add_content_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.VARIATION_PRACTICE:
                    # 变式训练页：复用练习页逻辑
                    self._add_exercise_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )
                elif page_type == PPTPageType.COMMON_MISTAKES:
                    # 易错警示页（向后兼容）：调用错题分析页方法
                    self._add_mistake_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        _COLOR_RED,
                        COLOR_GREEN_1
                    )
                elif page_type == PPTPageType.MISTAKE_ANALYSIS:  # feat-041 智能错题分析页
                    # 错题分析页：左右分栏（红色错误 vs 绿色正确）
                    self._add_mistake_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        _COLOR_RED,  # 红色表示错误
                        COLOR_GREEN_1  # 绿色表示正确
                    )
                # 先行组织者与脚手架策略（feat-037）
                elif page_type == PPTPageType.BRIDGE:
                    # 概念桥接页：连接新旧知识的双色对比布局
                    self._add_bridge_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 情境创设自动生成（feat-039）
                elif page_type == PPTPageType.SCENARIO_IMPORT:
                    # 情境导入页：贴近学生生活的情境导入
                    self._add_scenario_import_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 互动课堂页面类型（feat-017）
                elif page_type == PPTPageType.QUIZ:
                    self._add_quiz_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.CLICK_TO_REVEAL:
                    self._add_click_to_reveal_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.MATCHING:
                    self._add_matching_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.QUIZ_GAME:
                    self._add_quiz_game_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.FILL_BLANK:
                    self._add_fill_blank_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 理科实验页面类型（feat-031 物理/化学/生物）
                elif page_type == PPTPageType.EXPERIMENT_DESIGN:
                    self._add_experiment_design_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.EXPERIMENT_OBSERVATION:
                    self._add_experiment_observation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.DATA_ANALYSIS:
                    self._add_data_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.LAW_SUMMARY:
                    self._add_law_summary_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.APPLICATION_TRANSFER:
                    self._add_application_transfer_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.EXPERIMENT_STEPS:
                    self._add_experiment_steps_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.REACTION_PHENOMENA:
                    self._add_reaction_phenomena_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.CHEMICAL_EQUATION:
                    self._add_chemical_equation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.MICROSCOPIC_EXPLANATION:
                    self._add_microscopic_explanation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.PERIODIC_TREND:
                    self._add_periodic_trend_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.STRUCTURE_OBSERVATION:
                    self._add_structure_observation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.FUNCTION_ANALYSIS:
                    self._add_function_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.PROCESS_DESCRIPTION:
                    self._add_process_description_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.SYSTEM_THINKING:
                    self._add_system_thinking_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.EXPERIMENT_INQUIRY:
                    self._add_experiment_inquiry_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 文科教学页面类型（feat-032 历史/政治/地理）
                elif page_type == PPTPageType.TIMELINE:
                    self._add_timeline_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.CAUSE_ANALYSIS:
                    self._add_cause_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.HISTORICAL_MATERIAL:
                    self._add_historical_material_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.IMPACT_EVALUATION:
                    self._add_impact_evaluation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.CONCEPT_DEFINITION:
                    self._add_concept_definition_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.PRINCIPLE_EXPLANATION:
                    self._add_principle_explanation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.CASE_ANALYSIS:
                    self._add_case_analysis_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.LOCATION:
                    self._add_location_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.NATURAL_ENVIRONMENT:
                    self._add_natural_environment_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.REGIONAL_COMPARISON:
                    self._add_regional_comparison_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.HUMAN_LAND_RELATION:
                    self._add_human_land_relation_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                # 问题链与元认知提示系统（feat-038）
                elif page_type == PPTPageType.QUESTION_CHAIN:
                    self._add_question_chain_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color,
                        secondary_color
                    )
                elif page_type == PPTPageType.METACOGNITIVE_BUBBLE:
                    self._add_metacognitive_bubble_slide(
                        prs,
                        slide_data,
                        content_size,
                        primary_color
                    )
                elif page_type == PPTPageType.REFLECTION_CHECKPOINT:
                    self._add_reflection_checkpoint_slide(
                        prs,
                        slide_data,
                        content_size,
                        secondary_color
                    )

            # 保存文件
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            logger.info(f"PPT 文件已生成: {output_path}")

            # 添加动画效果（如果动画模块可用）
            if ANIMATION_AVAILABLE:
                try:
                    from .pptx_animator import add_animation_to_ppt
                    logger.info(f"正在为 PPT 添加动画效果（年级：{grade}）...")
                    add_animation_to_ppt(str(output_path), str(output_path), grade, subject)
                    logger.info(f"动画效果已添加：{output_path}")
                except Exception as e:
                    logger.warning(f"动画添加失败：{e}，PPT 文件仍可正常使用")
            # 添加互动触发器效果（如果互动模块可用）
            if INTERACTIVE_AVAILABLE:
                try:
                    interactive_gen = InteractivePPTGenerator()
                    logger.info(f"正在为 PPT 添加点击触发器效果...")
                    interactive_gen.inject_click_trigger(output_path, output_path)
                    logger.info(f"点击触发器已添加：{output_path}")
                except Exception as e:
                    logger.warning(f"触发器添加失败：{e}，PPT 文件仍可正常使用")

            return output_path

        except Exception as e:
            logger.error(f"PPT 文件生成失败: {e}")
            raise RuntimeError(f"PPT 文件生成失败: {e}") from e

    def _add_cover_slide(
        self,
        prs: Presentation,
        title: str,
        title_size: int,
        color: RGBColor,
        use_pinyin: bool = False
    ):
        """添加封面页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 手动添加标题文本框（因为空白布局没有占位符）
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        text_frame = title_box.text_frame
        text_frame.text = title

        # 使用辅助方法设置文本框属性
        self._setup_text_frame(text_frame, vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

        # 设置标题样式
        title_para = text_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(title_size)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "AI 教学 PPT 生成器"

        # 设置副标题文本框属性
        self._setup_text_frame(subtitle_frame, vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(title_size - 10)
        subtitle_para.font.color.rgb = RGBColor(128, 128, 128)

    def _add_toc_slide(
        self,
        prs: Presentation,
        slides: List[Dict[str, Any]],
        title_size: int,
        content_size: int,
        color: RGBColor
    ):
        """添加目录页"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 标题和内容布局

        # 标题
        title_box = slide.shapes.title
        title_box.text = "目录"
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(title_size)
        title_para.font.color.rgb = color

        # 目录内容 - 手动添加文本框（使用辅助方法设置防溢出属性）
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        for idx, slide_data in enumerate(slides):
            if slide_data.get("page_type") == PPTPageType.TOC:
                continue

            p = text_frame.add_paragraph()
            p.text = f"{idx + 1}. {slide_data.get('title', '未知标题')}"
            p.font.size = Pt(content_size)
            p.level = 0

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加内容页（优化版：增加图示、表格、流程图支持）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框（使用辅助方法设置防溢出属性）
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)  # 统一设置 word_wrap、auto_fit、margins

        for item in slide_data.get("content", []):
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.level = 0

        # 添加互动环节
        if slide_data.get("interaction"):
            p = text_frame.add_paragraph()
            p.text = f"\n💡 互动：{slide_data['interaction']}"
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = RGBColor(0, 128, 0)
            p.level = 0

        # 添加口诀
        if slide_data.get("mnemonic"):
            p = text_frame.add_paragraph()
            p.text = f"\n📝 记忆口诀：{slide_data['mnemonic']}"
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = RGBColor(128, 0, 128)
            p.level = 0

    def _add_character_learning_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加生字学习页（语文专属：字 + 拼音 + 组词 + 造句）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "生字学习")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容（使用辅助方法设置防溢出属性）
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        char = slide_data.get("character", {})
        if char:
            # 生字（大字显示）
            p = text_frame.add_paragraph()
            p.text = f"\n【{char.get('word', '')}】"
            p.font.size = Pt(font_size + 10)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

            # 拼音
            if char.get("pinyin"):
                p = text_frame.add_paragraph()
                p.text = f"拼音：{char['pinyin']}"
                p.font.size = Pt(font_size)
                p.level = 0

            # 部首和笔画
            if char.get("radical") or char.get("strokes"):
                info = []
                if char.get("radical"):
                    info.append(f"部首：{char['radical']}")
                if char.get("strokes"):
                    info.append(f"笔画：{char['strokes']}")
                p = text_frame.add_paragraph()
                p.text = "  ".join(info)
                p.font.size = Pt(font_size - 2)
                p.level = 0

            # 组词
            if char.get("groups"):
                p = text_frame.add_paragraph()
                p.text = f"\n组词：{', '.join(char['groups'])}"
                p.font.size = Pt(font_size)
                p.level = 0

            # 造句
            if char.get("sentences"):
                p = text_frame.add_paragraph()
                p.text = "\n造句："
                p.font.size = Pt(font_size)
                p.level = 0
                for sentence in char.get("sentences", [])[:2]:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {sentence}"
                    p.font.size = Pt(font_size - 2)
                    p.level = 1

            # 近反义词
            if char.get("synonyms") or char.get("antonyms"):
                if char.get("synonyms"):
                    p = text_frame.add_paragraph()
                    p.text = f"\n近义词：{', '.join(char['synonyms'])}"
                    p.font.size = Pt(font_size - 2)
                    p.level = 0
                if char.get("antonyms"):
                    p = text_frame.add_paragraph()
                    p.text = f"反义词：{', '.join(char['antonyms'])}"
                    p.font.size = Pt(font_size - 2)
                    p.level = 0

    def _add_poem_appreciation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加古诗鉴赏页（语文专属：原文 + 注释 + 赏析）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "古诗鉴赏")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容（使用辅助方法设置防溢出属性）
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        poem = slide_data.get("poem", {})
        if poem:
            # 诗名和作者
            p = text_frame.add_paragraph()
            p.text = f"{poem.get('title', '')}  【{poem.get('dynasty', '')}】{poem.get('author', '')}"
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = secondary_color
            p.alignment = PP_ALIGN.CENTER

            # 诗句（居中显示）
            if poem.get("lines"):
                p = text_frame.add_paragraph()
                p.text = ""
                for line in poem.get("lines", []):
                    p.text += f"{line}\n"
                p.font.size = Pt(font_size + 4)
                p.alignment = PP_ALIGN.CENTER

            # 朗读节奏
            if poem.get("rhythm"):
                p = text_frame.add_paragraph()
                p.text = f"\n朗读节奏：{poem['rhythm']}"
                p.font.size = Pt(font_size - 2)
                p.font.color.rgb = RGBColor(128, 128, 128)
                p.level = 0

            # 赏析
            if poem.get("appreciation"):
                p = text_frame.add_paragraph()
                p.text = f"\n【赏析】{poem['appreciation']}"
                p.font.size = Pt(font_size)
                p.level = 0

            # 主旨
            if poem.get("theme"):
                p = text_frame.add_paragraph()
                p.text = f"\n【主旨】{poem['theme']}"
                p.font.size = Pt(font_size)
                p.level = 0

    def _add_reading_comprehension_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加阅读理解页（语文专属：段落分析 + 主题提炼）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "阅读理解")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        reading = slide_data.get("reading", {})
        if reading:
            # 段落内容
            if reading.get("paragraph"):
                p = text_frame.add_paragraph()
                p.text = f"【段落】{reading['paragraph'][:200]}{'...' if len(reading.get('paragraph', '')) > 200 else ''}"
                p.font.size = Pt(font_size - 2)
                p.level = 0

            # 层次划分
            if reading.get("structure"):
                p = text_frame.add_paragraph()
                p.text = f"\n【层次】{reading['structure']}"
                p.font.size = Pt(font_size)
                p.level = 0

            # 关键句
            if reading.get("key_sentence"):
                p = text_frame.add_paragraph()
                p.text = f"\n【关键句】{reading['key_sentence']}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
                p.level = 0

            # 修辞手法
            if reading.get("rhetoric"):
                p = text_frame.add_paragraph()
                p.text = f"\n【修辞】{reading['rhetoric']}"
                p.font.size = Pt(font_size - 2)
                p.level = 0

            # 主题思想
            if reading.get("theme"):
                p = text_frame.add_paragraph()
                p.text = f"\n【主题】{reading['theme']}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(128, 0, 128)
                p.level = 0

    def _add_writing_guidance_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加作文指导页（语文专属：审题 + 结构 + 素材）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "作文指导")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        writing = slide_data.get("writing", {})
        if writing:
            # 作文题目
            if writing.get("topic"):
                p = text_frame.add_paragraph()
                p.text = f"【题目】{writing['topic']}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.level = 0

            # 审题立意
            if writing.get("keywords"):
                p = text_frame.add_paragraph()
                p.text = f"\n【关键词】{', '.join(writing['keywords'])}"
                p.font.size = Pt(font_size)
                p.level = 0

            # 结构框架
            if writing.get("structure"):
                p = text_frame.add_paragraph()
                p.text = "\n【结构框架】"
                p.font.size = Pt(font_size)
                p.level = 0
                for i, part in enumerate(writing.get("structure", [])):
                    p = text_frame.add_paragraph()
                    p.text = f"  {i + 1}. {part}"
                    p.font.size = Pt(font_size - 2)
                    p.level = 1

            # 好词好句
            if writing.get("good_words") or writing.get("good_sentences"):
                if writing.get("good_words"):
                    p = text_frame.add_paragraph()
                    p.text = f"\n【好词】{', '.join(writing['good_words'][:6])}"
                    p.font.size = Pt(font_size - 2)
                    p.level = 0
                if writing.get("good_sentences"):
                    p = text_frame.add_paragraph()
                    p.text = f"\n【好句】{writing['good_sentences'][0] if writing['good_sentences'] else ''}"
                    p.font.size = Pt(font_size - 2)
                    p.level = 0

            # 写作提示
            if writing.get("tips"):
                p = text_frame.add_paragraph()
                p.text = f"\n【提示】{writing['tips']}"
                p.font.size = Pt(font_size - 2)
                p.font.color.rgb = RGBColor(255, 100, 100)
                p.level = 0

    def _add_paragraph_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加段落分析页（语文专属：层次划分 + 关键句）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "段落分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        reading = slide_data.get("reading", {})
        if reading:
            # 段落内容
            if reading.get("paragraph"):
                p = text_frame.add_paragraph()
                p.text = f"{reading['paragraph'][:150]}{'...' if len(reading.get('paragraph', '')) > 150 else ''}"
                p.font.size = Pt(font_size - 2)
                p.level = 0

            # 层次划分
            if reading.get("structure"):
                p = text_frame.add_paragraph()
                p.text = f"\n【层次划分】{reading['structure']}"
                p.font.size = Pt(font_size)
                p.level = 0

            # 关键句赏析
            if reading.get("key_sentence") and reading.get("analysis"):
                p = text_frame.add_paragraph()
                p.text = f"\n【关键句】{reading['key_sentence']}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
                p.level = 0
                p = text_frame.add_paragraph()
                p.text = f"【赏析】{reading['analysis']}"
                p.font.size = Pt(font_size - 2)
                p.level = 0

    def _add_rhetoric_device_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加修辞手法页（语文专属：比喻/拟人/排比等）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "修辞手法")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        for item in slide_data.get("content", []):
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.level = 0

        # 添加修辞示例
        if slide_data.get("rhetoric_examples"):
            p = text_frame.add_paragraph()
            p.text = "\n【示例】"
            p.font.size = Pt(font_size)
            p.font.color.rgb = secondary_color
            p.level = 0
            for example in slide_data.get("rhetoric_examples", [])[:3]:
                p = text_frame.add_paragraph()
                p.text = f"  • {example}"
                p.font.size = Pt(font_size - 2)
                p.level = 1

    def _add_diagram_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加图示页（优化版：支持流程图、结构图、概念图）"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 标题
        title_text = slide_data.get("title", "知识图示")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.alignment = PP_ALIGN.CENTER

        # 创建图示容器（用文本框模拟）
        diagram_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
        diagram_frame = diagram_box.text_frame
        diagram_frame.word_wrap = True

        # 渲染图示内容
        diagram_type = slide_data.get("diagram_type", "flowchart")
        diagram_data = slide_data.get("diagram_content", [])

        if diagram_type == "flowchart":
            # 流程图样式
            for i, item in enumerate(diagram_data):
                p = diagram_frame.add_paragraph()
                p.text = f"┌{'─' * 50}┐" if i == 0 else f"│ {item}"
                p.font.size = Pt(font_size)
                p.alignment = PP_ALIGN.CENTER
        elif diagram_type == "structure":
            # 结构图样式
            for item in diagram_data:
                p = diagram_frame.add_paragraph()
                p.text = f"  ● {item}"
                p.font.size = Pt(font_size)
        else:
            # 默认列表样式
            for item in diagram_data:
                p = diagram_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(font_size)

    def _add_table_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加表格页（优化版：支持对比表、分类表、数据表）"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 标题
        title_text = slide_data.get("title", "数据表格")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.alignment = PP_ALIGN.CENTER

        # 获取表格数据
        table_data = slide_data.get("table_data", [])
        if not table_data:
            return

        rows = len(table_data) + 1  # +1 for header
        cols = len(table_data[0]) if table_data else 4

        # 创建表格
        table_x, table_y = Inches(0.8), Inches(1.5)
        table_width, table_height = Inches(8.4), Inches(5)
        table = slide.shapes.add_table(rows, cols, table_x, table_y, table_width, table_height).table

        # 设置列宽
        for i in range(cols):
            table.columns[i].width = table_width / cols

        # 填充表头
        headers = slide_data.get("table_headers", [f"列{i+1}" for i in range(cols)])
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = color
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER

        # 填充数据行
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(font_size - 2)
                    paragraph.alignment = PP_ALIGN.CENTER
                    # 交替行颜色
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

    def _add_interaction_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加互动问答页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "互动问答")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        for item in slide_data.get("content", []):
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)

    def _add_exercise_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加课堂练习页（支持差异化教学难度标注 feat-040）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "课堂练习")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        for item in slide_data.get("content", []):
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)

        # 处理 classroom_exercise 字段（feat-040 差异化教学）
        exercises = slide_data.get("classroom_exercise", [])
        if exercises:
            y_position = 2.0
            for idx, exercise in enumerate(exercises):
                # 获取难度星级
                difficulty = exercise.get("difficulty", "basic")
                star_display = {"basic": "★", "intermediate": "★★", "advanced": "★★★"}.get(difficulty, "★")

                # 添加题目标题（带难度星级）
                question_title = slide.shapes.add_textbox(Inches(1), Inches(y_position), Inches(9), Inches(0.5))
                q_para = question_title.text_frame.add_paragraph()
                q_para.text = f"【{star_display}】{exercise.get('question', '练习题')}"
                q_para.font.size = Pt(font_size - 1)
                q_para.font.bold = True
                # 根据难度设置颜色
                if difficulty == "basic":
                    q_para.font.color.rgb = COLOR_GREEN_2  # 绿色 - 基础
                elif difficulty == "intermediate":
                    q_para.font.color.rgb = _COLOR_ORANGE  # 橙色 - 提高
                else:
                    q_para.font.color.rgb = _COLOR_RED  # 红色 - 拓展

                y_position += 0.6

                # 添加答案和解析（可折叠样式）
                if exercise.get("answer"):
                    answer_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_position), Inches(8.5), Inches(0.4))
                    a_para = answer_box.text_frame.add_paragraph()
                    a_para.text = f"答案：{exercise.get('answer')}"
                    a_para.font.size = Pt(font_size - 2)
                    a_para.font.color.rgb = RGBColor(100, 100, 100)  # 灰色
                    y_position += 0.5

    def _add_summary_slide(
        self,
        prs: Presentation,
        content: Dict[str, Any],
        title_size: int,
        font_size: int,
        color: RGBColor
    ):
        """添加总结回顾页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = "总结回顾"
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(title_size)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        p = text_frame.add_paragraph()
        p.text = content.get("summary", "")
        p.font.size = Pt(font_size)
        p.font.bold = True

        p = text_frame.add_paragraph()
        p.text = "\n重点："
        p.font.size = Pt(font_size)
        p.font.bold = True

        for point in content.get("key_points", []):
            p = text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(font_size)

    def _add_vocabulary_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加单词学习页（英语学科专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "单词学习")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        # 显示单词列表
        vocabulary = slide_data.get("vocabulary", [])
        if vocabulary:
            for vocab_item in vocabulary:
                word = vocab_item.get("word", "")
                phonetic = vocab_item.get("phonetic", "")
                meaning = vocab_item.get("meaning", "")
                example = vocab_item.get("example", "")

                # 单词（大号字体）
                p = text_frame.add_paragraph()
                p.text = word
                p.font.size = Pt(font_size + 4)
                p.font.color.rgb = color
                p.font.bold = True

                # 音标
                if phonetic:
                    p = text_frame.add_paragraph()
                    p.text = f"  [{phonetic}]"
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = RGBColor(128, 128, 128)

                # 释义
                if meaning:
                    p = text_frame.add_paragraph()
                    p.text = f"  {meaning}"
                    p.font.size = Pt(font_size)

                # 例句
                if example:
                    p = text_frame.add_paragraph()
                    p.text = f"  例：{example}"
                    p.font.size = Pt(font_size - 2)
                    p.font.color.rgb = RGBColor(0, 100, 200)

                # 空行分隔
                p = text_frame.add_paragraph()
                p.text = ""
        else:
            # 没有词汇数据时，显示普通内容
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_grammar_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加语法讲解页（英语学科专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "语法讲解")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        # 语法规则说明
        grammar = slide_data.get("grammar", "")
        if grammar:
            p = text_frame.add_paragraph()
            p.text = "语法规则："
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = color
            p.font.bold = True

            p = text_frame.add_paragraph()
            p.text = grammar
            p.font.size = Pt(font_size)

            p = text_frame.add_paragraph()
            p.text = ""
        else:
            # 显示普通内容
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

        # 例句
        if slide_data.get("examples"):
            p = text_frame.add_paragraph()
            p.text = "\n例句："
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = RGBColor(0, 100, 200)
            p.font.bold = True

            for example in slide_data["examples"]:
                p = text_frame.add_paragraph()
                p.text = f"• {example}"
                p.font.size = Pt(font_size)

    def _add_dialogue_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加情景对话页（英语学科专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "情景对话")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        # 对话内容
        dialogue = slide_data.get("dialogue", [])
        if dialogue:
            for line in dialogue:
                if isinstance(line, dict):
                    speaker = line.get("speaker", "")
                    text = line.get("text", "")
                    p = text_frame.add_paragraph()
                    p.text = f"{speaker}: {text}" if speaker else text
                    p.font.size = Pt(font_size)
                else:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.font.size = Pt(font_size)
        else:
            # 显示普通内容
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加课文分析页（英语学科专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "课文分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容 - 手动添加文本框
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        # 段落大意
        if slide_data.get("paragraph_summary"):
            p = text_frame.add_paragraph()
            p.text = "段落大意："
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = color
            p.font.bold = True

            p = text_frame.add_paragraph()
            p.text = slide_data["paragraph_summary"]
            p.font.size = Pt(font_size)

            p = text_frame.add_paragraph()
            p.text = ""

        # 关键句型
        if slide_data.get("key_sentences"):
            p = text_frame.add_paragraph()
            p.text = "关键句型："
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = RGBColor(0, 100, 200)
            p.font.bold = True

            for sentence in slide_data["key_sentences"]:
                p = text_frame.add_paragraph()
                p.text = f"• {sentence}"
                p.font.size = Pt(font_size)

            p = text_frame.add_paragraph()
            p.text = ""

        # 常用表达
        if slide_data.get("expressions"):
            p = text_frame.add_paragraph()
            p.text = "常用表达："
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = RGBColor(128, 0, 128)
            p.font.bold = True

            for expr in slide_data["expressions"]:
                p = text_frame.add_paragraph()
                p.text = f"• {expr}"
                p.font.size = Pt(font_size)
        else:
            # 显示普通内容
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_bridge_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加概念桥接页（先行组织者策略）：连接新旧知识的双色对比布局"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 标题
        title_text = slide_data.get('title', '概念桥接')
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.alignment = PP_ALIGN.CENTER

        # 获取桥接配置
        bridge_config = slide_data.get('bridge_config', {})
        known_concept = bridge_config.get('known_concept', '已知知识')
        new_concept = bridge_config.get('new_concept', '新知识')
        bridge_type = bridge_config.get('bridge_type', '类比')
        bridge_content = bridge_config.get('bridge_content', '')
        connection = bridge_config.get('connection', '')

        # 三色区域布局：左侧已知（主色），中间箭头，右侧新知（辅助色）
        left_x, left_y = Inches(0.3), Inches(1.3)
        right_x, right_y = Inches(5.7), Inches(1.3)
        box_width, box_height = Inches(4), Inches(2.5)

        # 左侧：已知概念（主色背景感）
        left_box = slide.shapes.add_textbox(left_x, left_y, box_width, Inches(0.6))
        left_title = left_box.text_frame
        left_title.text = f'📚 你已经知道'
        left_title.paragraphs[0].font.size = Pt(font_size + 1)
        left_title.paragraphs[0].font.color.rgb = color
        left_title.paragraphs[0].font.bold = True
        left_title.paragraphs[0].alignment = PP_ALIGN.CENTER

        left_content = slide.shapes.add_textbox(left_x, left_y + Inches(0.6), box_width, box_height)
        left_frame = left_content.text_frame
        left_frame.text = known_concept
        for p in left_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.alignment = PP_ALIGN.CENTER

        # 中间：箭头和关系说明
        arrow_x, arrow_y = Inches(4.3), Inches(2.3)
        arrow_box = slide.shapes.add_textbox(arrow_x, arrow_y, Inches(1.4), Inches(1.5))
        arrow_frame = arrow_box.text_frame
        arrow_frame.text = '→'
        arrow_para = arrow_frame.paragraphs[0]
        arrow_para.font.size = Pt(48)
        arrow_para.font.color.rgb = RGBColor(128, 128, 128)
        arrow_para.alignment = PP_ALIGN.CENTER

        # 关系说明在箭头下方
        if connection:
            conn_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(1))
            conn_frame = conn_box.text_frame
            conn_frame.text = f'💡 关系：{connection}'
            conn_para = conn_frame.paragraphs[0]
            conn_para.font.size = Pt(font_size - 1)
            conn_para.font.color.rgb = RGBColor(100, 100, 100)
            conn_para.font.italic = True
            conn_para.alignment = PP_ALIGN.CENTER

        # 右侧：新概念（辅助色）
        right_box = slide.shapes.add_textbox(right_x, right_y, box_width, Inches(0.6))
        right_title = right_box.text_frame
        right_title.text = f'🎯 今天学习'
        right_title.paragraphs[0].font.size = Pt(font_size + 1)
        right_title.paragraphs[0].font.color.rgb = secondary_color
        right_title.paragraphs[0].font.bold = True
        right_title.paragraphs[0].alignment = PP_ALIGN.CENTER

        right_content = slide.shapes.add_textbox(right_x, right_y + Inches(0.6), box_width, box_height)
        right_frame = right_content.text_frame
        right_frame.text = new_concept
        for p in right_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.alignment = PP_ALIGN.CENTER

        # 桥接说明（底部）
        if bridge_content:
            bridge_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1.5))
            bridge_frame = bridge_box.text_frame
            bridge_frame.word_wrap = True
            bridge_frame.text = f'{bridge_type}桥接：{bridge_content}'
            bridge_para = bridge_frame.paragraphs[0]
            bridge_para.font.size = Pt(font_size - 1)
            bridge_para.font.color.rgb = RGBColor(80, 80, 80)

    def _add_comparison_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加对比分析页（优化版：支持两种或多种概念对比）"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 标题
        title_text = slide_data.get("title", "对比分析")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.alignment = PP_ALIGN.CENTER

        # 获取对比数据
        comparison_data = slide_data.get("comparison_items", [])
        if len(comparison_data) >= 2:
            # 左右分栏布局
            left_x, left_y = Inches(0.5), Inches(1.5)
            right_x, right_y = Inches(5.5), Inches(1.5)
            col_width, col_height = Inches(4), Inches(5)

            # 左侧内容
            left_item = comparison_data[0]
            left_box = slide.shapes.add_textbox(left_x, left_y, col_width, Inches(0.8))
            left_title = left_box.text_frame
            left_title.text = left_item.get("title", "概念 A")
            left_title.paragraphs[0].font.size = Pt(font_size + 2)
            left_title.paragraphs[0].font.color.rgb = color
            left_title.paragraphs[0].font.bold = True
            left_title.paragraphs[0].alignment = PP_ALIGN.CENTER

            left_content = slide.shapes.add_textbox(left_x, left_y + Inches(0.8), col_width, col_height)
            left_frame = left_content.text_frame
            for point in left_item.get("points", []):
                p = left_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(font_size)

            # 右侧内容
            right_item = comparison_data[1] if len(comparison_data) > 1 else comparison_data[0]
            right_box = slide.shapes.add_textbox(right_x, right_y, col_width, Inches(0.8))
            right_title = right_box.text_frame
            right_title.text = right_item.get("title", "概念 B")
            right_title.paragraphs[0].font.size = Pt(font_size + 2)
            right_title.paragraphs[0].font.color.rgb = secondary_color
            right_title.paragraphs[0].font.bold = True
            right_title.paragraphs[0].alignment = PP_ALIGN.CENTER

            right_content = slide.shapes.add_textbox(right_x, right_y + Inches(0.8), col_width, col_height)
            right_frame = right_content.text_frame
            for point in right_item.get("points", []):
                p = right_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(font_size)
        else:
            # 没有对比数据时，显示普通内容
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            content_frame = content_box.text_frame
            for item in slide_data.get("content", []):
                p = content_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)


    def _add_mistake_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        primary_color: RGBColor,
        secondary_color: RGBColor
    ):
        """
        添加错题分析页 (feat-041 智能错题分析)

        左右分栏布局：左侧红色展示常见错误，右侧绿色展示正确方法
        帮助学生理解错误原因，预防易错点

        Args:
            prs: PPT 演示文稿对象
            slide_data: 幻灯片数据 (包含 common_mistakes 数组)
            font_size: 字号
            primary_color: 主色调 (用于标题)
            secondary_color: 辅助色 (用于正确方法)
        """
        from pptx.enum.shapes import MSO_SHAPE

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 标题
        title_text = slide_data.get("title", "易错警示")
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = primary_color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 获取错题数据
        common_mistakes = slide_data.get("common_mistakes", [])

        if common_mistakes:
            # 左右分栏布局
            left_x, right_x = Inches(0.5), Inches(5.5)
            col_width = Inches(4)

            # === 左侧：常见错误 (红色) ===
            # 左侧标题
            error_title = slide.shapes.add_textbox(left_x, Inches(1.2), col_width, Inches(0.6))
            error_title_frame = error_title.text_frame
            error_title_para = error_title_frame.paragraphs[0]
            error_title_para.text = "❌ 常见错误"
            error_title_para.font.size = Pt(font_size + 2)
            error_title_para.font.color.rgb = _COLOR_RED
            error_title_para.font.bold = True
            error_title_para.alignment = PP_ALIGN.CENTER

            # 左侧内容区域背景 (浅红色)
            error_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left_x, Inches(1.8), col_width, Inches(4.5)
            )
            error_bg.fill.solid()
            error_bg.fill.fore_color.rgb = RGBColor(255, 240, 240)
            error_bg.line.fill.background()

            # 左侧错题内容
            mistake = common_mistakes[0] if common_mistakes else {}
            error_content = slide.shapes.add_textbox(left_x + Inches(0.2), Inches(2), col_width - Inches(0.4), Inches(4.1))
            error_frame = error_content.text_frame
            error_frame.word_wrap = True

            mistake_example = mistake.get("mistake_example", mistake.get("mistake", "无"))
            p = error_frame.add_paragraph()
            p.text = "【错误示例】"
            p.font.size = Pt(font_size - 2)
            p.font.bold = True
            p.font.color.rgb = _COLOR_RED

            p = error_frame.add_paragraph()
            p.text = mistake_example
            p.font.size = Pt(font_size)
            p.space_after = Pt(8)

            reason = mistake.get("reason", mistake.get("reason_analysis", "无"))
            p = error_frame.add_paragraph()
            p.text = "【错误原因】"
            p.font.size = Pt(font_size - 2)
            p.font.bold = True
            p.font.color.rgb = _COLOR_RED

            p = error_frame.add_paragraph()
            p.text = reason
            p.font.size = Pt(font_size)

            # === 右侧：正确方法 (绿色) ===
            correct_title = slide.shapes.add_textbox(right_x, Inches(1.2), col_width, Inches(0.6))
            correct_title_frame = correct_title.text_frame
            correct_title_para = correct_title_frame.paragraphs[0]
            correct_title_para.text = "✓ 正确方法"
            correct_title_para.font.size = Pt(font_size + 2)
            correct_title_para.font.color.rgb = COLOR_GREEN_1
            correct_title_para.font.bold = True
            correct_title_para.alignment = PP_ALIGN.CENTER

            correct_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                right_x, Inches(1.8), col_width, Inches(4.5)
            )
            correct_bg.fill.solid()
            correct_bg.fill.fore_color.rgb = RGBColor(240, 255, 240)
            correct_bg.line.fill.background()

            correct_content = slide.shapes.add_textbox(right_x + Inches(0.2), Inches(2), col_width - Inches(0.4), Inches(4.1))
            correct_frame = correct_content.text_frame
            correct_frame.word_wrap = True

            correct_method = mistake.get("correct_method", mistake.get("correct", "无"))
            p = correct_frame.add_paragraph()
            p.text = "【正确解法】"
            p.font.size = Pt(font_size - 2)
            p.font.bold = True
            p.font.color.rgb = COLOR_GREEN_1

            p = correct_frame.add_paragraph()
            p.text = correct_method
            p.font.size = Pt(font_size)
            p.space_after = Pt(8)

            prevention = mistake.get("prevention_strategy", mistake.get("prevention", "无"))
            if prevention and prevention != "无":
                p = correct_frame.add_paragraph()
                p.text = "【预防策略】"
                p.font.size = Pt(font_size - 2)
                p.font.bold = True
                p.font.color.rgb = COLOR_GREEN_1

                p = correct_frame.add_paragraph()
                p.text = prevention
                p.font.size = Pt(font_size)
        else:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            content_frame = content_box.text_frame
            for item in slide_data.get("content", []):
                p = content_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(font_size)
                p.space_after = Pt(6)

    def _add_quiz_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加互动选择题页（带 ABCD 选项按钮）"""
        from pptx.enum.shapes import MSO_SHAPE

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 1. 添加标题
        title_text = slide_data.get("title", "互动问答")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加问题文本
        question = slide_data.get("question", "请选择正确答案")
        question_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1.2))
        question_frame = question_box.text_frame
        question_frame.word_wrap = True
        question_frame.text = question
        for para in question_frame.paragraphs:
            para.font.size = Pt(font_size)
            para.alignment = PP_ALIGN.CENTER

        # 3. 添加选项按钮
        options = slide_data.get("options", ["A. 选项一", "B. 选项二", "C. 选项三", "D. 选项四"])
        correct_index = slide_data.get("correct_index", 0)
        option_labels = ['A', 'B', 'C', 'D']

        button_width = Inches(3.5)
        button_height = Inches(0.8)
        button_spacing = Inches(0.3)
        start_y = Inches(2.6)

        for i, option in enumerate(options[:4]):
            col = i % 2
            row = i // 2
            left = Inches(1) + col * (button_width + button_spacing)
            top = start_y + row * (button_height + button_spacing)

            # 创建按钮形状
            button = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, button_width, button_height
            )
            button.fill.solid()
            button.fill.fore_color.rgb = secondary_color
            button.line.fill.background()

            # 添加选项文本
            tf = button.text_frame
            label = option_labels[i] if i < 4 else chr(ord('A') + i)
            tf.text = f"{label}. {option}"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].font.bold = True

        # 4. 添加答案提示
        answer_text = f"正确答案：{option_labels[correct_index]}" if correct_index < 4 else "查看答案"
        answer_box = slide.shapes.add_textbox(Inches(2.5), Inches(5.5), Inches(5), Inches(0.8))
        answer_frame = answer_box.text_frame
        answer_frame.text = answer_text
        answer_para = answer_frame.paragraphs[0]
        answer_para.font.size = Pt(font_size)
        answer_para.font.color.rgb = RGBColor(76, 175, 80)  # 绿色
        answer_para.font.bold = True
        answer_para.alignment = PP_ALIGN.CENTER

    def _add_click_to_reveal_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加点击显示答案页"""
        from pptx.enum.shapes import MSO_SHAPE

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. 添加标题
        title_text = slide_data.get("title", "点击显示答案")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加问题文本
        question = slide_data.get("question", "问题是什么？")
        question_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(2))
        question_frame = question_box.text_frame
        question_frame.word_wrap = True
        question_frame.text = question
        for para in question_frame.paragraphs:
            para.font.size = Pt(font_size)
            para.alignment = PP_ALIGN.CENTER

        # 3. 添加"点击显示"按钮区域
        reveal_btn = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3), Inches(3.5), Inches(4), Inches(1)
        )
        reveal_btn.fill.solid()
        reveal_btn.fill.fore_color.rgb = RGBColor(156, 39, 176)  # 紫色
        reveal_btn.line.fill.background()

        btn_tf = reveal_btn.text_frame
        btn_tf.text = "点击查看答案"
        btn_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        btn_tf.paragraphs[0].font.size = Pt(font_size)
        btn_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        btn_tf.paragraphs[0].font.bold = True

        # 4. 添加答案框
        hidden_content = slide_data.get("answer", "这是隐藏的答案内容")
        answer_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(4.8), Inches(8), Inches(2)
        )
        answer_box.fill.solid()
        answer_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        answer_box.line.color.rgb = RGBColor(76, 175, 80)
        answer_box.line.width = Pt(2)

        ans_tf = answer_box.text_frame
        ans_tf.word_wrap = True
        ans_tf.text = hidden_content
        for para in ans_tf.paragraphs:
            para.font.size = Pt(font_size)
            para.alignment = PP_ALIGN.CENTER

    def _add_matching_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加拖拽匹配游戏页"""
        from pptx.enum.shapes import MSO_SHAPE

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. 添加标题
        title_text = slide_data.get("title", "拖拽匹配游戏")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加左右列标题
        left_title = slide_data.get("left_title", "单词")
        right_title = slide_data.get("right_title", "释义")

        left_title_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(3.5), Inches(0.5))
        lt_frame = left_title_box.text_frame
        lt_frame.text = left_title
        lt_frame.paragraphs[0].font.size = Pt(font_size + 2)
        lt_frame.paragraphs[0].font.color.rgb = color
        lt_frame.paragraphs[0].font.bold = True
        lt_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        right_title_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(3.5), Inches(0.5))
        rt_frame = right_title_box.text_frame
        rt_frame.text = right_title
        rt_frame.paragraphs[0].font.size = Pt(font_size + 2)
        rt_frame.paragraphs[0].font.color.rgb = secondary_color
        rt_frame.paragraphs[0].font.bold = True
        rt_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 3. 添加匹配项
        pairs = slide_data.get("pairs", [])  # [(左，右), ...]
        left_y = Inches(2)
        item_height = Inches(0.7)
        spacing = Inches(0.2)

        for i, pair in enumerate(pairs[:6]):
            if isinstance(pair, (list, tuple)) and len(pair) >= 2:
                left_item, right_item = pair[0], pair[1]
            else:
                continue

            # 左侧项目
            left_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), left_y + i * (item_height + spacing),
                Inches(3.5), item_height
            )
            left_box.fill.solid()
            left_box.fill.fore_color.rgb = color
            left_box.line.fill.background()

            left_tf = left_box.text_frame
            left_tf.text = str(left_item)
            left_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            left_tf.paragraphs[0].font.size = Pt(font_size)
            left_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            # 右侧项目
            right_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.5), left_y + i * (item_height + spacing),
                Inches(3.5), item_height
            )
            right_box.fill.solid()
            right_box.fill.fore_color.rgb = RGBColor(200, 200, 200)
            right_box.line.fill.background()

            right_tf = right_box.text_frame
            right_tf.text = str(right_item)
            right_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            right_tf.paragraphs[0].font.size = Pt(font_size)
            right_tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)

    def _add_quiz_game_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加随堂测验页（多题展示）"""
        from pptx.enum.shapes import MSO_SHAPE

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        option_labels = ['A', 'B', 'C', 'D']

        # 1. 添加标题
        title_text = slide_data.get("title", "随堂测验")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加说明
        hint_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.5))
        hint_frame = hint_box.text_frame
        hint_frame.text = "请选择正确答案"
        hint_frame.paragraphs[0].font.size = Pt(font_size - 2)
        hint_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        hint_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 3. 添加问题
        questions = slide_data.get("questions", [])
        q_start_y = Inches(1.6)

        for i, q in enumerate(questions[:3]):  # 最多 3 题
            q_y = q_start_y + i * Inches(2.8)

            # 问题编号和文本
            q_text = q.get("question", f"问题 {i + 1}")
            q_box = slide.shapes.add_textbox(Inches(0.8), q_y, Inches(8.4), Inches(0.6))
            q_frame = q_box.text_frame
            q_frame.word_wrap = True
            q_frame.text = f"{i + 1}. {q_text}"
            q_frame.paragraphs[0].font.size = Pt(font_size)
            q_frame.paragraphs[0].font.bold = True

            # 选项
            options = q.get("options", [])
            for j, opt in enumerate(options[:4]):
                opt_x = Inches(1) + (j % 2) * Inches(4.2)
                opt_y = q_y + Inches(0.6) + (j // 2) * Inches(0.8)

                opt_btn = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    opt_x, opt_y, Inches(4), Inches(0.7)
                )
                opt_btn.fill.solid()
                opt_btn.fill.fore_color.rgb = color
                opt_btn.line.fill.background()

                opt_tf = opt_btn.text_frame
                label = option_labels[j] if j < 4 else chr(ord('A') + j)
                opt_tf.text = f"{label}. {opt}"
                opt_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                opt_tf.paragraphs[0].font.size = Pt(font_size - 2)
                opt_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    def _add_fill_blank_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加填空题页"""
        from pptx.enum.shapes import MSO_SHAPE

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. 添加标题
        title_text = slide_data.get("title", "填空题")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加句子（用___表示空白）
        sentences = slide_data.get("sentences", [])
        y_pos = Inches(1.5)

        for i, sentence in enumerate(sentences[:6]):
            sent_box = slide.shapes.add_textbox(Inches(1), y_pos + i * Inches(1), Inches(8), Inches(0.8))
            sent_frame = sent_box.text_frame
            sent_frame.word_wrap = True
            sent_frame.text = sentence
            sent_frame.paragraphs[0].font.size = Pt(font_size)

        # 3. 添加答案框
        answers = slide_data.get("answers", [])
        ans_text = "答案：" + " | ".join(str(a) for a in answers) if answers else "答案：见教师用书"

        ans_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(6), Inches(8), Inches(1.2)
        )
        ans_box.fill.solid()
        ans_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        ans_box.line.color.rgb = secondary_color
        ans_box.line.width = Pt(2)

        ans_tf = ans_box.text_frame
        ans_tf.text = ans_text
        ans_tf.paragraphs[0].font.size = Pt(font_size)
        ans_tf.paragraphs[0].font.color.rgb = secondary_color
        ans_tf.paragraphs[0].font.bold = True
        ans_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 理科实验页面渲染方法（feat-031 物理/化学/生物）==========

    def _add_experiment_design_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加实验设计页（物理/化学/生物通用）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "实验设计")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        exp_design = slide_data.get("experiment_design", {})
        if exp_design:
            # 探究目的
            purpose = exp_design.get("purpose", "")
            if purpose:
                p = text_frame.add_paragraph()
                p.text = "【探究目的】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = purpose
                p.font.size = Pt(font_size)

            # 实验器材
            materials = exp_design.get("materials", [])
            if materials:
                p = text_frame.add_paragraph()
                p.text = "\n【实验器材】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for mat in materials:
                    p = text_frame.add_paragraph()
                    p.text = f"• {mat}"
                    p.font.size = Pt(font_size)

            # 实验步骤
            steps = exp_design.get("steps", [])
            if steps:
                p = text_frame.add_paragraph()
                p.text = "\n【实验步骤】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for step in steps:
                    if isinstance(step, dict):
                        step_desc = step.get("description", "")
                        p = text_frame.add_paragraph()
                        p.text = f"{step.get('step_number', '')}. {step_desc}"
                        p.font.size = Pt(font_size)
                    else:
                        p = text_frame.add_paragraph()
                        p.text = str(step)
                        p.font.size = Pt(font_size)

            # 注意事项
            notes = exp_design.get("notes", [])
            if notes:
                p = text_frame.add_paragraph()
                p.text = "\n【注意事项】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                for note in notes:
                    p = text_frame.add_paragraph()
                    p.text = f"⚠ {note}"
                    p.font.size = Pt(font_size)
        else:
            # 没有数据时显示普通内容
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_experiment_observation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加现象观察页（物理/化学通用）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "现象观察")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容 - 三段式对比布局
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        obs = slide_data.get("experiment_observation", {})
        if obs:
            # 实验前
            before = obs.get("before", "")
            if before:
                p = text_frame.add_paragraph()
                p.text = "【实验前】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_BLUE_3
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = before
                p.font.size = Pt(font_size)

            # 实验过程现象
            process = obs.get("process", [])
            if process:
                p = text_frame.add_paragraph()
                p.text = "\n【实验现象】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for phen in process:
                    p = text_frame.add_paragraph()
                    p.text = f"→ {phen}"
                    p.font.size = Pt(font_size)

            # 实验后
            after = obs.get("after", "")
            if after:
                p = text_frame.add_paragraph()
                p.text = "\n【实验后】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = after
                p.font.size = Pt(font_size)

            # 关键现象
            key_points = obs.get("key_points", [])
            if key_points:
                p = text_frame.add_paragraph()
                p.text = "\n【关键现象】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                for kp in key_points:
                    p = text_frame.add_paragraph()
                    p.text = f"★ {kp}"
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_data_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加数据分析页（物理/化学通用）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "数据分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        data = slide_data.get("data_analysis", {})
        if data:
            # 数据处理方法
            method = data.get("processing_method", "")
            if method:
                p = text_frame.add_paragraph()
                p.text = "【数据处理方法】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = method
                p.font.size = Pt(font_size)

            # 发现的规律
            pattern = data.get("pattern_discovered", "")
            if pattern:
                p = text_frame.add_paragraph()
                p.text = "\n【发现的规律】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = pattern
                p.font.size = Pt(font_size)

            # 误差分析
            error = data.get("error_analysis", "")
            if error:
                p = text_frame.add_paragraph()
                p.text = "\n【误差分析】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = error
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_law_summary_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加规律归纳页（物理/化学通用）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "规律归纳")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        law = slide_data.get("physics_law", {})
        if law:
            # 规律名称
            name = law.get("name", "")
            if name:
                p = text_frame.add_paragraph()
                p.text = name
                p.font.size = Pt(font_size + 4)
                p.font.color.rgb = color
                p.font.bold = True

            # 文字表述
            statement = law.get("statement", "")
            if statement:
                p = text_frame.add_paragraph()
                p.text = "\n【文字表述】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = statement
                p.font.size = Pt(font_size)

            # 公式表达
            formula = law.get("formula", "")
            if formula:
                p = text_frame.add_paragraph()
                p.text = "\n【公式表达】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = formula
                p.font.size = Pt(font_size + 2)
                p.font.bold = True

            # 适用条件
            conditions = law.get("conditions", "")
            if conditions:
                p = text_frame.add_paragraph()
                p.text = "\n【适用条件】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = conditions
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_application_transfer_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加应用迁移页（物理/化学/生物通用）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "应用迁移")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        app = slide_data.get("application", {})
        if app:
            # 应用场景
            scenario = app.get("scenario", "")
            if scenario:
                p = text_frame.add_paragraph()
                p.text = "【应用场景】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = scenario
                p.font.size = Pt(font_size)

            # 实际问题
            problem = app.get("problem", "")
            if problem:
                p = text_frame.add_paragraph()
                p.text = "\n【实际问题】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = problem
                p.font.size = Pt(font_size)

            # 解决方案
            solution = app.get("solution", "")
            if solution:
                p = text_frame.add_paragraph()
                p.text = "\n【解决方案】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = solution
                p.font.size = Pt(font_size)

            # 拓展思考
            extension = app.get("extension", "")
            if extension:
                p = text_frame.add_paragraph()
                p.text = "\n【拓展思考】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = extension
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_experiment_steps_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加实验步骤页（化学专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "实验步骤")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        exp_steps = slide_data.get("experiment_steps", {})
        if exp_steps:
            # 实验目的
            purpose = exp_steps.get("purpose", "")
            if purpose:
                p = text_frame.add_paragraph()
                p.text = "【实验目的】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = purpose
                p.font.size = Pt(font_size)

            # 仪器和药品
            instruments = exp_steps.get("instruments", [])
            reagents = exp_steps.get("reagents", [])
            if instruments or reagents:
                p = text_frame.add_paragraph()
                p.text = "\n【实验用品】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                if instruments:
                    p = text_frame.add_paragraph()
                    p.text = "仪器：" + ", ".join(instruments)
                    p.font.size = Pt(font_size)
                if reagents:
                    p = text_frame.add_paragraph()
                    p.text = "药品：" + ", ".join(reagents)
                    p.font.size = Pt(font_size)

            # 操作步骤
            procedures = exp_steps.get("procedures", [])
            if procedures:
                p = text_frame.add_paragraph()
                p.text = "\n【操作步骤】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for proc in procedures:
                    if isinstance(proc, dict):
                        desc = proc.get("description", "")
                        p = text_frame.add_paragraph()
                        p.text = f"{proc.get('step_number', '')}. {desc}"
                        p.font.size = Pt(font_size)
                    else:
                        p = text_frame.add_paragraph()
                        p.text = str(proc)
                        p.font.size = Pt(font_size)

            # 安全事项
            safety = exp_steps.get("safety_notes", [])
            if safety:
                p = text_frame.add_paragraph()
                p.text = "\n【安全注意事项】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_RED
                p.font.bold = True
                for s in safety:
                    p = text_frame.add_paragraph()
                    p.text = f"⚠ {s}"
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_reaction_phenomena_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加反应现象页（化学专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "反应现象")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容 - 反应前后对比
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        phen = slide_data.get("reaction_phenomena", {})
        if phen:
            # 反应前
            before = phen.get("before", {})
            if before:
                p = text_frame.add_paragraph()
                p.text = "【反应前】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_BLUE_3
                p.font.bold = True
                if isinstance(before, dict):
                    substances = before.get("substances", [])
                    appearance = before.get("appearance", "")
                    for sub in substances:
                        p = text_frame.add_paragraph()
                        p.text = f"• {sub}"
                        p.font.size = Pt(font_size)
                    if appearance:
                        p = text_frame.add_paragraph()
                        p.text = f"外观：{appearance}"
                        p.font.size = Pt(font_size)
                else:
                    p = text_frame.add_paragraph()
                    p.text = str(before)
                    p.font.size = Pt(font_size)

            # 反应过程现象
            process = phen.get("process", [])
            if process:
                p = text_frame.add_paragraph()
                p.text = "\n【反应现象】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                for p_item in process:
                    p = text_frame.add_paragraph()
                    p.text = f"→ {p_item}"
                    p.font.size = Pt(font_size)

            # 反应后
            after = phen.get("after", {})
            if after:
                p = text_frame.add_paragraph()
                p.text = "\n【反应后】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                if isinstance(after, dict):
                    products = after.get("products", [])
                    appearance = after.get("appearance", "")
                    for prod in products:
                        p = text_frame.add_paragraph()
                        p.text = f"• {prod}"
                        p.font.size = Pt(font_size)
                    if appearance:
                        p = text_frame.add_paragraph()
                        p.text = f"外观：{appearance}"
                        p.font.size = Pt(font_size)
                else:
                    p = text_frame.add_paragraph()
                    p.text = str(after)
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_chemical_equation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加化学方程式页（化学专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "化学方程式")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        eq = slide_data.get("chemical_equation", {})
        if eq:
            # 反应物
            reactants = eq.get("reactants", [])
            if reactants:
                p = text_frame.add_paragraph()
                p.text = "【反应物】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for r in reactants:
                    p = text_frame.add_paragraph()
                    p.text = f"• {r}"
                    p.font.size = Pt(font_size)

            # 生成物
            products = eq.get("products", [])
            if products:
                p = text_frame.add_paragraph()
                p.text = "\n【生成物】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for prod in products:
                    p = text_frame.add_paragraph()
                    p.text = f"• {prod}"
                    p.font.size = Pt(font_size)

            # 配平方程式
            balanced = eq.get("balanced_equation", "")
            if balanced:
                p = text_frame.add_paragraph()
                p.text = "\n【配平的化学方程式】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = balanced
                p.font.size = Pt(font_size + 4)
                p.font.bold = True

            # 反应条件
            condition = eq.get("condition", "")
            if condition:
                p = text_frame.add_paragraph()
                p.text = f"\n【反应条件】{condition}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True

            # 方程式意义
            meaning = eq.get("meaning", "")
            if meaning:
                p = text_frame.add_paragraph()
                p.text = "\n【方程式意义】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = meaning
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_microscopic_explanation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加微观解释页（化学/生物通用）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "微观解释")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        micro = slide_data.get("microscopic_explanation", {})
        if micro:
            # 微粒模型描述
            model = micro.get("model_description", "")
            if model:
                p = text_frame.add_paragraph()
                p.text = "【微粒模型】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = model
                p.font.size = Pt(font_size)

            # 微粒变化过程
            changes = micro.get("particle_changes", "")
            if changes:
                p = text_frame.add_paragraph()
                p.text = "\n【微粒变化】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = changes
                p.font.size = Pt(font_size)

            # 化学键变化
            bond_breaking = micro.get("bond_breaking", "")
            bond_forming = micro.get("bond_forming", "")
            if bond_breaking or bond_forming:
                p = text_frame.add_paragraph()
                p.text = "\n【化学键变化】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                if bond_breaking:
                    p = text_frame.add_paragraph()
                    p.text = f"断裂：{bond_breaking}"
                    p.font.size = Pt(font_size)
                if bond_forming:
                    p = text_frame.add_paragraph()
                    p.text = f"形成：{bond_forming}"
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_periodic_trend_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加元素周期律页（化学专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "元素周期律")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        trend = slide_data.get("periodic_trend", {})
        if trend:
            # 元素名称和符号
            name = trend.get("element_name", "")
            symbol = trend.get("symbol", "")
            if name or symbol:
                p = text_frame.add_paragraph()
                p.text = f"{name} ({symbol})" if name and symbol else (name or symbol)
                p.font.size = Pt(font_size + 4)
                p.font.color.rgb = color
                p.font.bold = True

            # 周期表位置
            position = trend.get("position", {})
            if position:
                p = text_frame.add_paragraph()
                p.text = f"\n【周期表位置】第{position.get('period', '')}周期 第{position.get('group', '')}族"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

            # 原子结构
            atomic = trend.get("atomic_structure", {})
            if atomic:
                p = text_frame.add_paragraph()
                p.text = "\n【原子结构】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                protons = atomic.get("protons", "")
                electrons = atomic.get("electrons", "")
                config = atomic.get("electron_configuration", "")
                if protons:
                    p = text_frame.add_paragraph()
                    p.text = f"质子数：{protons}"
                    p.font.size = Pt(font_size)
                if electrons:
                    p = text_frame.add_paragraph()
                    p.text = f"电子数：{electrons}"
                    p.font.size = Pt(font_size)
                if config:
                    p = text_frame.add_paragraph()
                    p.text = f"电子排布：{config}"
                    p.font.size = Pt(font_size)

            # 性质递变规律
            props_trend = trend.get("properties_trend", "")
            if props_trend:
                p = text_frame.add_paragraph()
                p.text = "\n【性质递变规律】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = props_trend
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_structure_observation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加结构观察页（生物专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "结构观察")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        struct = slide_data.get("structure_observation", {})
        if struct:
            # 观察对象
            obj = struct.get("object", "")
            if obj:
                p = text_frame.add_paragraph()
                p.text = "【观察对象】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = obj
                p.font.size = Pt(font_size)

            # 结构组成部分
            parts = struct.get("parts", [])
            if parts:
                p = text_frame.add_paragraph()
                p.text = "\n【结构组成】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for part in parts:
                    if isinstance(part, dict):
                        p = text_frame.add_paragraph()
                        p.text = f"• {part.get('name', '')}: {part.get('description', '')}"
                        p.font.size = Pt(font_size)
                    else:
                        p = text_frame.add_paragraph()
                        p.text = f"• {part}"
                        p.font.size = Pt(font_size)

            # 观察方法
            method = struct.get("observation_method", "")
            if method:
                p = text_frame.add_paragraph()
                p.text = "\n【观察方法】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = method
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_function_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加功能分析页（生物专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "功能分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        func = slide_data.get("function_analysis", {})
        if func:
            # 结构名称
            struct_name = func.get("structure_name", "")
            if struct_name:
                p = text_frame.add_paragraph()
                p.text = "【结构名称】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = struct_name
                p.font.size = Pt(font_size)

            # 功能描述
            function = func.get("function", "")
            if function:
                p = text_frame.add_paragraph()
                p.text = "\n【功能】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = function
                p.font.size = Pt(font_size)

            # 功能实现机制
            mechanism = func.get("mechanism", "")
            if mechanism:
                p = text_frame.add_paragraph()
                p.text = "\n【功能实现机制】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = mechanism
                p.font.size = Pt(font_size)

            # 结构与功能相适应的证据
            adaptation = func.get("adaptation_evidence", "")
            if adaptation:
                p = text_frame.add_paragraph()
                p.text = "\n【结构与功能的适应性】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = adaptation
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_process_description_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加过程描述页（生物专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "过程描述")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        proc = slide_data.get("process_description", {})
        if proc:
            # 过程名称
            name = proc.get("process_name", "")
            if name:
                p = text_frame.add_paragraph()
                p.text = "【生命过程】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = name
                p.font.size = Pt(font_size)

            # 过程意义
            significance = proc.get("significance", "")
            if significance:
                p = text_frame.add_paragraph()
                p.text = "\n【生物学意义】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = significance
                p.font.size = Pt(font_size)

            # 阶段
            stages = proc.get("stages", [])
            if stages:
                p = text_frame.add_paragraph()
                p.text = "\n【过程阶段】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for stage in stages:
                    if isinstance(stage, dict):
                        stage_name = stage.get("name", "")
                        desc = stage.get("description", "")
                        p = text_frame.add_paragraph()
                        p.text = f"  {stage_name}: {desc}"
                        p.font.size = Pt(font_size)
                    else:
                        p = text_frame.add_paragraph()
                        p.text = str(stage)
                        p.font.size = Pt(font_size)

            # 调控机制
            regulation = proc.get("regulation", "")
            if regulation:
                p = text_frame.add_paragraph()
                p.text = "\n【调控机制】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = regulation
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_system_thinking_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加系统思维页（生物专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "系统思维")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        system = slide_data.get("system_thinking", {})
        if system:
            # 结构层次
            level = system.get("level", "")
            if level:
                p = text_frame.add_paragraph()
                p.text = "【结构层次】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = level
                p.font.size = Pt(font_size)

            # 与其他系统的联系
            connections = system.get("connections", [])
            if connections:
                p = text_frame.add_paragraph()
                p.text = "\n【系统联系】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for conn in connections:
                    if isinstance(conn, dict):
                        conn_sys = conn.get("connected_system", "")
                        rel = conn.get("relationship", "")
                        p = text_frame.add_paragraph()
                        p.text = f"• {conn_sys}: {rel}"
                        p.font.size = Pt(font_size)
                    else:
                        p = text_frame.add_paragraph()
                        p.text = f"• {conn}"
                        p.font.size = Pt(font_size)

            # 整体协调机制
            coordination = system.get("coordination", "")
            if coordination:
                p = text_frame.add_paragraph()
                p.text = "\n【整体协调】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = coordination
                p.font.size = Pt(font_size)

            # 稳态维持
            homeostasis = system.get("homeostasis", "")
            if homeostasis:
                p = text_frame.add_paragraph()
                p.text = "\n【稳态维持】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = homeostasis
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_experiment_inquiry_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加实验探究页（生物/化学通用）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "实验探究")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        inquiry = slide_data.get("experiment_inquiry", {})
        if inquiry:
            # 探究问题
            question = inquiry.get("question", "")
            if question:
                p = text_frame.add_paragraph()
                p.text = "【探究问题】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = question
                p.font.size = Pt(font_size)

            # 假设
            hypothesis = inquiry.get("hypothesis", "")
            if hypothesis:
                p = text_frame.add_paragraph()
                p.text = "\n【假设】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = hypothesis
                p.font.size = Pt(font_size)

            # 实验设计
            design = inquiry.get("design", {})
            if design:
                p = text_frame.add_paragraph()
                p.text = "\n【实验设计】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

                # 变量
                variables = design.get("variables", {})
                if variables:
                    p = text_frame.add_paragraph()
                    p.text = "变量控制："
                    p.font.size = Pt(font_size)
                    p.font.bold = True
                    indep = variables.get("independent", "")
                    dep = variables.get("dependent", "")
                    controlled = variables.get("controlled", [])
                    if indep:
                        p = text_frame.add_paragraph()
                        p.text = f"  自变量：{indep}"
                        p.font.size = Pt(font_size)
                    if dep:
                        p = text_frame.add_paragraph()
                        p.text = f"  因变量：{dep}"
                        p.font.size = Pt(font_size)
                    if controlled:
                        p = text_frame.add_paragraph()
                        p.text = f"  控制变量：{', '.join(controlled)}"
                        p.font.size = Pt(font_size)

                # 组别
                groups = design.get("groups", [])
                if groups:
                    p = text_frame.add_paragraph()
                    p.text = "\n实验组别："
                    p.font.size = Pt(font_size)
                    p.font.bold = True
                    for g in groups:
                        p = text_frame.add_paragraph()
                        p.text = f"  • {g}"
                        p.font.size = Pt(font_size)

            # 预期结果
            expected = inquiry.get("expected_result", "")
            if expected:
                p = text_frame.add_paragraph()
                p.text = "\n【预期结果】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = expected
                p.font.size = Pt(font_size)

            # 结论推导
            conclusion = inquiry.get("conclusion", "")
            if conclusion:
                p = text_frame.add_paragraph()
                p.text = "\n【结论推导】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = conclusion
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_timeline_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加时间轴页（历史专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "时间轴")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        timeline = slide_data.get("timeline", {})
        if timeline:
            # 时间范围
            period = timeline.get("period", "")
            if period:
                p = text_frame.add_paragraph()
                p.text = f"【时间范围】{period}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

            # 关键事件
            events = timeline.get("events", [])
            if events:
                p = text_frame.add_paragraph()
                p.text = "\n【关键事件】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                for event in events:
                    p = text_frame.add_paragraph()
                    year = event.get("year", "")
                    event_name = event.get("event", "")
                    p.text = f"  • {year} - {event_name}"
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_cause_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加因果分析页（历史专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "因果分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        cause = slide_data.get("cause_analysis", {})
        if cause:
            root = cause.get("root_cause", "")
            if root:
                p = text_frame.add_paragraph()
                p.text = "【根本原因】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = root
                p.font.size = Pt(font_size)

            direct = cause.get("direct_cause", "")
            if direct:
                p = text_frame.add_paragraph()
                p.text = "\n【直接原因】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = direct
                p.font.size = Pt(font_size)

            trigger = cause.get("trigger", "")
            if trigger:
                p = text_frame.add_paragraph()
                p.text = "\n【导火线】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = trigger
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_historical_material_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加史料解读页（历史专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "史料解读")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        material = slide_data.get("historical_material", {})
        if material:
            text = material.get("text", "")
            if text:
                p = text_frame.add_paragraph()
                p.text = "【史料原文】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = text
                p.font.size = Pt(font_size)
                p.space_after = Pt(10)

            source = material.get("source", "")
            if source:
                p = text_frame.add_paragraph()
                p.text = f"\n【出处】{source}"
                p.font.size = Pt(font_size)
                p.font.italic = True

            value = material.get("value", "")
            if value:
                p = text_frame.add_paragraph()
                p.text = "\n【史料价值】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = value
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_impact_evaluation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加影响评价页（历史专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "影响评价")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        impact = slide_data.get("impact_evaluation", {})
        if impact:
            short_term = impact.get("short_term", "")
            if short_term:
                p = text_frame.add_paragraph()
                p.text = "【短期影响】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = short_term
                p.font.size = Pt(font_size)

            long_term = impact.get("long_term", "")
            if long_term:
                p = text_frame.add_paragraph()
                p.text = "\n【长期影响】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = long_term
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_concept_definition_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加概念界定页（政治专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "概念界定")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        concept = slide_data.get("concept_definition", {})
        if concept:
            name = concept.get("name", "")
            if name:
                p = text_frame.add_paragraph()
                p.text = f"【概念名称】{name}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

            connotation = concept.get("connotation", "")
            if connotation:
                p = text_frame.add_paragraph()
                p.text = "\n【内涵】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = connotation
                p.font.size = Pt(font_size)

            extension = concept.get("extension", "")
            if extension:
                p = text_frame.add_paragraph()
                p.text = "\n【外延】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = extension
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_principle_explanation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加原理阐述页（政治专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "原理阐述")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        principle = slide_data.get("principle_explanation", {})
        if principle:
            content = principle.get("content", "")
            if content:
                p = text_frame.add_paragraph()
                p.text = "【原理内容】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = content
                p.font.size = Pt(font_size)

            basis = principle.get("theoretical_basis", "")
            if basis:
                p = text_frame.add_paragraph()
                p.text = "\n【理论依据】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = basis
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_case_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加案例分析页（政治专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "案例分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        case = slide_data.get("case_analysis", {})
        if case:
            background = case.get("background", "")
            if background:
                p = text_frame.add_paragraph()
                p.text = "【案例背景】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = background
                p.font.size = Pt(font_size)

            description = case.get("description", "")
            if description:
                p = text_frame.add_paragraph()
                p.text = "\n【案例描述】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = description
                p.font.size = Pt(font_size)

            analysis = case.get("theoretical_analysis", "")
            if analysis:
                p = text_frame.add_paragraph()
                p.text = "\n【理论分析】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = analysis
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_location_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加位置定位页（地理专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "位置定位")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        location = slide_data.get("location", {})
        if location:
            lat_lon = location.get("latitude_longitude", "")
            if lat_lon:
                p = text_frame.add_paragraph()
                p.text = f"【经纬度位置】{lat_lon}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

            sea_land = location.get("sea_land_position", "")
            if sea_land:
                p = text_frame.add_paragraph()
                p.text = f"\n【海陆位置】{sea_land}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = sea_land
                p.font.size = Pt(font_size)

            relative = location.get("relative_position", "")
            if relative:
                p = text_frame.add_paragraph()
                p.text = "\n【相对位置】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = relative
                p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_natural_environment_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加自然环境页（地理专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "自然环境")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        env = slide_data.get("natural_environment", {})
        if env:
            terrain = env.get("terrain", {})
            if isinstance(terrain, dict):
                terrain_type = terrain.get("type", "")
                if terrain_type:
                    p = text_frame.add_paragraph()
                    p.text = f"【地形】{terrain_type}"
                    p.font.size = Pt(font_size + 2)
                    p.font.color.rgb = color
                    p.font.bold = True

            climate = env.get("climate", {})
            if isinstance(climate, dict):
                climate_type = climate.get("type", "")
                if climate_type:
                    p = text_frame.add_paragraph()
                    p.text = f"\n【气候】{climate_type}"
                    p.font.size = Pt(font_size + 2)
                    p.font.color.rgb = secondary_color
                    p.font.bold = True
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_regional_comparison_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加区域比较页（地理专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "区域比较")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        comparison = slide_data.get("regional_comparison", {})
        if comparison:
            objects = comparison.get("comparison_objects", [])
            if objects:
                p = text_frame.add_paragraph()
                p.text = f"【比较对象】{' vs '.join(objects)}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

            similarities = comparison.get("similarities", [])
            if similarities:
                p = text_frame.add_paragraph()
                p.text = "\n【相似点】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                for sim in similarities:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {sim.get('description', '') if isinstance(sim, dict) else sim}"
                    p.font.size = Pt(font_size)

            differences = comparison.get("differences", [])
            if differences:
                p = text_frame.add_paragraph()
                p.text = "\n【差异点】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = _COLOR_ORANGE
                p.font.bold = True
                for diff in differences:
                    if isinstance(diff, dict):
                        p = text_frame.add_paragraph()
                        p.text = f"  • {diff.get('aspect', '')}: {diff.get('region_a', '')} vs {diff.get('region_b', '')}"
                        p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_human_land_relation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加人地关系页（地理专属）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "人地关系")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 内容
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        relation = slide_data.get("human_land_relationship", {})
        if relation:
            env_impact = relation.get("environment_impact_on_human", "")
            if env_impact:
                p = text_frame.add_paragraph()
                p.text = "【环境对人类的影响】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = env_impact
                p.font.size = Pt(font_size)

            human_impact = relation.get("human_impact_on_environment", "")
            if human_impact:
                p = text_frame.add_paragraph()
                p.text = "\n【人类对环境的影响】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = human_impact
                p.font.size = Pt(font_size)

            measures = relation.get("coordination_measures", [])
            if measures:
                p = text_frame.add_paragraph()
                p.text = "\n【协调措施】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = COLOR_GREEN_1
                p.font.bold = True
                for m in measures:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {m}"
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    # ========== 英语学科深度增强页面类型（feat-033）==========

    def _add_word_root_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加词根词缀页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "词根词缀分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        word_analysis = slide_data.get("word_analysis", {})
        if word_analysis:
            # 词根
            root = word_analysis.get("root", "")
            if root:
                p = text_frame.add_paragraph()
                p.text = f"【词根】{root}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

                root_meaning = word_analysis.get("root_meaning", "")
                if root_meaning:
                    p = text_frame.add_paragraph()
                    p.text = f"  含义：{root_meaning}"
                    p.font.size = Pt(font_size)

                root_origin = word_analysis.get("root_origin", "")
                if root_origin:
                    p = text_frame.add_paragraph()
                    p.text = f"  来源：{root_origin}"
                    p.font.size = Pt(font_size)

                p = text_frame.add_paragraph()
                p.text = ""

            # 前缀
            prefix = word_analysis.get("prefix", "")
            if prefix:
                p = text_frame.add_paragraph()
                p.text = f"【前缀】{prefix}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(106, 90, 205)
                p.font.bold = True

                prefix_meaning = word_analysis.get("prefix_meaning", "")
                if prefix_meaning:
                    p = text_frame.add_paragraph()
                    p.text = f"  含义：{prefix_meaning}"
                    p.font.size = Pt(font_size)

                p = text_frame.add_paragraph()
                p.text = ""

            # 后缀
            suffix = word_analysis.get("suffix", "")
            if suffix:
                p = text_frame.add_paragraph()
                p.text = f"【后缀】{suffix}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(63, 81, 181)
                p.font.bold = True

                suffix_function = word_analysis.get("suffix_function", "")
                if suffix_function:
                    p = text_frame.add_paragraph()
                    p.text = f"  功能：{suffix_function}"
                    p.font.size = Pt(font_size)

                p = text_frame.add_paragraph()
                p.text = ""

            # 词族
            word_family = word_analysis.get("word_family", [])
            if word_family:
                p = text_frame.add_paragraph()
                p.text = "【同根词族】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(76, 175, 80)
                p.font.bold = True

                for word in word_family:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {word}"
                    p.font.size = Pt(font_size)

            # 构词法图解
            structure_diagram = word_analysis.get("structure_diagram", "")
            if structure_diagram:
                p = text_frame.add_paragraph()
                p.text = ""
                p = text_frame.add_paragraph()
                p.text = f"【构词图解】{structure_diagram}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(255, 159, 67)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_word_structure_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加构词法图解页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # 标题
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "构词法图解")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        word_analysis = slide_data.get("word_analysis", {})
        target_word = slide_data.get("target_word", "单词")

        if word_analysis:
            p = text_frame.add_paragraph()
            p.text = target_word
            p.font.size = Pt(font_size + 8)
            p.font.color.rgb = color
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            p = text_frame.add_paragraph()
            p.text = ""

            prefix = word_analysis.get("prefix", "")
            root = word_analysis.get("root", "")
            suffix = word_analysis.get("suffix", "")

            decomposition = []
            if prefix:
                decomposition.append(f"前缀 ({prefix})")
            if root:
                decomposition.append(f"词根 ({root})")
            if suffix:
                decomposition.append(f"后缀 ({suffix})")

            if decomposition:
                p = text_frame.add_paragraph()
                p.text = " = ".join(decomposition)
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(128, 128, 128)

            p = text_frame.add_paragraph()
            p.text = ""

            memory_tip = word_analysis.get("memory_tip", "")
            if memory_tip:
                p = text_frame.add_paragraph()
                p.text = f"💡 记忆技巧：{memory_tip}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(255, 159, 67)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_etymology_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加词源故事页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "词源故事")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        word_analysis = slide_data.get("word_analysis", {})
        if word_analysis:
            etymology = word_analysis.get("etymology", "")
            if etymology:
                p = text_frame.add_paragraph()
                p.text = etymology
                p.font.size = Pt(font_size)
                p.line_spacing = 1.5
            else:
                for item in slide_data.get("content", []):
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_collocation_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加固定搭配页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "固定搭配")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        collocations = slide_data.get("collocations", {})
        if collocations:
            phrasal_verbs = collocations.get("phrasal_verbs", [])
            if phrasal_verbs:
                p = text_frame.add_paragraph()
                p.text = "【动词短语】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

                for pv in phrasal_verbs:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {pv}"
                    p.font.size = Pt(font_size)

                p = text_frame.add_paragraph()
                p.text = ""

            prepositional_phrases = collocations.get("prepositional_phrases", [])
            if prepositional_phrases:
                p = text_frame.add_paragraph()
                p.text = "【介词短语】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(106, 90, 205)
                p.font.bold = True

                for pp in prepositional_phrases:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {pp}"
                    p.font.size = Pt(font_size)

                p = text_frame.add_paragraph()
                p.text = ""

            fixed_expressions = collocations.get("fixed_expressions", [])
            if fixed_expressions:
                p = text_frame.add_paragraph()
                p.text = "【固定表达】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(63, 81, 181)
                p.font.bold = True

                for fe in fixed_expressions:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {fe}"
                    p.font.size = Pt(font_size)

                p = text_frame.add_paragraph()
                p.text = ""

            idioms = collocations.get("idioms", [])
            if idioms:
                p = text_frame.add_paragraph()
                p.text = "【惯用语】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(255, 159, 67)
                p.font.bold = True

                for idiom in idioms:
                    p = text_frame.add_paragraph()
                    p.text = f"  • {idiom}"
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_vocab_network_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加词汇网络页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "词汇网络")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        collocations = slide_data.get("collocations", {})

        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(5))
        left_frame = left_box.text_frame

        synonyms = collocations.get("synonyms", [])
        if synonyms:
            p = left_frame.add_paragraph()
            p.text = "【同义词】"
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = RGBColor(67, 160, 71)
            p.font.bold = True

            for syn in synonyms:
                if isinstance(syn, dict):
                    word = syn.get("word", "")
                    usage = syn.get("usage_diff", "")
                    p = left_frame.add_paragraph()
                    p.text = f"  • {word}"
                    p.font.size = Pt(font_size)
                    if usage:
                        p = left_frame.add_paragraph()
                        p.text = f"    （{usage}）"
                        p.font.size = Pt(font_size - 2)
                        p.font.color.rgb = RGBColor(128, 128, 128)
                else:
                    p = left_frame.add_paragraph()
                    p.text = f"  • {syn}"
                    p.font.size = Pt(font_size)

        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(5))
        right_frame = right_box.text_frame

        antonyms = collocations.get("antonyms", [])
        if antonyms:
            p = right_frame.add_paragraph()
            p.text = "【反义词】"
            p.font.size = Pt(font_size + 2)
            p.font.color.rgb = RGBColor(255, 107, 107)
            p.font.bold = True

            for ant in antonyms:
                if isinstance(ant, dict):
                    word = ant.get("word", "")
                    p = right_frame.add_paragraph()
                    p.text = f"  • {word}"
                    p.font.size = Pt(font_size)
                else:
                    p = right_frame.add_paragraph()
                    p.text = f"  • {ant}"
                    p.font.size = Pt(font_size)

    def _add_collocation_matrix_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加搭配矩阵页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "搭配矩阵")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        collocations = slide_data.get("collocations", {})
        matrix = collocations.get("collocation_matrix", {})

        if matrix:
            verb_noun = matrix.get("verb_noun", [])
            if verb_noun:
                rows = len(verb_noun) + 1
                table_shape = slide.shapes.add_table(rows, 2, Inches(1), Inches(2.5), Inches(8), Inches(0.8 * rows))
                table = table_shape.table

                table.cell(0, 0).text = "动词"
                table.cell(0, 1).text = "名词搭配"
                for cell in [table.cell(0, 0), table.cell(0, 1)]:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = color
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(font_size)
                        para.font.color.rgb = RGBColor(255, 255, 255)
                        para.font.bold = True

                for i, (v, n) in enumerate(verb_noun):
                    table.cell(i + 1, 0).text = str(v)
                    table.cell(i + 1, 1).text = str(n)
                    for cell in [table.cell(i + 1, 0), table.cell(i + 1, 1)]:
                        for para in cell.text_frame.paragraphs:
                            para.font.size = Pt(font_size)
                            para.alignment = PP_ALIGN.CENTER
        else:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
            text_frame = content_box.text_frame
            self._setup_text_frame(text_frame)
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_syntax_tree_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加语法树页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "句子成分分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        syntax_tree = slide_data.get("syntax_tree", {})
        if syntax_tree:
            sentence = syntax_tree.get("sentence", "")
            if sentence:
                p = text_frame.add_paragraph()
                p.text = f"【原句】{sentence}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True
                p = text_frame.add_paragraph()
                p.text = ""

            sentence_type = syntax_tree.get("sentence_type", "")
            if sentence_type:
                p = text_frame.add_paragraph()
                p.text = f"【句子类型】{sentence_type}"
                p.font.size = Pt(font_size)

            main_structure = syntax_tree.get("main_structure", "")
            if main_structure:
                p = text_frame.add_paragraph()
                p.text = f"【主干结构】{main_structure}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(106, 90, 205)
                p = text_frame.add_paragraph()
                p.text = ""

            components = syntax_tree.get("components", [])
            if components:
                p = text_frame.add_paragraph()
                p.text = "【句子成分】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = RGBColor(63, 81, 181)
                p.font.bold = True

                for comp in components:
                    if isinstance(comp, dict):
                        name = comp.get("name", "")
                        comp_type = comp.get("type", "")
                        text_content = comp.get("text", "")
                        function = comp.get("function", "")

                        p = text_frame.add_paragraph()
                        p.text = f"  • {name} ({comp_type})"
                        p.font.size = Pt(font_size)
                        p.font.color.rgb = RGBColor(76, 175, 80)

                        if text_content:
                            p = text_frame.add_paragraph()
                            p.text = f"    内容：{text_content}"
                            p.font.size = Pt(font_size - 2)

                        if function:
                            p = text_frame.add_paragraph()
                            p.text = f"    功能：{function}"
                            p.font.size = Pt(font_size - 2)
                            p.font.color.rgb = RGBColor(128, 128, 128)

                p = text_frame.add_paragraph()
                p.text = ""

            tree_diagram = syntax_tree.get("tree_diagram", "")
            if tree_diagram:
                p = text_frame.add_paragraph()
                p.text = f"【树状图】{tree_diagram}"
                p.font.size = Pt(font_size - 2)
                p.font.color.rgb = RGBColor(128, 128, 128)

            clause_analysis = syntax_tree.get("clause_analysis", {})
            if clause_analysis:
                p = text_frame.add_paragraph()
                p.text = ""
                main_clause = clause_analysis.get("main_clause", "")
                if main_clause:
                    p = text_frame.add_paragraph()
                    p.text = f"【主句】{main_clause}"
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = RGBColor(255, 159, 67)

                subordinate_clauses = clause_analysis.get("subordinate_clauses", [])
                for clause in subordinate_clauses:
                    if isinstance(clause, dict):
                        clause_type = clause.get("type", "")
                        introducer = clause.get("introducer", "")
                        clause_text = clause.get("text", "")

                        p = text_frame.add_paragraph()
                        p.text = f"  └─ {clause_type} ({introducer}): {clause_text}"
                        p.font.size = Pt(font_size - 2)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_syntax_structure_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加句法结构页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "句法结构图解")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        syntax_tree = slide_data.get("syntax_tree", {})
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        if syntax_tree:
            tree_diagram = syntax_tree.get("tree_diagram", "")
            if tree_diagram:
                p = text_frame.add_paragraph()
                p.text = tree_diagram
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(128, 128, 128)
            else:
                for item in slide_data.get("content", []):
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.font.size = Pt(font_size)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_clause_analysis_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加从句分析页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "从句分析")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        syntax_tree = slide_data.get("syntax_tree", {})
        if syntax_tree:
            clause_analysis = syntax_tree.get("clause_analysis", {})
            if clause_analysis:
                main_clause = clause_analysis.get("main_clause", "")
                if main_clause:
                    p = text_frame.add_paragraph()
                    p.text = f"【主句】{main_clause}"
                    p.font.size = Pt(font_size + 2)
                    p.font.color.rgb = color
                    p.font.bold = True
                    p = text_frame.add_paragraph()
                    p.text = ""

                subordinate_clauses = clause_analysis.get("subordinate_clauses", [])
                if subordinate_clauses:
                    p = text_frame.add_paragraph()
                    p.text = "【从句】"
                    p.font.size = Pt(font_size + 2)
                    p.font.color.rgb = RGBColor(106, 90, 205)
                    p.font.bold = True

                    for clause in subordinate_clauses:
                        if isinstance(clause, dict):
                            clause_type = clause.get("type", "")
                            introducer = clause.get("introducer", "")
                            function = clause.get("function", "")
                            clause_text = clause.get("text", "")

                            p = text_frame.add_paragraph()
                            p.text = f"  • {clause_type}"
                            p.font.size = Pt(font_size)
                            p.font.bold = True

                            if introducer:
                                p = text_frame.add_paragraph()
                                p.text = f"    引导词：{introducer}"
                                p.font.size = Pt(font_size - 2)

                            if function:
                                p = text_frame.add_paragraph()
                                p.text = f"    功能：{function}"
                                p.font.size = Pt(font_size - 2)

                            if clause_text:
                                p = text_frame.add_paragraph()
                                p.text = f"    内容：{clause_text}"
                                p.font.size = Pt(font_size - 2)
                                p.font.color.rgb = RGBColor(128, 128, 128)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_special_sentence_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加特殊句型页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "特殊句型")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        special_sentences = slide_data.get("special_sentences", [])
        if special_sentences:
            for sentence_item in special_sentences:
                if isinstance(sentence_item, dict):
                    name = sentence_item.get("name", "")
                    structure = sentence_item.get("structure", "")
                    example = sentence_item.get("example", "")

                    p = text_frame.add_paragraph()
                    p.text = f"【{name}】"
                    p.font.size = Pt(font_size + 2)
                    p.font.color.rgb = color
                    p.font.bold = True

                    if structure:
                        p = text_frame.add_paragraph()
                        p.text = f"  结构：{structure}"
                        p.font.size = Pt(font_size)

                    if example:
                        p = text_frame.add_paragraph()
                        p.text = f"  例句：{example}"
                        p.font.size = Pt(font_size - 2)
                        p.font.color.rgb = RGBColor(0, 100, 200)

                    p = text_frame.add_paragraph()
                    p.text = ""
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_tense_timeline_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加时态时间轴页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "时态时间轴")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        self._setup_text_frame(text_frame)

        tense_timeline = slide_data.get("tense_timeline", {})
        if tense_timeline:
            tense_name = tense_timeline.get("tense_name", "")
            if tense_name:
                p = text_frame.add_paragraph()
                p.text = f"【时态】{tense_name}"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

            tense_group = tense_timeline.get("tense_group", "")
            if tense_group:
                p = text_frame.add_paragraph()
                p.text = f"【时态组】{tense_group}"
                p.font.size = Pt(font_size)

            aspect = tense_timeline.get("aspect", "")
            if aspect:
                p = text_frame.add_paragraph()
                p.text = f"【状态】{aspect}"
                p.font.size = Pt(font_size)

            p = text_frame.add_paragraph()
            p.text = ""

            time_axis = tense_timeline.get("time_axis_description", "")
            if time_axis:
                p = text_frame.add_paragraph()
                p.text = f"【时间轴】{time_axis}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(106, 90, 205)

            reference_point = tense_timeline.get("reference_point", "")
            if reference_point:
                p = text_frame.add_paragraph()
                p.text = f"【参照点】{reference_point}"
                p.font.size = Pt(font_size)

            action_point = tense_timeline.get("action_point", "")
            if action_point:
                p = text_frame.add_paragraph()
                p.text = f"【动作点】{action_point}"
                p.font.size = Pt(font_size)

            duration = tense_timeline.get("duration", "")
            if duration:
                p = text_frame.add_paragraph()
                p.text = f"【持续段】{duration}"
                p.font.size = Pt(font_size)

            p = text_frame.add_paragraph()
            p.text = ""

            usage = tense_timeline.get("usage", "")
            if usage:
                p = text_frame.add_paragraph()
                p.text = f"【用法】{usage}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(76, 175, 80)

            time_markers = tense_timeline.get("time_markers", [])
            if time_markers:
                p = text_frame.add_paragraph()
                p.text = ""
                p = text_frame.add_paragraph()
                p.text = f"【常用时间状语】{', '.join(time_markers)}"
                p.font.size = Pt(font_size - 2)
                p.font.color.rgb = RGBColor(128, 128, 128)

            comparison = tense_timeline.get("comparison", {})
            if comparison:
                similar = comparison.get("similar_tense", "")
                diff = comparison.get("difference", "")

                p = text_frame.add_paragraph()
                p.text = ""
                p = text_frame.add_paragraph()
                p.text = f"【易混时态】{similar}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(255, 159, 67)

                if diff:
                    p = text_frame.add_paragraph()
                    p.text = f"  区别：{diff}"
                    p.font.size = Pt(font_size - 2)
        else:
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_tense_overview_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """添加时态概览页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "时态概览")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        tense_overview = slide_data.get("tense_overview", [])

        if tense_overview:
            rows = min(len(tense_overview) + 1, 17)
            table_shape = slide.shapes.add_table(rows, 3, Inches(0.5), Inches(2), Inches(9), Inches(0.6 * rows))
            table = table_shape.table

            headers = ["时态", "用法", "例句"]
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = color
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(font_size - 2)
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.font.bold = True

            for i, tense_item in enumerate(tense_overview[:16]):
                if isinstance(tense_item, dict):
                    tense_name = tense_item.get("tense", "")
                    usage = tense_item.get("usage", "")
                    example = tense_item.get("example", "")

                    table.cell(i + 1, 0).text = tense_name
                    table.cell(i + 1, 1).text = usage
                    table.cell(i + 1, 2).text = example

                    for j in range(3):
                        for para in table.cell(i + 1, j).text_frame.paragraphs:
                            para.font.size = Pt(font_size - 2)
        else:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
            text_frame = content_box.text_frame
            self._setup_text_frame(text_frame)
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    def _add_tense_comparison_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加时态对比页（英语学科深度增强 - feat-033）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "时态对比")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color

        tense_overview = slide_data.get("tense_overview", [])

        if len(tense_overview) >= 2:
            tense1 = tense_overview[0]
            left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(5))
            left_frame = left_box.text_frame

            if isinstance(tense1, dict):
                p = left_frame.add_paragraph()
                p.text = f"【{tense1.get('tense', '时态 1')}】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = color
                p.font.bold = True

                usage1 = tense1.get("usage", "")
                if usage1:
                    p = left_frame.add_paragraph()
                    p.text = f"用法：{usage1}"
                    p.font.size = Pt(font_size)

                example1 = tense1.get("example", "")
                if example1:
                    p = left_frame.add_paragraph()
                    p.text = f"例句：{example1}"
                    p.font.size = Pt(font_size - 2)
                    p.font.color.rgb = RGBColor(0, 100, 200)

            tense2 = tense_overview[1]
            right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4), Inches(5))
            right_frame = right_box.text_frame

            if isinstance(tense2, dict):
                p = right_frame.add_paragraph()
                p.text = f"【{tense2.get('tense', '时态 2')}】"
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = secondary_color
                p.font.bold = True

                usage2 = tense2.get("usage", "")
                if usage2:
                    p = right_frame.add_paragraph()
                    p.text = f"用法：{usage2}"
                    p.font.size = Pt(font_size)

                example2 = tense2.get("example", "")
                if example2:
                    p = right_frame.add_paragraph()
                    p.text = f"例句：{example2}"
                    p.font.size = Pt(font_size - 2)
                    p.font.color.rgb = RGBColor(0, 100, 200)

            diff_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(8.4), Inches(1))
            diff_frame = diff_box.text_frame
            p = diff_frame.add_paragraph()
            p.text = slide_data.get("difference", "请分析以上两种时态的区别")
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = RGBColor(255, 159, 67)
        else:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
            text_frame = content_box.text_frame
            self._setup_text_frame(text_frame)
            for item in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)

    # ========== 问题链与元认知提示系统（feat-038）==========

    def _add_question_chain_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        primary_color: RGBColor,
        secondary_color: RGBColor
    ):
        """
        添加问题链递进页（feat-038）
        展示递进式问题序列：是什么 → 为什么 → 怎么用 → 如果...会怎样 → 还有其他方法吗
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "问题链思考")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = primary_color
        title_para.font.bold = True

        # 问题链图标映射
        question_icons = {
            "what": "❓",
            "why": "🤔",
            "how": "✏️",
            "what_if": "💡",
            "what_else": "🚀",
        }

        # 层级颜色映射（从低阶到高阶）
        level_colors = {
            "what": primary_color,
            "why": secondary_color,
            "how": primary_color,
            "what_if": secondary_color,
            "what_else": _COLOR_GOLD,
        }

        question_chain = slide_data.get("question_chain", [])
        start_y = 2.0
        box_height = 1.2

        for i, question_item in enumerate(question_chain):
            if isinstance(question_item, dict):
                level = question_item.get("level", "what")
                question_text = question_item.get("question", "")
                icon = question_icons.get(level, "❓")
                level_color = level_colors.get(level, primary_color)

                # 计算位置
                y_pos = start_y + (i * (box_height + 0.3))

                # 添加问题框
                left = Inches(0.8)
                top = Inches(y_pos)
                width = Inches(8.4)
                height = Inches(box_height)

                shape = slide.shapes.add_shape(
                    1,  # msoShapeRectangle
                    left, top, width, height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = level_color
                shape.fill.transparency = 0.85
                shape.line.color.rgb = level_color
                shape.line.width = Pt(2)

                # 添加问题文本
                text_frame = shape.text_frame
                text_frame.word_wrap = True

                p = text_frame.paragraphs[0]
                p.text = f"{icon} 【{question_item.get('chinese_name', level)}】{question_text}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(50, 50, 50)

                # 添加层级说明
                if question_item.get("description"):
                    p = text_frame.add_paragraph()
                    p.text = f"  {question_item.get('description')}"
                    p.font.size = Pt(font_size - 4)
                    p.font.color.rgb = RGBColor(100, 100, 100)
                    p.space_before = Pt(4)

    def _add_metacognitive_bubble_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """
        添加思考气泡页（feat-038）
        元认知提示：云朵形状思考气泡 + 特殊配色，与正常内容视觉区分
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "想一想")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # 元认知提示类型图标
        meta_icons = {
            "reflect": "💭",
            "connect": "🔗",
            "predict": "🔮",
            "evaluate": "✓",
            "extend": "🚀",
            "monitor": "📊",
        }

        # 背景色（淡色）
        bg_color = RGBColor(240, 248, 255)  # 淡蓝色背景

        # 获取元认知提示数据
        meta_prompt = slide_data.get("meta_prompt", {})
        icon = meta_icons.get(meta_prompt.get("type", "reflect"), "💭")
        content = meta_prompt.get("content", slide_data.get("content", ["停下来思考一下..."])[0] if isinstance(slide_data.get("content"), list) else "停下来思考一下...")

        # 添加云朵形状背景（使用椭圆近似）
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        # 添加圆角矩形作为思考气泡背景
        bubble = slide.shapes.add_shape(
            5,  # msoShapeRoundedRectangle
            left, top, width, height
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = bg_color
        bubble.fill.transparency = 0.3
        bubble.line.color.rgb = color
        bubble.line.width = Pt(3)
        bubble.line.dash_style = 2  # 虚线边框

        # 添加图标和文本
        text_frame = bubble.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = f"{icon} {content}"
        p.font.size = Pt(font_size + 2)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.alignment = 1  # 居中

        # 如果是反思型提示，添加额外装饰
        if meta_prompt.get("type") == "reflect":
            # 添加小云朵装饰在左上角
            small_cloud = slide.shapes.add_shape(
                5,
                Inches(0.8), Inches(1.7), Inches(1), Inches(0.6)
            )
            small_cloud.fill.solid()
            small_cloud.fill.fore_color.rgb = bg_color
            small_cloud.line.color.rgb = color
            small_cloud.line.width = Pt(2)

    def _add_reflection_checkpoint_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor
    ):
        """
        添加反思 checkpoint 页（feat-038）
        学习 checkpoints，促进学生自我监控理解程度
        """
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "学习检查点")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = color
        title_para.font.bold = True

        # checkpoint 配置
        checkpoint_color = RGBColor(106, 90, 205)  # 紫色表示检查点
        bg_color = RGBColor(250, 250, 250)  # 浅灰背景

        # 获取检查点列表
        checkpoints = slide_data.get("checkpoints", [])

        if not checkpoints and slide_data.get("content"):
            # 兼容 content 格式
            checkpoints = [{"text": item} for item in slide_data.get("content", [])]

        start_y = 1.8
        box_height = 1.0

        for i, checkpoint_item in enumerate(checkpoints):
            if isinstance(checkpoint_item, dict):
                checkpoint_text = checkpoint_item.get("text", checkpoint_item.get("content", ""))
                checkpoint_type = checkpoint_item.get("type", "check")
                difficulty = checkpoint_item.get("difficulty", 1)

                # 计算位置
                y_pos = start_y + (i * (box_height + 0.4))

                # 添加检查点框
                left = Inches(0.8)
                top = Inches(y_pos)
                width = Inches(8.4)
                height = Inches(box_height)

                box = slide.shapes.add_shape(
                    1,  # msoShapeRectangle
                    left, top, width, height
                )
                box.fill.solid()
                box.fill.fore_color.rgb = bg_color
                box.fill.transparency = 0.5
                box.line.color.rgb = checkpoint_color
                box.line.width = Pt(2)

                # 添加检查点图标
                icon = "✓" if checkpoint_type == "check" else "⚠" if checkpoint_type == "warn" else "?"
                icon_box = slide.shapes.add_shape(
                    24,  # msoShapeEllipse
                    Inches(0.9), Inches(y_pos + 0.15), Inches(0.5), Inches(0.5)
                )
                icon_box.fill.solid()
                icon_box.fill.fore_color.rgb = checkpoint_color
                icon_box.line.color.rgb = checkpoint_color

                icon_text = icon_box.text_frame
                icon_text.paragraphs[0].text = icon
                icon_text.paragraphs[0].alignment = 1
                icon_text.paragraphs[0].font.size = Pt(font_size)
                icon_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                icon_text.paragraphs[0].font.bold = True

                # 添加检查点文本
                text_frame = box.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.7)

                p = text_frame.paragraphs[0]
                p.text = checkpoint_text
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(50, 50, 50)

                # 添加难度星级
                stars = "★" * difficulty + "☆" * (3 - difficulty)
                p = text_frame.add_paragraph()
                p.text = f"难度：{stars}"
                p.font.size = Pt(font_size - 4)
                p.font.color.rgb = RGBColor(150, 150, 150)
                p.space_before = Pt(4)

    def _add_scenario_import_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        font_size: int,
        color: RGBColor,
        secondary_color: RGBColor
    ):
        """添加情境导入页（feat-039）：贴近学生生活的情境导入，全幅背景色 + 大字引导问题"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 获取情境配置
        scenario = slide_data.get('scenario', {})
        scenario_description = scenario.get('scenario_description', slide_data.get('content', [''])[0] if slide_data.get('content') else '')
        guiding_question = scenario.get('guiding_question', '思考问题')
        visual_suggestion = scenario.get('visual_suggestion', '')
        scenario_type = scenario.get('scenario_type', '生活实例')

        # 标题
        title_text = slide_data.get('title', '情境导入')
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 4)
        title_para.font.color.rgb = color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 情境类型标签
        type_icons = {
            '生活实例': '🏠',
            '童话故事': '📚',
            '动画人物': '🎬',
            '校园场景': '🏫',
            '社会实践': '🤝',
            '科研探索': '🔬',
            '职业体验': '💼',
            '购物消费': '🛒',
            '体育运动': '⚽',
            '旅行游记': '✈️',
            '游乐场': '🎢',
            '交通工具': '🚗',
        }
        type_icon = type_icons.get(scenario_type, '💡')

        type_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.5))
        type_frame = type_box.text_frame
        type_frame.text = f'{type_icon} {scenario_type}'
        type_para = type_frame.paragraphs[0]
        type_para.font.size = Pt(font_size - 2)
        type_para.font.color.rgb = secondary_color
        type_para.alignment = PP_ALIGN.CENTER

        # 情境描述区域（带背景色块）
        scenario_box = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            Inches(0.5), Inches(1.8), Inches(9), Inches(2.5)
        )
        scenario_box.fill.solid()
        scenario_box.fill.fore_color.rgb = color
        scenario_box.fill.transparency = 0.85  # 高透明度背景
        scenario_box.line.color.rgb = color
        scenario_box.line.width = Pt(2)

        scenario_frame = scenario_box.text_frame
        scenario_frame.word_wrap = True
        scenario_frame.margin_left = Inches(0.5)
        scenario_frame.margin_right = Inches(0.5)
        scenario_frame.margin_top = Inches(0.3)
        scenario_frame.margin_bottom = Inches(0.3)

        scenario_para = scenario_frame.paragraphs[0]
        scenario_para.text = scenario_description
        scenario_para.font.size = Pt(font_size + 1)
        scenario_para.font.color.rgb = RGBColor(50, 50, 50)

        # 引导问题（大字醒目）
        question_box = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            Inches(0.5), Inches(4.5), Inches(9), Inches(2)
        )
        question_box.fill.solid()
        question_box.fill.fore_color.rgb = secondary_color
        question_box.fill.transparency = 0.8
        question_box.line.color.rgb = secondary_color
        question_box.line.width = Pt(3)

        question_frame = question_box.text_frame
        question_frame.word_wrap = True
        question_frame.margin_left = Inches(0.5)
        question_frame.margin_right = Inches(0.5)
        question_frame.margin_top = Inches(0.3)
        question_frame.margin_bottom = Inches(0.3)

        question_para = question_frame.paragraphs[0]
        question_para.text = f'❓ {guiding_question}'
        question_para.font.size = Pt(font_size + 3)
        question_para.font.color.rgb = RGBColor(50, 50, 50)
        question_para.font.bold = True
        question_para.alignment = PP_ALIGN.CENTER

        # 视觉建议（底部小字）
        if visual_suggestion:
            visual_box = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(8), Inches(0.5))
            visual_frame = visual_box.text_frame
            visual_frame.text = f'🎨 视觉建议：{visual_suggestion}'
            visual_para = visual_frame.paragraphs[0]
            visual_para.font.size = Pt(font_size - 3)
            visual_para.font.color.rgb = RGBColor(120, 120, 120)
            visual_para.font.italic = True
            visual_para.alignment = PP_ALIGN.CENTER


    def merge_ppts(self, ppt_paths: List[Path], output_path: Path, title: str = "合并课件") -> Path:
        """
        合并多个 PPT 文件为一个新课件

        Args:
            ppt_paths: 要合并的 PPT 文件路径列表
            output_path: 输出文件路径
            title: 合并后课件的标题（用于封面页）

        Returns:
            合并后的 PPT 文件路径
        """
        try:
            from pptx import Presentation

            # 创建新的 PPT
            merged_prs = Presentation()
            merged_prs.slide_width = Inches(10)
            merged_prs.slide_height = Inches(7.5)

            # 添加封面页
            self._add_cover_slide(
                merged_prs,
                title,
                40,  # 标题字号
                COLOR_BLUE_6,  # 主色
                False  # 不需要拼音
            )

            # 遍历每个 PPT 文件，复制所有幻灯片
            total_slides = 0
            for ppt_path in ppt_paths:
                if not ppt_path.exists():
                    logger.warning(f"PPT 文件不存在：{ppt_path}")
                    continue

                try:
                    source_prs = Presentation(ppt_path)

                    # 跳过第一个幻灯片（假设是封面），避免重复
                    start_idx = 1 if len(source_prs.slides) > 1 else 0

                    for slide_idx in range(start_idx, len(source_prs.slides)):
                        source_slide = source_prs.slides[slide_idx]

                        # 复制幻灯片内容到新 PPT
                        # 注意：python-pptx 不直接支持复制幻灯片，需要手动复制内容
                        self._copy_slide_content(source_slide, merged_prs)
                        total_slides += 1

                except Exception as e:
                    logger.error(f"复制 PPT {ppt_path} 时出错：{e}")
                    continue

            # 保存合并后的 PPT
            output_path.parent.mkdir(parents=True, exist_ok=True)
            merged_prs.save(output_path)

            logger.info(f"成功合并 {len(ppt_paths)} 个 PPT 文件，共 {total_slides} 页幻灯片")
            return output_path

        except Exception as e:
            logger.error(f"合并 PPT 失败：{e}")
            raise

    def _copy_slide_content(self, source_slide, target_prs: Presentation):
        """
        复制幻灯片内容到新 PPT

        Args:
            source_slide: 源幻灯片
            target_prs: 目标 PPT 对象
        """
        try:
            # 添加空白幻灯片
            blank_layout = target_prs.slide_layouts[6]  # 空白布局
            target_slide = target_prs.slides.add_slide(blank_layout)

            # 复制所有形状
            for shape in source_slide.shapes:
                if shape.shape_type == 14:  # msoShapeType['placeHolder']
                    # 跳过占位符
                    continue
                elif shape.shape_type == 5:  # msoShapeType['textBox']
                    # 复制文本框
                    self._copy_textbox(shape, target_slide)
                elif shape.shape_type == 13:  # msoShapeType['picture']
                    # 复制图片（暂不支持，跳过）
                    logger.debug("跳过图片复制")
                    continue
                # 其他形状暂不支持

        except Exception as e:
            logger.error(f"复制幻灯片内容失败：{e}")

    def _copy_textbox(self, source_shape, target_slide):
        """
        复制文本框到目标幻灯片

        Args:
            source_shape: 源文本框形状
            target_slide: 目标幻灯片
        """
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            # 获取源文本框属性
            src_left = source_shape.left
            src_top = source_shape.top
            src_width = source_shape.width
            src_height = source_shape.height

            # 在目标幻灯片创建文本框
            target_box = target_slide.shapes.add_textbox(
                src_left, src_top, src_width, src_height
            )
            target_frame = target_box.text_frame

            # 复制文本内容和格式
            source_frame = source_shape.text_frame
            if source_frame.text:
                target_frame.text = source_frame.text

                # 复制段落格式
                for i, src_para in enumerate(source_frame.paragraphs):
                    if i < len(target_frame.paragraphs):
                        tgt_para = target_frame.paragraphs[i]
                        tgt_para.font.size = src_para.font.size
                        tgt_para.font.bold = src_para.font.bold
                        tgt_para.font.italic = src_para.font.italic
                        if src_para.font.color.rgb:
                            tgt_para.font.color.rgb = src_para.font.color.rgb
                        tgt_para.alignment = src_para.alignment
                        tgt_para.level = src_para.level

        except Exception as e:
            logger.error(f"复制文本框失败：{e}")

    def smart_merge_ppts(
        self,
        ppt_a_path: Path,
        ppt_b_path: Path,
        output_path: Path,
        merge_strategy: Dict[str, Any],
        title: str = "智能合并课件"
    ) -> Path:
        """
        根据 LLM 生成的合并策略，智能合并两个 PPT 文件

        Args:
            ppt_a_path: PPT A 文件路径
            ppt_b_path: PPT B 文件路径
            output_path: 输出文件路径
            merge_strategy: 合并策略，包含：
                - slides_to_merge: 要合并的页面对象列表
                - slides_to_skip_a: 从 A 中跳过的页码
                - slides_to_skip_b: 从 B 中跳过的页码
                - global_adjustments: 全局调整说明
            title: 合并后课件的标题

        Returns:
            合并后的 PPT 文件路径
        """
        try:
            from pptx import Presentation

            # 创建新的 PPT
            merged_prs = Presentation()
            merged_prs.slide_width = Inches(10)
            merged_prs.slide_height = Inches(7.5)

            # 添加封面页
            self._add_cover_slide(
                merged_prs,
                title,
                40,
                COLOR_BLUE_6,
                False
            )

            # 加载两个 PPT
            prs_a = Presentation(ppt_a_path) if ppt_a_path.exists() else None
            prs_b = Presentation(ppt_b_path) if ppt_b_path.exists() else None

            if not prs_a and not prs_b:
                raise ValueError("两个 PPT 文件都不存在")

            # 解析合并策略
            slides_to_merge = merge_strategy.get("slides_to_merge", [])
            skip_a = set(merge_strategy.get("slides_to_skip_a", []))
            skip_b = set(merge_strategy.get("slides_to_skip_b", []))

            # 执行合并操作
            total_slides = 0

            # 1. 处理策略中明确指定的合并操作
            for merge_item in slides_to_merge:
                from_a = merge_item.get("from_a", [])
                from_b = merge_item.get("from_b", [])
                action = merge_item.get("action", "append")
                instruction = merge_item.get("instruction", "")

                if action == "combine":
                    # 合并 A 和 B 的页面内容到一页
                    self._merge_slide_content(
                        merged_prs,
                        prs_a, from_a,
                        prs_b, from_b,
                        instruction
                    )
                    total_slides += 1
                elif action == "append_a":
                    # 追加 A 的页面
                    for slide_idx in from_a:
                        if slide_idx not in skip_a and prs_a:
                            self._copy_slide_from_prs(prs_a, slide_idx - 1, merged_prs)
                            total_slides += 1
                elif action == "append_b":
                    # 追加 B 的页面
                    for slide_idx in from_b:
                        if slide_idx not in skip_b and prs_b:
                            self._copy_slide_from_prs(prs_b, slide_idx - 1, merged_prs)
                            total_slides += 1

            # 2. 处理未被策略覆盖的页面（默认追加）
            if prs_a:
                for i in range(len(prs_a.slides)):
                    slide_num = i + 1
                    # 跳过封面和已处理的页面
                    if slide_num == 1 or slide_num in skip_a:
                        continue
                    # 检查是否已在策略中处理
                    already_processed = any(
                        slide_num in merge_item.get("from_a", [])
                        for merge_item in slides_to_merge
                    )
                    if not already_processed:
                        self._copy_slide_from_prs(prs_a, i, merged_prs)
                        total_slides += 1

            if prs_b:
                for i in range(len(prs_b.slides)):
                    slide_num = i + 1
                    # 跳过封面和已处理的页面
                    if slide_num == 1 or slide_num in skip_b:
                        continue
                    # 检查是否已在策略中处理
                    already_processed = any(
                        slide_num in merge_item.get("from_b", [])
                        for merge_item in slides_to_merge
                    )
                    if not already_processed:
                        self._copy_slide_from_prs(prs_b, i, merged_prs)
                        total_slides += 1

            # 保存合并后的 PPT
            output_path.parent.mkdir(parents=True, exist_ok=True)
            try:
                merged_prs.save(output_path)
                # 验证文件可以被正确读取
                from pptx import Presentation
                test_prs = Presentation(output_path)
                logger.info(f"智能合并完成：{total_slides} 页幻灯片，文件验证成功 ({output_path.name})")
            except Exception as save_error:
                logger.error(f"文件保存或验证失败: {save_error}")
                raise

            return output_path

        except Exception as e:
            logger.error(f"智能合并 PPT 失败：{e}")
            raise

    def _merge_slide_content(
        self,
        target_prs: Presentation,
        prs_a: Optional[Presentation], indices_a: List[int],
        prs_b: Optional[Presentation], indices_b: List[int],
        instruction: str
    ):
        """
        将多个幻灯片内容合并到一页

        Args:
            target_prs: 目标 PPT
            prs_a: PPT A
            indices_a: A 中要合并的页码（从 1 开始）
            prs_b: PPT B
            indices_b: B 中要合并的页码（从 1 开始）
            instruction: 合并指令说明
        """
        try:
            # 添加空白幻灯片
            blank_layout = target_prs.slide_layouts[6]
            target_slide = target_prs.slides.add_slide(blank_layout)

            # 设置合并后的标题（使用指令中的说明或默认标题）
            title_box = target_slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = instruction[:50] if instruction else "合并内容"
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.font.color.rgb = COLOR_BLUE_6

            # 垂直分割线
            divider = target_slide.shapes.add_shape(
                1,  # msoShapeRectangle
                Inches(0.5), Inches(1.1), Inches(9), Inches(0.02)
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = COLOR_BLUE_1
            divider.line.fill.background()

            # 复制 A 的内容（上半部分）
            content_y = Inches(1.3)
            if prs_a:
                for idx in indices_a:
                    if 0 < idx <= len(prs_a.slides):
                        source_slide = prs_a.slides[idx - 1]
                        for shape in source_slide.shapes:
                            if shape.has_text_frame:
                                box = target_slide.shapes.add_textbox(
                                    Inches(0.5), content_y, Inches(8.5), Inches(1.5)
                                )
                                frame = box.text_frame
                                frame.text = shape.text[:200]  # 限制长度
                                for para in frame.paragraphs:
                                    para.font.size = Pt(14)
                                content_y += Inches(1.7)

            # 复制 B 的内容（下半部分）
            if prs_b:
                for idx in indices_b:
                    if 0 < idx <= len(prs_b.slides):
                        source_slide = prs_b.slides[idx - 1]
                        for shape in source_slide.shapes:
                            if shape.has_text_frame:
                                box = target_slide.shapes.add_textbox(
                                    Inches(0.5), content_y, Inches(8.5), Inches(1.5)
                                )
                                frame = box.text_frame
                                frame.text = shape.text[:200]
                                for para in frame.paragraphs:
                                    para.font.size = Pt(14)
                                content_y += Inches(1.7)

        except Exception as e:
            logger.error(f"合并幻灯片内容失败：{e}")

    def _copy_slide_from_prs(
        self,
        source_prs: Presentation,
        slide_index: int,
        target_prs: Presentation
    ):
        """
        从 PPT 复制指定索引的幻灯片到目标 PPT

        Args:
            source_prs: 源 PPT
            slide_index: 幻灯片索引（从 0 开始）
            target_prs: 目标 PPT
        """
        try:
            if slide_index < 0 or slide_index >= len(source_prs.slides):
                return

            source_slide = source_prs.slides[slide_index]
            self._copy_slide_content(source_slide, target_prs)

        except Exception as e:
            logger.error(f"复制幻灯片失败：{e}")


# 全局 PPT 生成器实例
_ppt_generator_instance: Optional[PPTGenerator] = None

def get_ppt_generator() -> PPTGenerator:
    """获取 PPT 生成器单例"""
    global _ppt_generator_instance
    if _ppt_generator_instance is None:
        _ppt_generator_instance = PPTGenerator()
    return _ppt_generator_instance


def generate_ppt_from_versions(
    merged_slides: List[Dict[str, Any]],
    title: str,
    grade: str = "6",
    subject: str = "general",
    style: str = "simple"
) -> Path:
    """
    基于版本数据生成最终 PPT

    Args:
        merged_slides: 合并的页面列表，每项包含：
            - source_pptx: 源 PPTX 路径
            - slide_index: 页码（0-indexed）
            - content_snapshot: AI 修改的内容快照（v2+ 版本）
            - version: 版本号（v1, v2, ...）
        title: 最终 PPT 标题
        grade: 年级
        subject: 学科
        style: 风格

    Returns:
        生成的 PPT 文件路径

    策略：
    - v1 版本（原始上传）：从源 PPTX 复制页面
    - v2+ 版本（AI 修改）：从 content_snapshot 重建页面
    """
    from pptx import Presentation
    from uuid import uuid4

    logger.info(f"生成最终 PPT: {len(merged_slides)} 页，title={title}")

    # 创建空白 PPT
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 获取样式配置
    grade_style = PPTStyle.get_style(grade)
    title_size = grade_style["title_size"]
    content_size = grade_style["font_size"]

    # 添加封面页
    generator = get_ppt_generator()
    generator._add_cover_slide(prs, title, title_size, grade_style["primary"], False)

    # 处理每个页面
    for slide_info in merged_slides:
        source_pptx = slide_info["source_pptx"]
        slide_index = slide_info["slide_index"]
        content_snapshot = slide_info.get("content_snapshot")
        version = slide_info.get("version", "v1")

        try:
            if content_snapshot:
                # v2+ 版本：从 content_snapshot 重建页面
                logger.info(f"从 content_snapshot 重建页面 {slide_index} ({version})")
                generator._add_slide_from_snapshot(prs, content_snapshot, content_size)
            else:
                # v1 版本：从源 PPTX 复制页面
                logger.info(f"从源 PPTX 复制页面 {slide_index} ({version})")
                generator._copy_slide_from_pptx(prs, source_pptx, slide_index)
        except Exception as e:
            logger.warning(f"处理页面 {slide_index} 失败：{e}，使用占位页面")
            # 使用占位页面
            _add_placeholder_slide(prs, f"页面 {slide_index + 1} (版本 {version})", content_size)

    # 保存文件
    from ..config import settings
    output_dir = settings.UPLOAD_DIR / "generated"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"final_{uuid4().hex[:8]}.pptx"

    prs.save(output_path)
    logger.info(f"最终 PPT 已保存：{output_path}")

    return output_path


def _add_placeholder_slide(prs: Presentation, text: str, font_size: int):
    """添加占位页面（当页面复制失败时使用）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "页面加载失败"
    content = slide.placeholders[1]
    content.text = text
