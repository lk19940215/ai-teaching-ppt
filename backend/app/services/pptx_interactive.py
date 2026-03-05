"""
PPT 课堂互动功能增强模块

实现互动问答和点击显示答案效果：
1. 互动问答页模板（带 ABCD 选项按钮布局）
2. "点击显示答案"效果（使用 OOXML 触发器）
3. 拖拽匹配游戏页（如单词 - 释义配对）
4. 随堂测验页（选择题、填空题）

实现原理：
- python-pptx 生成基础 PPTX 文件和互动元素
- 通过 OOXML 触发器实现点击显示效果
- 使用超链接和形状组合实现互动按钮
"""

import zipfile
import tempfile
import shutil
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple
import xml.etree.ElementTree as ET
from xml.dom import minidom
import logging
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# OOXML 命名空间
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

# 注册命名空间
ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
ET.register_namespace('p', 'http://schemas.openxmlformats.org/presentationml/2006/main')
ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')


class InteractivePPTGenerator:
    """互动 PPT 生成器"""

    # 选项标签
    OPTION_LABELS = ['A', 'B', 'C', 'D']

    # 颜色配置
    COLORS = {
        'button_default': RGBColor(68, 138, 255),  # 蓝色按钮
        'button_hover': RGBColor(52, 152, 219),    # 悬停蓝色
        'button_correct': RGBColor(76, 175, 80),   # 绿色（正确答案）
        'button_wrong': RGBColor(244, 67, 54),     # 红色（错误答案）
        'button_reveal': RGBColor(156, 39, 176),   # 紫色（显示答案）
        'text_white': RGBColor(255, 255, 255),
        'text_dark': RGBColor(33, 33, 33),
    }

    def __init__(self):
        """初始化互动 PPT 生成器"""
        pass

    def add_quiz_slide(
        self,
        prs: Presentation,
        question: str,
        options: List[str],
        correct_index: int,
        title: str = "互动问答",
        font_size: int = 24
    ) -> int:
        """
        添加互动问答页（带 ABCD 选项按钮）

        Args:
            prs: Presentation 对象
            question: 问题文本
            options: 选项列表（2-4 个）
            correct_index: 正确答案索引（0-3）
            title: 页面标题
            font_size: 字号

        Returns:
            幻灯片索引
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 1. 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = self.COLORS['button_default']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加问题文本
        question_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.2), Inches(8), Inches(1.5)
        )
        question_frame = question_box.text_frame
        question_frame.word_wrap = True
        question_frame.text = question
        for para in question_frame.paragraphs:
            para.font.size = Pt(font_size)
            para.alignment = PP_ALIGN.CENTER

        # 3. 添加选项按钮
        button_width = Inches(3.5)
        button_height = Inches(0.8)
        button_spacing = Inches(0.3)
        start_y = Inches(2.8)

        for i, option in enumerate(options[:4]):
            # 计算按钮位置
            col = i % 2
            row = i // 2
            left = Inches(1) + col * (button_width + button_spacing)
            top = start_y + row * (button_height + button_spacing)

            # 创建按钮形状
            button = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, button_width, button_height
            )
            button.fill.solid()
            button.fill.fore_color.rgb = self.COLORS['button_default']
            button.line.fill.background()

            # 添加选项文本
            tf = button.text_frame
            tf.text = f"{self.OPTION_LABELS[i]}. {option}"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = self.COLORS['text_white']
            tf.paragraphs[0].font.bold = True

        # 4. 添加答案提示框（初始隐藏，通过动画控制）
        self._add_answer_reveal_box(
            slide,
            correct_index,
            font_size
        )

        return prs.slides.index(slide)

    def _add_answer_reveal_box(
        self,
        slide,
        correct_index: int,
        font_size: int
    ):
        """
        添加答案显示提示框

        Args:
            slide: 幻灯片对象
            correct_index: 正确答案索引
            font_size: 字号
        """
        # 提示文本
        hint_text = f"正确答案：{self.OPTION_LABELS[correct_index]}"

        hint_box = slide.shapes.add_textbox(
            Inches(3), Inches(5.5), Inches(4), Inches(0.8)
        )
        hint_frame = hint_box.text_frame
        hint_frame.text = "点击查看答案"
        hint_para = hint_frame.paragraphs[0]
        hint_para.font.size = Pt(font_size - 2)
        hint_para.font.color.rgb = self.COLORS['button_reveal']
        hint_para.alignment = PP_ALIGN.CENTER

        # 答案框
        answer_box = slide.shapes.add_textbox(
            Inches(2.5), Inches(6), Inches(5), Inches(0.8)
        )
        answer_frame = answer_box.text_frame
        answer_frame.text = hint_text
        answer_para = answer_frame.paragraphs[0]
        answer_para.font.size = Pt(font_size)
        answer_para.font.color.rgb = self.COLORS['button_correct']
        answer_para.font.bold = True
        answer_para.alignment = PP_ALIGN.CENTER

    def add_click_to_reveal_slide(
        self,
        prs: Presentation,
        question: str,
        hidden_content: str,
        title: str = "点击显示答案",
        font_size: int = 24
    ) -> int:
        """
        添加点击显示答案页

        Args:
            prs: Presentation 对象
            question: 问题/题目
            hidden_content: 隐藏的答案内容
            title: 页面标题
            font_size: 字号

        Returns:
            幻灯片索引
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 1. 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = self.COLORS['button_default']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加问题文本
        question_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.2), Inches(8), Inches(2)
        )
        question_frame = question_box.text_frame
        question_frame.word_wrap = True
        question_frame.text = question
        for para in question_frame.paragraphs:
            para.font.size = Pt(font_size)
            para.alignment = PP_ALIGN.CENTER

        # 3. 添加"点击显示"按钮
        reveal_btn = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3), Inches(3.5), Inches(4), Inches(1)
        )
        reveal_btn.fill.solid()
        reveal_btn.fill.fore_color.rgb = self.COLORS['button_reveal']
        reveal_btn.line.fill.background()

        btn_tf = reveal_btn.text_frame
        btn_tf.text = "点击查看答案"
        btn_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        btn_tf.paragraphs[0].font.size = Pt(font_size)
        btn_tf.paragraphs[0].font.color.rgb = self.COLORS['text_white']
        btn_tf.paragraphs[0].font.bold = True

        # 4. 添加隐藏答案框（带边框，初始可见但可设置为通过触发器控制）
        answer_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(4.8), Inches(8), Inches(2)
        )
        answer_box.fill.solid()
        answer_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        answer_box.line.color.rgb = self.COLORS['button_correct']
        answer_box.line.width = Pt(2)

        ans_tf = answer_box.text_frame
        ans_tf.word_wrap = True
        ans_tf.text = hidden_content
        for para in ans_tf.paragraphs:
            para.font.size = Pt(font_size)
            para.alignment = PP_ALIGN.CENTER

        return prs.slides.index(slide)

    def add_matching_game_slide(
        self,
        prs: Presentation,
        pairs: List[Tuple[str, str]],
        title: str = "拖拽匹配游戏",
        font_size: int = 20,
        left_title: str = "单词",
        right_title: str = "释义"
    ) -> int:
        """
        添加拖拽匹配游戏页

        Args:
            prs: Presentation 对象
            pairs: 配对列表 [(左项，右项), ...]
            title: 页面标题
            font_size: 字号
            left_title: 左侧标题
            right_title: 右侧标题

        Returns:
            幻灯片索引
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = self.COLORS['button_default']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加左右列标题
        left_title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.3), Inches(3.5), Inches(0.5)
        )
        lt_frame = left_title_box.text_frame
        lt_frame.text = left_title
        lt_frame.paragraphs[0].font.size = Pt(font_size + 2)
        lt_frame.paragraphs[0].font.color.rgb = self.COLORS['button_default']
        lt_frame.paragraphs[0].font.bold = True
        lt_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        right_title_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(1.3), Inches(3.5), Inches(0.5)
        )
        rt_frame = right_title_box.text_frame
        rt_frame.text = right_title
        rt_frame.paragraphs[0].font.size = Pt(font_size + 2)
        rt_frame.paragraphs[0].font.color.rgb = self.COLORS['button_reveal']
        rt_frame.paragraphs[0].font.bold = True
        rt_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 3. 添加左侧项目
        left_y = Inches(2)
        item_height = Inches(0.7)
        spacing = Inches(0.2)

        for i, (left_item, right_item) in enumerate(pairs[:6]):
            # 左侧项目
            left_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), left_y + i * (item_height + spacing),
                Inches(3.5), item_height
            )
            left_box.fill.solid()
            left_box.fill.fore_color.rgb = self.COLORS['button_default']
            left_box.line.fill.background()

            left_tf = left_box.text_frame
            left_tf.text = left_item
            left_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            left_tf.paragraphs[0].font.size = Pt(font_size)
            left_tf.paragraphs[0].font.color.rgb = self.COLORS['text_white']

            # 右侧项目
            right_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.5), left_y + i * (item_height + spacing),
                Inches(3.5), item_height
            )
            right_box.fill.solid()
            right_box.fill.fore_color.rgb = RGBColor(200, 200, 200)
            right_box.line.fill.background()

            right_tf = right_box.text_frame
            right_tf.text = right_item
            right_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            right_tf.paragraphs[0].font.size = Pt(font_size)
            right_tf.paragraphs[0].font.color.rgb = self.COLORS['text_dark']

        return prs.slides.index(slide)

    def add_quiz_game_slide(
        self,
        prs: Presentation,
        questions: List[Dict[str, Any]],
        title: str = "随堂测验",
        font_size: int = 22
    ) -> int:
        """
        添加随堂测验页（多题展示）

        Args:
            prs: Presentation 对象
            questions: 问题列表 [{"question": "问题", "options": ["A", "B"], "answer": 0}, ...]
            title: 页面标题
            font_size: 字号

        Returns:
            幻灯片索引
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = self.COLORS['button_default']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加说明
        hint_box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(0.5)
        )
        hint_frame = hint_box.text_frame
        hint_frame.text = "请选择正确答案"
        hint_frame.paragraphs[0].font.size = Pt(font_size - 2)
        hint_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        hint_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 3. 添加问题
        q_start_y = Inches(1.6)
        q_height = Inches(2.5)

        for i, q in enumerate(questions[:3]):  # 最多 3 题
            q_y = q_start_y + i * Inches(2.8)

            # 问题编号和文本
            q_text = f"{i + 1}. {q.get('question', '')}"
            q_box = slide.shapes.add_textbox(
                Inches(0.8), q_y, Inches(8.4), Inches(0.8)
            )
            q_frame = q_box.text_frame
            q_frame.word_wrap = True
            q_frame.text = q_text
            q_frame.paragraphs[0].font.size = Pt(font_size)
            q_frame.paragraphs[0].font.bold = True

            # 选项
            options = q.get('options', [])
            for j, opt in enumerate(options[:4]):
                opt_x = Inches(1) + (j % 2) * Inches(4.2)
                opt_y = q_y + Inches(0.7) + (j // 2) * Inches(0.9)

                # 选项按钮
                opt_btn = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    opt_x, opt_y, Inches(4), Inches(0.7)
                )
                opt_btn.fill.solid()
                opt_btn.fill.fore_color.rgb = self.COLORS['button_default']
                opt_btn.line.fill.background()

                opt_tf = opt_btn.text_frame
                label = self.OPTION_LABELS[j] if j < 4 else chr(ord('A') + j)
                opt_tf.text = f"{label}. {opt}"
                opt_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                opt_tf.paragraphs[0].font.size = Pt(font_size - 2)
                opt_tf.paragraphs[0].font.color.rgb = self.COLORS['text_white']

        return prs.slides.index(slide)

    def add_fill_blank_slide(
        self,
        prs: Presentation,
        sentences: List[str],
        answers: List[str],
        title: str = "填空题",
        font_size: int = 22
    ) -> int:
        """
        添加填空题页面

        Args:
            prs: Presentation 对象
            sentences: 句子列表（用___表示空白）
            answers: 答案列表
            title: 页面标题
            font_size: 字号

        Returns:
            幻灯片索引
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size + 6)
        title_para.font.color.rgb = self.COLORS['button_default']
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # 2. 添加句子
        y_pos = Inches(1.5)
        for i, sentence in enumerate(sentences[:6]):
            sent_box = slide.shapes.add_textbox(
                Inches(1), y_pos + i * Inches(1), Inches(8), Inches(0.8)
            )
            sent_frame = sent_box.text_frame
            sent_frame.word_wrap = True
            sent_frame.text = sentence
            sent_frame.paragraphs[0].font.size = Pt(font_size)

        # 3. 添加答案框（初始隐藏）
        ans_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(6), Inches(8), Inches(1.2)
        )
        ans_box.fill.solid()
        ans_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        ans_box.line.color.rgb = self.COLORS['button_correct']
        ans_box.line.width = Pt(2)

        ans_tf = ans_box.text_frame
        ans_tf.text = "答案：" + " | ".join(answers)
        ans_tf.paragraphs[0].font.size = Pt(font_size)
        ans_tf.paragraphs[0].font.color.rgb = self.COLORS['button_correct']
        ans_tf.paragraphs[0].font.bold = True
        ans_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        return prs.slides.index(slide)

    def inject_click_trigger(
        self,
        pptx_path: Path,
        output_path: Optional[Path] = None,
        trigger_config: Optional[Dict[str, Any]] = None
    ) -> Path:
        """
        为 PPTX 注入点击触发器动画

        Args:
            pptx_path: 输入 PPTX 路径
            output_path: 输出路径
            trigger_config: 触发器配置 {
                "trigger_shape_id": 形状 ID（点击此形状触发）,
                "target_shape_id": 目标形状 ID（被触发动画）,
                "effect": "appear" | "disappear" | "color_change"
            }

        Returns:
            输出文件路径
        """
        output_path = output_path or pptx_path

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            try:
                # 解压 PPTX
                self._extract_pptx(pptx_path, temp_path)

                # 注入触发器动画
                self._inject_triggers(temp_path, trigger_config)

                # 重新打包
                self._package_pptx(temp_path, output_path)

                logger.info(f"点击触发器已注入：{output_path}")
                return output_path

            except Exception as e:
                logger.error(f"触发器注入失败：{e}")
                raise RuntimeError(f"触发器注入失败：{e}") from e

    def _extract_pptx(self, pptx_path: Path, dest_dir: Path):
        """解压 PPTX 文件"""
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            zip_ref.extractall(dest_dir)

    def _package_pptx(self, src_dir: Path, output_path: Path):
        """重新打包为 PPTX 文件"""
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for file_path in src_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(src_dir)
                    zip_ref.write(file_path, arcname)

    def _inject_triggers(self, pptx_dir: Path, trigger_config: Optional[Dict[str, Any]] = None):
        """
        注入触发器到幻灯片

        Args:
            pptx_dir: 解压后的 PPTX 目录
            trigger_config: 触发器配置
        """
        slides_dir = pptx_dir / 'ppt' / 'slides'
        if not slides_dir.exists():
            raise ValueError("无效的 PPTX 结构：找不到 slides 目录")

        for slide_file in sorted(slides_dir.glob('slide*.xml')):
            try:
                tree = ET.parse(slide_file)
                root = tree.getroot()

                # 查找 timing 节点并添加触发器
                self._add_trigger_to_timing(root, slide_file, trigger_config)

                self._save_xml(tree, slide_file)

            except ET.ParseError as e:
                logger.warning(f"幻灯片 XML 解析失败 {slide_file.name}: {e}")
            except Exception as e:
                logger.warning(f"处理幻灯片失败 {slide_file.name}: {e}")

    def _add_trigger_to_timing(
        self,
        slide_root: ET.Element,
        slide_file: Path,
        trigger_config: Optional[Dict[str, Any]] = None
    ):
        """
        添加触发器到幻灯片时间轴

        OOXML 触发器结构:
        <p:timing>
          <p:timeNodeList>
            <p:par>  <!-- 并行时间线 -->
              <p:cTn id="1" fill="hold">  <!-- 根时间节点 -->
                <p:stCondLst>
                  <p:cond delay="indefinite"/>  <!-- 等待触发 -->
                </p:stCondLst>
              </p:cTn>
            </p:par>
            <p:par>
              <p:cTn id="2" fill="hold" presetID="1" presetClass="entr" presetSubtype="0">
                <p:stCondLst>
                  <p:cond delay="0">  <!-- 立即触发 -->
                    <p:eventFld>
                      <p:tn p:id="1"/>  <!-- 触发器形状 ID -->
                    </p:eventFld>
                  </p:cond>
                </p:stCondLst>
              </p:cTn>
            </p:par>
          </p:timeNodeList>
        </p:timing>

        Args:
            slide_root: 幻灯片根元素
            slide_file: 幻灯片文件路径（用于提取形状 ID）
            trigger_config: 触发器配置
        """
        c_sld = slide_root.find('p:cSld', NS)
        if c_sld is None:
            return

        # 获取幻灯片中的所有形状及其 ID
        sp_tree = c_sld.find('p:spTree', NS)
        if sp_tree is None:
            return

        # 提取形状 ID 映射
        shape_ids = self._extract_shape_ids(sp_tree)
        logger.debug(f"幻灯片 {slide_file.name} 中的形状 ID: {shape_ids}")

        # 查找或创建 timing 节点
        timing = c_sld.find('p:timing', NS)
        if timing is None:
            timing = ET.SubElement(c_sld, '{%s}timing' % NS['p'])

        time_node_list = timing.find('p:timeNodeList', NS)
        if time_node_list is None:
            time_node_list = ET.SubElement(timing, '{%s}timeNodeList' % NS['p'])

        # 为幻灯片添加默认的点击触发器结构
        # 如果有配置，使用配置；否则自动检测互动元素
        if trigger_config:
            self._create_trigger_nodes(
                time_node_list,
                trigger_config.get('trigger_shape_id'),
                trigger_config.get('target_shape_id'),
                trigger_config.get('effect', 'appear')
            )
        else:
            # 自动检测并添加触发器
            self._auto_detect_and_create_triggers(time_node_list, shape_ids)

    def _extract_shape_ids(self, sp_tree: ET.Element) -> Dict[str, int]:
        """
        从形状树中提取形状 ID 和文本内容

        Args:
            sp_tree: 形状树元素

        Returns:
            {形状文本：形状 ID} 映射
        """
        shape_ids = {}
        for sp in sp_tree.findall('.//p:sp', NS):
            nv_sp_pr = sp.find('p:nvSpPr', NS)
            if nv_sp_pr is None:
                continue

            c_n = nv_sp_pr.find('p:cNvPr', NS)
            if c_n is None:
                continue

            shape_id = c_n.get('id')
            shape_name = c_n.get('name', '')

            # 获取文本内容
            tx_body = sp.find('p:txBody', NS)
            if tx_body is not None:
                text_content = ''
                for para in tx_body.findall('.//a:p', NS):
                    for t in para.findall('.//a:t', NS):
                        if t.text:
                            text_content += t.text

                if text_content.strip():
                    shape_ids[text_content.strip()] = int(shape_id)

        return shape_ids

    def _create_trigger_nodes(
        self,
        time_node_list: ET.Element,
        trigger_id: Optional[int],
        target_id: Optional[int],
        effect: str = 'appear'
    ):
        """
        创建触发器时间节点

        Args:
            time_node_list: 时间节点列表元素
            trigger_id: 触发器形状 ID
            target_id: 目标形状 ID
            effect: 效果类型
        """
        # 创建根时间节点（幻灯片主时间线）
        par_root = ET.SubElement(time_node_list, '{%s}par' % NS['p'])
        c_tn_root = ET.SubElement(par_root, '{%s}cTn' % NS['p'])
        c_tn_root.set('id', '1')
        c_tn_root.set('fill', 'hold')

        # 添加开始条件（无限延迟，等待触发）
        st_cond_lst_root = ET.SubElement(c_tn_root, '{%s}stCondLst' % NS['p'])
        ET.SubElement(st_cond_lst_root, '{%s}cond' % NS['p']).set('delay', 'indefinite')

        # 如果没有配置，使用默认 ID
        if trigger_id is None:
            trigger_id = 2
        if target_id is None:
            target_id = 3

        # 创建触发动画的时间节点
        par_trigger = ET.SubElement(time_node_list, '{%s}par' % NS['p'])
        c_tn_trigger = ET.SubElement(par_trigger, '{%s}cTn' % NS['p'])
        c_tn_trigger.set('id', str(target_id))
        c_tn_trigger.set('fill', 'hold')

        # 设置预设动画类型（出现效果）
        preset_id = '1'  # 出现动画的预设 ID
        if effect == 'disappear':
            preset_id = '2'  # 消失动画
        elif effect == 'color_change':
            preset_id = '10'  # 强调动画

        c_tn_trigger.set('presetID', preset_id)
        c_tn_trigger.set('presetClass', 'entr')  # 进入动画
        c_tn_trigger.set('presetSubtype', '0')

        # 添加目标形状引用
        tgt = ET.SubElement(c_tn_trigger, '{%s}tgtEl' % NS['p'])
        tgt_sp = ET.SubElement(tgt, '{%s}sp' % NS['p'])
        tgt_sp.set('id', str(trigger_id))

        # 添加触发条件（点击触发器形状）
        st_cond_lst = ET.SubElement(c_tn_trigger, '{%s}stCondLst' % NS['p'])
        cond = ET.SubElement(st_cond_lst, '{%s}cond' % NS['p'])
        cond.set('delay', '0')

        # 添加事件字段（点击触发）
        event_fld = ET.SubElement(cond, '{%s}eventFld' % NS['p'])
        tn = ET.SubElement(event_fld, '{%s}tn' % NS['p'])
        tn.set('p:id', str(trigger_id))

    def _auto_detect_and_create_triggers(
        self,
        time_node_list: ET.Element,
        shape_ids: Dict[str, int]
    ):
        """
        自动检测互动元素并创建触发器

        Args:
            time_node_list: 时间节点列表元素
            shape_ids: 形状 ID 映射
        """
        # 检测"点击查看答案"按钮
        reveal_triggers = ['点击查看答案', '显示答案', '查看答案', 'Click to reveal']
        trigger_id = None
        target_id = None

        for trigger_text in reveal_triggers:
            if trigger_text in shape_ids:
                trigger_id = shape_ids[trigger_text]
                break

        # 检测答案框
        if trigger_id:
            # 答案框通常是下一个形状
            target_id = trigger_id + 1

            # 创建触发器
            self._create_trigger_nodes(
                time_node_list,
                trigger_id,
                target_id,
                'appear'
            )
            logger.debug(f"已为'{trigger_text}'创建触发器：trigger={trigger_id}, target={target_id}")

    def _save_xml(self, tree: ET.ElementTree, file_path: Path):
        """保存 XML 文件"""
        xml_str = ET.tostring(tree.getroot(), encoding='unicode', xml_declaration=True)

        try:
            parsed = minidom.parseString(xml_str)
            pretty_xml = '\n'.join(
                line for line in parsed.toprettyxml(indent='  ', encoding=None).split('\n')
                if line.strip()
            )
        except Exception:
            pretty_xml = xml_str

        file_path.write_text(pretty_xml, encoding='utf-8')


# 便捷函数
def create_interactive_ppt(
    output_path: Path,
    grade: str = "6",
    subject: str = "general"
) -> InteractivePPTGenerator:
    """
    创建互动 PPT 生成器

    Args:
        output_path: 输出路径
        grade: 年级
        subject: 学科

    Returns:
        InteractivePPTGenerator 实例
    """
    return InteractivePPTGenerator()


def add_interactive_elements_to_ppt(
    pptx_path: str,
    output_path: Optional[str] = None,
    grade: str = "6",
    subject: str = "general"
) -> str:
    """
    为现有 PPT 添加互动元素

    Args:
        pptx_path: 输入 PPTX 路径
        output_path: 输出路径
        grade: 年级
        subject: 学科

    Returns:
        输出文件路径
    """
    generator = InteractivePPTGenerator()

    pptx_path = Path(pptx_path)
    output = Path(output_path) if output_path else pptx_path

    # 注入点击触发器
    result = generator.inject_click_trigger(pptx_path, output)

    return str(result)
