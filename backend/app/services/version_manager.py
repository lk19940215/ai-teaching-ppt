# -*- coding: utf-8 -*-
"""
版本管理服务
管理 PPT 会话、版本历史、页面状态

设计文档：.claude-coder/plans/ppt-merge-technical-design.md#5-版本化管理设计
"""

import os
import uuid
import shutil
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass, field, asdict
import logging

from ..models.ppt_structure import (
    SlideVersion, SlideStatus, SlideState, DocumentState, SessionData,
    DocumentData, SlideData
)
from .ppt_to_image import PptToImageConverter, Resolution, LibreOfficeDetector

logger = logging.getLogger(__name__)


class VersionManager:
    """版本管理器"""

    # 存储根目录
    VERSIONS_ROOT = Path("uploads/versions")

    def __init__(self):
        # 内存中的会话存储
        self._sessions: Dict[str, SessionData] = {}
        # 确保存储目录存在
        self.VERSIONS_ROOT.mkdir(parents=True, exist_ok=True)

    def create_session(self, documents: Dict[str, str]) -> SessionData:
        """
        创建新会话

        Args:
            documents: {document_id: pptx_path}，如 {"ppt_a": "uploads/xxx.pptx"}

        Returns:
            SessionData 会话数据
        """
        session_id = str(uuid.uuid4())[:8]
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        session = SessionData(
            session_id=session_id,
            documents={},
            created_at=now,
            last_updated=now
        )

        # 为每个文档创建初始状态
        for doc_id, pptx_path in documents.items():
            doc_state = self._create_document_state(doc_id, pptx_path, session_id)
            session.documents[doc_id] = doc_state

        # 保存会话
        self._sessions[session_id] = session
        logger.info(f"创建会话 {session_id}，包含 {len(documents)} 个文档")

        return session

    def _create_document_state(
        self,
        doc_id: str,
        pptx_path: str,
        session_id: str
    ) -> DocumentState:
        """
        创建文档初始状态

        Args:
            doc_id: 文档 ID (ppt_a, ppt_b)
            pptx_path: PPTX 文件路径
            session_id: 会话 ID

        Returns:
            DocumentState 文档状态
        """
        from pptx import Presentation

        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise ValueError(f"文件不存在：{pptx_path}")

        # 获取总页数
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)

        # 创建会话目录
        session_dir = self.VERSIONS_ROOT / session_id
        session_dir.mkdir(parents=True, exist_ok=True)

        # 转换 PPT 为图片（使用已有的 PptToImageConverter）
        converter = PptToImageConverter(session_dir, Resolution.HIGH)
        result = converter.convert(pptx_path)

        if not result.success:
            logger.error(f"PPT 转图片失败：{result.error}")
            # 继续执行，但版本图片可能缺失

        # 为每页创建初始版本 v1
        slides = {}
        now = datetime.now().strftime("%H:%M:%S")

        for slide_idx in range(total_slides):
            # 查找对应的图片
            image_url = None
            for img in result.images:
                if img["page"] == slide_idx:
                    image_url = img["url"]
                    break

            version = SlideVersion(
                version="v1",
                image_url=image_url or "",
                created_at=now,
                operation="原始上传",
                source_pptx=str(pptx_path)
            )

            slides[slide_idx] = SlideState(
                current_version="v1",
                status=SlideStatus.ACTIVE,
                versions=[version]
            )

        logger.info(f"创建文档 {doc_id} 状态，共 {total_slides} 页")

        return DocumentState(
            source_file=str(pptx_path),
            slides=slides
        )

    def get_session(self, session_id: str) -> Optional[SessionData]:
        """获取会话数据"""
        return self._sessions.get(session_id)

    def create_version(
        self,
        session_id: str,
        document_id: str,
        slide_index: int,
        operation: str,
        prompt: Optional[str] = None,
        new_pptx: Optional[Path] = None,
        content_snapshot: Optional[Dict[str, Any]] = None
    ) -> SlideVersion:
        """
        创建新版本

        Args:
            session_id: 会话 ID
            document_id: 文档 ID (ppt_a, ppt_b)
            slide_index: 页码（0-indexed）
            operation: 操作类型 (ai_polish, ai_expand, ai_rewrite, ai_extract)
            prompt: AI 提示语
            new_pptx: 新生成的单页 PPTX 路径（用于转换新图片）
            content_snapshot: AI 修改的内容快照（用于最终生成 PPT）

        Returns:
            SlideVersion 新版本数据
        """
        session = self._sessions.get(session_id)
        if not session:
            raise ValueError(f"会话不存在：{session_id}")

        doc = session.documents.get(document_id)
        if not doc:
            raise ValueError(f"文档不存在：{document_id}")

        slide = doc.slides.get(slide_index)
        if not slide:
            raise ValueError(f"页面不存在：slide {slide_index}")

        # 计算新版本号
        current_version_num = int(slide.versions[-1].version[1:])  # v1 -> 1
        new_version_num = current_version_num + 1
        new_version_str = f"v{new_version_num}"

        # 转换新 PPTX 为图片（如果提供了）
        image_url = ""
        if new_pptx and new_pptx.exists():
            session_dir = self.VERSIONS_ROOT / session_id
            converter = PptToImageConverter(session_dir, Resolution.HIGH)

            # 创建临时目录
            temp_dir = session_dir / "temp"
            temp_dir.mkdir(exist_ok=True)

            # 转换单页
            result = converter.convert(new_pptx, pages=[0])
            if result.success and len(result.images) > 0:
                # 移动图片到会话目录
                img_path = Path(result.images[0]["path"])
                new_img_name = f"{document_id}_slide{slide_index}_{new_version_str}.png"
                new_img_path = session_dir / new_img_name
                shutil.move(str(img_path), str(new_img_path))
                image_url = f"/static/versions/{session_id}/{new_img_name}"

            # 清理临时目录
            try:
                shutil.rmtree(temp_dir)
            except Exception:
                pass

        # 创建新版本
        now = datetime.now().strftime("%H:%M:%S")
        version = SlideVersion(
            version=new_version_str,
            image_url=image_url,
            created_at=now,
            operation=operation,
            prompt=prompt,
            content_snapshot=content_snapshot
        )

        # 添加到版本列表
        slide.versions.append(version)
        slide.current_version = new_version_str
        slide.status = SlideStatus.ACTIVE
        session.last_updated = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        logger.info(f"创建版本 {new_version_str} for {document_id} slide {slide_index} ({operation})")

        return version

    def restore_version(
        self,
        session_id: str,
        document_id: str,
        slide_index: int,
        target_version: str
    ) -> bool:
        """
        恢复历史版本

        Args:
            session_id: 会话 ID
            document_id: 文档 ID
            slide_index: 页码
            target_version: 目标版本号 (v1, v2, ...)

        Returns:
            bool 是否成功
        """
        session = self._sessions.get(session_id)
        if not session:
            raise ValueError(f"会话不存在：{session_id}")

        doc = session.documents.get(document_id)
        if not doc:
            raise ValueError(f"文档不存在：{document_id}")

        slide = doc.slides.get(slide_index)
        if not slide:
            raise ValueError(f"页面不存在：slide {slide_index}")

        # 查找目标版本
        target = None
        for v in slide.versions:
            if v.version == target_version:
                target = v
                break

        if not target:
            raise ValueError(f"版本不存在：{target_version}")

        # 更新当前版本
        slide.current_version = target_version
        slide.status = SlideStatus.ACTIVE
        session.last_updated = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        logger.info(f"恢复 {document_id} slide {slide_index} 到版本 {target_version}")

        return True

    def toggle_slide(
        self,
        session_id: str,
        document_id: str,
        slide_index: int,
        action: str
    ) -> SlideStatus:
        """
        删除/恢复页面

        Args:
            session_id: 会话 ID
            document_id: 文档 ID
            slide_index: 页码
            action: "delete" 或 "restore"

        Returns:
            SlideStatus 操作后的状态
        """
        session = self._sessions.get(session_id)
        if not session:
            raise ValueError(f"会话不存在：{session_id}")

        doc = session.documents.get(document_id)
        if not doc:
            raise ValueError(f"文档不存在：{document_id}")

        slide = doc.slides.get(slide_index)
        if not slide:
            raise ValueError(f"页面不存在：slide {slide_index}")

        if action == "delete":
            slide.status = SlideStatus.DELETED
            slide.current_version = None
            logger.info(f"删除 {document_id} slide {slide_index}")
        elif action == "restore":
            slide.status = SlideStatus.ACTIVE
            # 恢复到最后一个版本
            if slide.versions:
                slide.current_version = slide.versions[-1].version
            logger.info(f"恢复 {document_id} slide {slide_index}")
        else:
            raise ValueError(f"未知操作：{action}")

        session.last_updated = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        return slide.status

    def get_version_history(
        self,
        session_id: str,
        document_id: str,
        slide_index: int
    ) -> List[SlideVersion]:
        """获取版本历史"""
        session = self._sessions.get(session_id)
        if not session:
            return []

        doc = session.documents.get(document_id)
        if not doc:
            return []

        slide = doc.slides.get(slide_index)
        if not slide:
            return []

        return slide.versions

    def get_active_slides(self, session_id: str, document_id: str) -> List[int]:
        """获取所有活跃页面索引"""
        session = self._sessions.get(session_id)
        if not session:
            return []

        doc = session.documents.get(document_id)
        if not doc:
            return []

        return [
            idx for idx, slide in doc.slides.items()
            if slide.status == SlideStatus.ACTIVE
        ]

    def get_final_merge_plan(
        self,
        session_id: str
    ) -> Dict[str, Dict[int, str]]:
        """
        获取最终合并计划

        Returns:
            {document_id: {slide_index: version}}
        """
        session = self._sessions.get(session_id)
        if not session:
            return {}

        plan = {}
        for doc_id, doc in session.documents.items():
            plan[doc_id] = {
                idx: slide.current_version
                for idx, slide in doc.slides.items()
                if slide.status == SlideStatus.ACTIVE and slide.current_version
            }

        return plan

    def cleanup_session(self, session_id: str) -> bool:
        """清理会话（删除版本图片）"""
        session_dir = self.VERSIONS_ROOT / session_id
        if session_dir.exists():
            try:
                shutil.rmtree(session_dir)
                logger.info(f"清理会话 {session_id}")
            except Exception as e:
                logger.error(f"清理会话失败：{e}")
                return False

        if session_id in self._sessions:
            del self._sessions[session_id]

        return True


# 全局单例
_version_manager: Optional[VersionManager] = None


def get_version_manager() -> VersionManager:
    """获取版本管理器单例"""
    global _version_manager
    if _version_manager is None:
        _version_manager = VersionManager()
    return _version_manager
